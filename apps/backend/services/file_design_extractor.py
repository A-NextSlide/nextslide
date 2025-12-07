"""
File Design & Content Extractor
Analyzes uploaded files (PPTX, PDF, images) to extract:
- Design: colors, fonts, layout patterns, visual style
- Content: text, data, key points, structure
- Intent: what the user wants to do with the file
"""

import os
import base64
import asyncio
import logging
import tempfile
import subprocess
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from io import BytesIO
from dataclasses import dataclass, field
from enum import Enum

logger = logging.getLogger(__name__)


class FileIntent(Enum):
    """What the user wants to do with the uploaded file"""
    RECREATE_EXACT = "recreate_exact"      # "recreate this", "make this exact slide"
    USE_DESIGN_ONLY = "use_design_only"    # "use this style", "match this design"
    USE_CONTENT_ONLY = "use_content_only"  # "use this content", "summarize this"
    USE_BOTH = "use_both"                  # "make slides like this", "similar to this"
    REFERENCE_ONLY = "reference_only"      # "here's context", "for reference"


class SlideStyle(Enum):
    """Type of slides to generate"""
    TRADITIONAL = "traditional"    # Standard components (text blocks, images, shapes)
    INTERACTIVE = "interactive"    # CustomComponents with HTML/CSS/animations
    AUTO = "auto"                  # Let AI decide based on content


@dataclass
class ExtractedDesign:
    """Design elements extracted from a file"""
    primary_color: str = "#1a1a2e"
    secondary_color: str = "#4a4a6a"
    accent_color: str = "#0066ff"
    background_color: str = "#ffffff"
    text_color: str = "#1a1a2e"
    color_palette: List[str] = field(default_factory=list)

    hero_font: str = "Inter"
    body_font: str = "Inter"
    font_weights: List[int] = field(default_factory=lambda: [400, 600, 700])

    layout_style: str = "modern"  # modern, corporate, playful, minimal, bold
    visual_density: str = "balanced"  # sparse, balanced, dense
    has_animations: bool = False
    has_gradients: bool = False
    has_shadows: bool = False

    style_keywords: List[str] = field(default_factory=list)
    raw_analysis: str = ""


@dataclass
class ExtractedContent:
    """Content extracted from a file"""
    title: str = ""
    subtitle: str = ""
    main_points: List[str] = field(default_factory=list)
    sections: List[Dict[str, Any]] = field(default_factory=list)
    data_points: List[Dict[str, Any]] = field(default_factory=list)
    images_described: List[str] = field(default_factory=list)
    total_slides: int = 0
    raw_text: str = ""
    summary: str = ""


@dataclass
class FileAnalysis:
    """Complete analysis of an uploaded file"""
    filename: str
    file_type: str
    intent: FileIntent
    slide_style: SlideStyle
    design: ExtractedDesign
    content: ExtractedContent
    confidence: float = 0.8
    analysis_notes: str = ""
    extracted_images: List[str] = field(default_factory=list)  # Uploaded image URLs from the file
    slide_screenshots: List[str] = field(default_factory=list)  # Base64 screenshots of slides for visual design reference


class FileDesignExtractor:
    """Extracts design and content from uploaded files using vision AI"""

    def __init__(self):
        self.temp_dir = None

    async def analyze_file(
        self,
        file_content: str,  # base64 encoded
        filename: str,
        file_type: str,
        user_message: str = ""
    ) -> FileAnalysis:
        """
        Analyze a file to extract design, content, and determine intent.

        Args:
            file_content: Base64 encoded file content
            filename: Original filename
            file_type: MIME type
            user_message: User's message to determine intent
        """
        logger.info(f"[FileExtractor] Analyzing: {filename} ({file_type})")

        # Determine intent from user message
        intent = self._detect_intent(user_message)
        slide_style = self._detect_slide_style(user_message)

        logger.info(f"[FileExtractor] Detected intent: {intent.value}, style: {slide_style.value}")

        # Convert file to images for vision analysis
        images = await self._convert_to_images(file_content, filename, file_type)

        if not images:
            logger.warning(f"[FileExtractor] Could not convert {filename} to images")
            return FileAnalysis(
                filename=filename,
                file_type=file_type,
                intent=intent,
                slide_style=slide_style,
                design=ExtractedDesign(),
                content=ExtractedContent(),
                confidence=0.3,
                analysis_notes="Could not extract images from file"
            )

        logger.info(f"[FileExtractor] Converted to {len(images)} images")

        # Extract design, content, and images in parallel
        design_task = self._extract_design(images[:5])  # First 5 slides for design analysis
        content_task = self._extract_content(images, file_content, filename, file_type)
        images_task = self._extract_and_upload_images(file_content, filename, file_type)

        design, content, extracted_images = await asyncio.gather(design_task, content_task, images_task)

        logger.info(f"[FileExtractor] Extracted {len(extracted_images)} images from file")

        # Store slide screenshots for visual design reference
        # For recreate/remake requests, store ALL slides so AI can replicate each one
        # For other intents, store up to 10 slides (first 5 + sample from rest)
        if intent == FileIntent.RECREATE_EXACT:
            # Store ALL slides for exact recreation
            slide_screenshots = images if images else []
            logger.info(f"[FileExtractor] RECREATE mode: Stored ALL {len(slide_screenshots)} slide screenshots")
        else:
            # For design inspiration, sample slides: first 5 + middle + last 2
            if len(images) <= 10:
                slide_screenshots = images
            else:
                # First 5, middle sample, last 2
                middle_idx = len(images) // 2
                slide_screenshots = images[:5] + [images[middle_idx]] + images[-2:]
            logger.info(f"[FileExtractor] Stored {len(slide_screenshots)} slide screenshots for visual reference")

        return FileAnalysis(
            filename=filename,
            file_type=file_type,
            intent=intent,
            slide_style=slide_style,
            design=design,
            content=content,
            confidence=0.85,
            analysis_notes=f"Analyzed {len(images)} slides/pages",
            extracted_images=extracted_images,
            slide_screenshots=slide_screenshots
        )

    def _detect_intent(self, message: str) -> FileIntent:
        """Detect user intent from their message"""
        message_lower = message.lower()

        # Recreate exact
        recreate_phrases = [
            "recreate", "remake", "copy this", "make this exact",
            "same as this", "duplicate", "replicate", "convert this"
        ]
        if any(phrase in message_lower for phrase in recreate_phrases):
            return FileIntent.RECREATE_EXACT

        # Design indicators
        design_phrases = [
            "use this style", "match this design", "this aesthetic",
            "like the look of", "same colors", "this theme", "design from",
            "style from", "look and feel", "same design", "same style",
            "keep the design", "match the style", "use the design"
        ]
        has_design_intent = any(phrase in message_lower for phrase in design_phrases)

        # Content indicators
        content_phrases = [
            "use this content", "summarize this", "based on this data",
            "from this document", "extract from", "content from",
            "information from", "about this topic", "summarize", "use the content"
        ]
        has_content_intent = any(phrase in message_lower for phrase in content_phrases)

        # If BOTH design and content mentioned, use both
        if has_design_intent and has_content_intent:
            return FileIntent.USE_BOTH

        # Design only
        if has_design_intent:
            return FileIntent.USE_DESIGN_ONLY

        # Content only
        if has_content_intent:
            return FileIntent.USE_CONTENT_ONLY

        # Both (default for "like this")
        both_phrases = [
            "like this", "similar to", "inspired by", "based on this",
            "something like", "make slides", "create presentation"
        ]
        if any(phrase in message_lower for phrase in both_phrases):
            return FileIntent.USE_BOTH

        # Reference only
        reference_phrases = [
            "for reference", "here's context", "fyi", "background info",
            "attached", "see attached"
        ]
        if any(phrase in message_lower for phrase in reference_phrases):
            return FileIntent.REFERENCE_ONLY

        # Default to using both if file is attached with a request
        return FileIntent.USE_BOTH

    def _detect_slide_style(self, message: str) -> SlideStyle:
        """Detect whether to use traditional or interactive slides"""
        message_lower = message.lower()

        # Interactive indicators
        interactive_phrases = [
            "interactive", "animated", "dynamic", "modern",
            "engaging", "cool", "fancy", "impressive",
            "wow", "standout", "unique", "creative",
            "html", "custom", "special effects"
        ]
        if any(phrase in message_lower for phrase in interactive_phrases):
            return SlideStyle.INTERACTIVE

        # Traditional indicators
        traditional_phrases = [
            "simple", "basic", "traditional", "standard",
            "classic", "professional", "corporate", "formal",
            "clean", "minimal", "straightforward"
        ]
        if any(phrase in message_lower for phrase in traditional_phrases):
            return SlideStyle.TRADITIONAL

        # Default to auto
        return SlideStyle.AUTO

    async def _convert_to_images(
        self,
        file_content: str,
        filename: str,
        file_type: str
    ) -> List[str]:
        """Convert file to base64 images for vision analysis"""
        images = []

        try:
            raw_bytes = base64.b64decode(file_content)

            # Create temp directory
            self.temp_dir = tempfile.mkdtemp(prefix="file_extract_")

            if file_type == "application/pdf" or filename.lower().endswith(".pdf"):
                images = await self._pdf_to_images(raw_bytes)
            elif "presentation" in file_type or filename.lower().endswith((".pptx", ".ppt")):
                images = await self._pptx_to_images(raw_bytes)
            elif file_type.startswith("image/"):
                # Already an image
                images = [file_content]

        except Exception as e:
            logger.error(f"[FileExtractor] Error converting file: {e}")
        finally:
            # Cleanup temp dir
            if self.temp_dir:
                try:
                    import shutil
                    shutil.rmtree(self.temp_dir, ignore_errors=True)
                except:
                    pass

        return images

    async def _pdf_to_images(self, pdf_bytes: bytes) -> List[str]:
        """Convert PDF pages to images"""
        images = []
        try:
            from pdf2image import convert_from_bytes
            pil_images = convert_from_bytes(pdf_bytes, first_page=1, last_page=10, dpi=150)

            for img in pil_images:
                buffered = BytesIO()
                img.save(buffered, format="PNG")
                img_b64 = base64.b64encode(buffered.getvalue()).decode()
                images.append(img_b64)

        except ImportError:
            logger.warning("[FileExtractor] pdf2image not installed, trying alternative")
            # Fallback: use subprocess with pdftoppm if available
            try:
                pdf_path = os.path.join(self.temp_dir, "input.pdf")
                with open(pdf_path, "wb") as f:
                    f.write(pdf_bytes)

                output_prefix = os.path.join(self.temp_dir, "page")
                subprocess.run(
                    ["pdftoppm", "-png", "-r", "150", "-l", "10", pdf_path, output_prefix],
                    capture_output=True, timeout=60
                )

                for img_path in sorted(Path(self.temp_dir).glob("page-*.png")):
                    with open(img_path, "rb") as f:
                        img_b64 = base64.b64encode(f.read()).decode()
                        images.append(img_b64)
            except Exception as e:
                logger.error(f"[FileExtractor] PDF conversion failed: {e}")

        return images

    async def _pptx_to_images(self, pptx_bytes: bytes) -> List[str]:
        """Convert PPTX slides to images"""
        images = []
        try:
            pptx_path = os.path.join(self.temp_dir, "input.pptx")
            with open(pptx_path, "wb") as f:
                f.write(pptx_bytes)

            # Convert PPTX to PDF first
            pdf_path = os.path.join(self.temp_dir, "input.pdf")
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", self.temp_dir, pptx_path],
                capture_output=True, timeout=120
            )

            if os.path.exists(pdf_path):
                with open(pdf_path, "rb") as f:
                    images = await self._pdf_to_images(f.read())
            else:
                logger.warning(f"[FileExtractor] LibreOffice conversion failed: {result.stderr}")

        except Exception as e:
            logger.error(f"[FileExtractor] PPTX conversion failed: {e}")

        return images

    async def _extract_design(self, images: List[str]) -> ExtractedDesign:
        """Extract design elements using vision AI"""
        if not images:
            return ExtractedDesign()

        try:
            from google import genai
            from agents.config import GEMINI_FLASH_LITE

            client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY"))

            # Build image parts
            image_parts = []
            for i, img_b64 in enumerate(images[:3]):  # First 3 slides for design
                image_parts.append({
                    "inline_data": {
                        "mime_type": "image/png",
                        "data": img_b64
                    }
                })

            prompt = """Analyze these presentation slides and extract the DESIGN SYSTEM.

Return a JSON object with:
{
    "primary_color": "#hex - main brand/heading color",
    "secondary_color": "#hex - secondary color",
    "accent_color": "#hex - accent/highlight color",
    "background_color": "#hex - main background",
    "text_color": "#hex - main text color",
    "color_palette": ["#hex", "#hex", ...] - all colors used,

    "hero_font": "Font name for headings (guess from style if unsure)",
    "body_font": "Font name for body text",
    "font_weights": [400, 600, 700] - weights used,

    "layout_style": "modern|corporate|playful|minimal|bold|elegant",
    "visual_density": "sparse|balanced|dense",
    "has_animations": false,
    "has_gradients": true/false,
    "has_shadows": true/false,

    "style_keywords": ["keyword1", "keyword2"] - words describing the style
}

Be precise with hex colors. Analyze actual colors from the slides."""

            contents = image_parts + [{"text": prompt}]

            response = client.models.generate_content(
                model=GEMINI_FLASH_LITE,
                contents=contents
            )

            response_text = response.text.strip()
            logger.info(f"[FileExtractor] Design response: {response_text[:500]}")

            # Parse JSON from response
            import json
            import re

            # Extract JSON from response
            json_match = re.search(r'\{[\s\S]*\}', response_text)
            if json_match:
                design_data = json.loads(json_match.group())

                return ExtractedDesign(
                    primary_color=design_data.get("primary_color", "#1a1a2e"),
                    secondary_color=design_data.get("secondary_color", "#4a4a6a"),
                    accent_color=design_data.get("accent_color", "#0066ff"),
                    background_color=design_data.get("background_color", "#ffffff"),
                    text_color=design_data.get("text_color", "#1a1a2e"),
                    color_palette=design_data.get("color_palette", []),
                    hero_font=design_data.get("hero_font", "Inter"),
                    body_font=design_data.get("body_font", "Inter"),
                    font_weights=design_data.get("font_weights", [400, 600, 700]),
                    layout_style=design_data.get("layout_style", "modern"),
                    visual_density=design_data.get("visual_density", "balanced"),
                    has_animations=design_data.get("has_animations", False),
                    has_gradients=design_data.get("has_gradients", False),
                    has_shadows=design_data.get("has_shadows", False),
                    style_keywords=design_data.get("style_keywords", []),
                    raw_analysis=response_text
                )

        except Exception as e:
            logger.error(f"[FileExtractor] Design extraction failed: {e}")

        return ExtractedDesign()

    async def _extract_content(
        self,
        images: List[str],
        file_content: str,
        filename: str,
        file_type: str
    ) -> ExtractedContent:
        """Extract content from file"""
        content = ExtractedContent()
        content.total_slides = len(images)

        try:
            # Try text extraction first (faster, more accurate for text)
            raw_bytes = base64.b64decode(file_content)

            if file_type == "application/pdf" or filename.lower().endswith(".pdf"):
                content.raw_text = await self._extract_pdf_text(raw_bytes)
            elif "presentation" in file_type or filename.lower().endswith((".pptx", ".ppt")):
                content.raw_text = await self._extract_pptx_text(raw_bytes)

            # Use vision for structure/summary if we have images
            if images:
                await self._analyze_content_with_vision(images, content)
            elif content.raw_text:
                await self._analyze_content_from_text(content)

        except Exception as e:
            logger.error(f"[FileExtractor] Content extraction failed: {e}")

        return content

    async def _extract_pdf_text(self, pdf_bytes: bytes) -> str:
        """Extract text from PDF"""
        try:
            import pypdf
            reader = pypdf.PdfReader(BytesIO(pdf_bytes))
            text_parts = []
            for i, page in enumerate(reader.pages[:50]):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"--- Page {i+1} ---\n{page_text}")
            return "\n\n".join(text_parts)[:100000]
        except Exception as e:
            logger.warning(f"[FileExtractor] PDF text extraction failed: {e}")
            return ""

    async def _extract_pptx_text(self, pptx_bytes: bytes) -> str:
        """Extract text from PPTX"""
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(pptx_bytes))
            text_parts = []

            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                if slide_texts:
                    text_parts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_texts))

            return "\n\n".join(text_parts)[:100000]
        except Exception as e:
            logger.warning(f"[FileExtractor] PPTX text extraction failed: {e}")
            return ""

    async def _analyze_content_with_vision(self, images: List[str], content: ExtractedContent):
        """Use vision to analyze content structure"""
        try:
            from google import genai
            from agents.config import GEMINI_FLASH_LITE

            client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY"))

            # Sample slides for analysis
            sample_images = images[:5] + images[-2:] if len(images) > 5 else images

            image_parts = []
            for img_b64 in sample_images:
                image_parts.append({
                    "inline_data": {
                        "mime_type": "image/png",
                        "data": img_b64
                    }
                })

            prompt = """Analyze these presentation slides and extract the CONTENT STRUCTURE.

Return a JSON object:
{
    "title": "Main presentation title",
    "subtitle": "Subtitle if any",
    "summary": "2-3 sentence summary of what this presentation is about",
    "main_points": ["Key point 1", "Key point 2", ...],
    "sections": [
        {"title": "Section name", "points": ["point 1", "point 2"]}
    ],
    "data_points": [
        {"label": "Metric name", "value": "value", "context": "what it means"}
    ],
    "images_described": ["Description of key images/charts used"]
}

Extract the actual content, data, and structure from these slides."""

            contents = image_parts + [{"text": prompt}]

            response = client.models.generate_content(
                model=GEMINI_FLASH_LITE,
                contents=contents
            )

            response_text = response.text.strip()

            import json
            import re

            json_match = re.search(r'\{[\s\S]*\}', response_text)
            if json_match:
                data = json.loads(json_match.group())
                content.title = data.get("title", "")
                content.subtitle = data.get("subtitle", "")
                content.summary = data.get("summary", "")
                content.main_points = data.get("main_points", [])
                content.sections = data.get("sections", [])
                content.data_points = data.get("data_points", [])
                content.images_described = data.get("images_described", [])

        except Exception as e:
            logger.error(f"[FileExtractor] Vision content analysis failed: {e}")

    async def _analyze_content_from_text(self, content: ExtractedContent):
        """Analyze content from extracted text using Claude"""
        if not content.raw_text:
            return

        try:
            from agents.ai.clients import get_client, invoke
            from agents.config import FILE_ANALYSIS_MODEL_FAST

            client, actual_model = get_client(FILE_ANALYSIS_MODEL_FAST)

            prompt = f"""Analyze this presentation text and extract the content structure.

TEXT:
{content.raw_text[:30000]}

Return a JSON object:
{{
    "title": "Main presentation title",
    "subtitle": "Subtitle if any",
    "summary": "2-3 sentence summary",
    "main_points": ["Key point 1", "Key point 2", ...],
    "sections": [{{"title": "Section", "points": ["point"]}}],
    "data_points": [{{"label": "Metric", "value": "value", "context": "meaning"}}]
}}"""

            response = invoke(
                client=client,
                model=actual_model,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=4000,
                temperature=0.3
            )

            import json
            import re

            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                data = json.loads(json_match.group())
                content.title = data.get("title", "")
                content.subtitle = data.get("subtitle", "")
                content.summary = data.get("summary", "")
                content.main_points = data.get("main_points", [])
                content.sections = data.get("sections", [])
                content.data_points = data.get("data_points", [])

        except Exception as e:
            logger.error(f"[FileExtractor] Text content analysis failed: {e}")

    async def _extract_and_upload_images(
        self,
        file_content: str,
        filename: str,
        file_type: str
    ) -> List[str]:
        """Extract images from PPTX/PDF and upload to storage. Returns list of image URLs."""
        from services.image_storage_service import ImageStorageService

        extracted_urls = []

        try:
            raw_bytes = base64.b64decode(file_content)

            # Only extract from PPTX for now (PPTX has embedded images we can extract)
            if "presentation" in file_type or filename.lower().endswith((".pptx", ".ppt")):
                extracted_urls = await self._extract_pptx_images(raw_bytes, filename)
            # PDF image extraction is more complex, skip for now
            # Images uploaded directly are already handled separately

        except Exception as e:
            logger.error(f"[FileExtractor] Image extraction failed: {e}")

        return extracted_urls

    async def _extract_pptx_images(self, pptx_bytes: bytes, filename: str) -> List[str]:
        """Extract all images from a PPTX file and upload to storage."""
        from services.image_storage_service import ImageStorageService

        image_urls = []

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(BytesIO(pptx_bytes))

            async with ImageStorageService() as storage:
                img_counter = 0

                for slide_idx, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        try:
                            # Check if shape has an image
                            shape_type = getattr(shape, 'shape_type', None)
                            has_picture = shape_type == getattr(MSO_SHAPE_TYPE, 'PICTURE', None)

                            if has_picture and hasattr(shape, 'image'):
                                image = shape.image
                                image_bytes = image.blob

                                # Determine content type and extension
                                content_type = image.content_type or 'image/png'
                                ext = '.png'
                                if 'jpeg' in content_type or 'jpg' in content_type:
                                    ext = '.jpg'
                                elif 'gif' in content_type:
                                    ext = '.gif'
                                elif 'webp' in content_type:
                                    ext = '.webp'

                                # Convert bytes to base64
                                image_b64 = base64.b64encode(image_bytes).decode('utf-8')

                                # Upload to storage
                                img_counter += 1
                                upload_filename = f"extracted_slide{slide_idx + 1}_img{img_counter}{ext}"
                                result = await storage.upload_image_from_base64(
                                    image_b64,
                                    filename=upload_filename,
                                    content_type=content_type,
                                    folder="pptx-extracted"
                                )

                                if result and result.get('url'):
                                    url = result['url'].split('?')[0]  # Remove query params
                                    image_urls.append(url)
                                    logger.info(f"[FileExtractor] Uploaded image {img_counter} from slide {slide_idx + 1}")

                        except Exception as e:
                            logger.warning(f"[FileExtractor] Failed to extract image from shape: {e}")
                            continue

            logger.info(f"[FileExtractor] Extracted {len(image_urls)} images from PPTX")

        except ImportError:
            logger.warning("[FileExtractor] python-pptx not installed, skipping image extraction")
        except Exception as e:
            logger.error(f"[FileExtractor] PPTX image extraction error: {e}")

        return image_urls


def design_to_theme_context(design: ExtractedDesign) -> Dict[str, Any]:
    """Convert ExtractedDesign to theme context for generation"""
    return {
        "color_palette": {
            "primary": design.primary_color,
            "secondary": design.secondary_color,
            "accent": design.accent_color,
            "background": design.background_color,
            "text": design.text_color,
            "all_colors": design.color_palette
        },
        "typography": {
            "hero_font": design.hero_font,
            "body_font": design.body_font,
            "font_weights": design.font_weights
        },
        "style": {
            "layout": design.layout_style,
            "density": design.visual_density,
            "keywords": design.style_keywords,
            "has_gradients": design.has_gradients,
            "has_shadows": design.has_shadows
        }
    }


def content_to_outline_context(content: ExtractedContent) -> str:
    """Convert ExtractedContent to context string for outline generation"""
    parts = []

    if content.title:
        parts.append(f"TITLE: {content.title}")
    if content.subtitle:
        parts.append(f"SUBTITLE: {content.subtitle}")
    if content.summary:
        parts.append(f"SUMMARY: {content.summary}")

    if content.main_points:
        parts.append("KEY POINTS:")
        for point in content.main_points:
            parts.append(f"  - {point}")

    if content.sections:
        parts.append("SECTIONS:")
        for section in content.sections:
            parts.append(f"  {section.get('title', 'Section')}:")
            for point in section.get('points', []):
                parts.append(f"    - {point}")

    if content.data_points:
        parts.append("DATA/METRICS:")
        for dp in content.data_points:
            parts.append(f"  - {dp.get('label')}: {dp.get('value')} ({dp.get('context', '')})")

    return "\n".join(parts)
