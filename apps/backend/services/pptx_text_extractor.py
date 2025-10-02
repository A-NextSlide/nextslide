"""
Lightweight PPTX text extractor

Extracts per-slide text content and notes from a .pptx file using python-pptx.
Designed for outline generation where we want titles + main text per slide.
"""

from io import BytesIO
from typing import Any, Dict, List, Optional

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore


def _safe_get_title(slide) -> Optional[str]:
    try:
        # python-pptx exposes a convenience property
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and getattr(title_shape, "text", None):
            text = str(title_shape.text).strip()
            if text:
                return text
    except Exception:
        pass
    # Fallback: first text-bearing shape
    try:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame and str(shape.text_frame.text or "").strip():
                txt = str(shape.text_frame.text).strip()
                if txt:
                    return txt
    except Exception:
        pass
    return None


def _collect_text_items(slide) -> List[str]:
    items: List[str] = []
    try:
        for shape in slide.shapes:
            # Tables
            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        row_texts: List[str] = []
                        for cell in row.cells:
                            cell_text = str(getattr(cell, "text", "") or "").strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            items.append(" | ".join(row_texts))
                except Exception:
                    # Ignore malformed tables
                    pass
                continue

            # Text frames
            if getattr(shape, "has_text_frame", False) and getattr(shape, "text_frame", None):
                try:
                    txt = str(shape.text_frame.text or "").strip()
                    if txt:
                        items.append(txt)
                except Exception:
                    pass
    except Exception:
        pass
    return items


def _get_notes(slide) -> str:
    try:
        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide and getattr(notes_slide, "notes_text_frame", None):
            txt = str(notes_slide.notes_text_frame.text or "").strip()
            return txt
    except Exception:
        pass
    return ""


def extract_pptx_text_from_bytes(file_bytes: bytes) -> Dict[str, Any]:
    """
    Extract per-slide title, text items, consolidated text, and notes from PPTX bytes.

    Returns:
        {
          "slide_count": int,
          "slides": [
            {"index": int, "title": str, "text_items": [str], "text": str, "notes": str}
          ]
        }
    """
    if Presentation is None:
        raise RuntimeError("python-pptx is not available. Please install 'python-pptx'.")

    prs = Presentation(BytesIO(file_bytes))
    result: Dict[str, Any] = {
        "slide_count": len(prs.slides),
        "slides": []
    }

    for idx, slide in enumerate(prs.slides, start=1):
        title = _safe_get_title(slide) or f"Slide {idx}"
        text_items = _collect_text_items(slide)
        # Avoid duplicating title if it appears in first text item
        filtered_items: List[str] = []
        for i, item in enumerate(text_items):
            if i == 0 and isinstance(title, str) and item.strip() == str(title).strip():
                continue
            filtered_items.append(item)
        full_text = "\n".join(filtered_items).strip()
        notes = _get_notes(slide)
        result["slides"].append({
            "index": idx,
            "title": title,
            "text_items": filtered_items,
            "text": full_text,
            "notes": notes
        })

    return result


