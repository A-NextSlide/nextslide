"""
Clean PPTX Import Service
Converts PowerPoint presentations to our internal format using python-pptx
"""

import uuid
import logging
from typing import Dict, Any, List, Optional, Tuple, Union
from io import BytesIO
import base64

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu, Length
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

class PPTXImporter:
    """
    Clean implementation of PPTX import functionality.
    Converts PowerPoint files to our slide format.
    """
    
    def __init__(self):
        self.stats = {
            "slides": 0,
            "components": 0,
            "images": 0,
            "text_blocks": 0,
            "shapes": 0,
            "errors": 0
        }
        
    async def import_file(self, file_path: str) -> Dict[str, Any]:
        """Import a PPTX file from disk"""
        try:
            prs = Presentation(file_path)
            # Store file path for theme extraction
            self._file_path = file_path
            return await self._process_presentation(prs)
        except Exception as e:
            logger.error(f"Failed to import PPTX file: {e}")
            raise
            
    async def import_bytes(self, file_bytes: bytes) -> Dict[str, Any]:
        """Import a PPTX file from bytes"""
        try:
            prs = Presentation(BytesIO(file_bytes))
            return await self._process_presentation(prs)
        except Exception as e:
            logger.error(f"Failed to import PPTX bytes: {e}")
            raise
            
    async def _process_presentation(self, prs: Presentation) -> Dict[str, Any]:
        """Process a PowerPoint presentation object"""
        logger.info(f"Processing presentation with {len(prs.slides)} slides")
        
        # Reset stats
        self.stats = {
            "slides": len(prs.slides),
            "components": 0,
            "images": 0,
            "text_blocks": 0,
            "shapes": 0,
            "errors": 0
        }
        
        # Extract theme data once for the entire presentation
        self.theme_colors = self._extract_theme_colors(prs)
        self.theme_fonts = self._extract_theme_fonts(getattr(self, '_file_path', None))
        logger.debug(f"Extracted theme colors: {self.theme_colors}")
        logger.debug(f"Extracted theme fonts: {getattr(self, 'theme_fonts', {})}")
        
        # Convert slide dimensions from EMU to pixels (96 DPI)
        width = self._emu_to_px(prs.slide_width)
        height = self._emu_to_px(prs.slide_height)
        
        # Normalize to our standard 1920x1080 if needed
        scale_x = 1920 / width if width > 0 else 1
        scale_y = 1080 / height if height > 0 else 1
        # Use uniform scaling to maintain aspect ratio
        scale = min(scale_x, scale_y)
        
        slides = []
        for idx, slide in enumerate(prs.slides):
            try:
                if idx == 0:
                    logger.info(f"=== PROCESSING SLIDE 1 with {len(slide.shapes)} shapes ===")
                processed_slide = await self._process_slide(slide, idx, scale)
                slides.append(processed_slide)
            except Exception as e:
                logger.error(f"Error processing slide {idx}: {e}")
                self.stats["errors"] += 1
                # Add empty slide on error
                slides.append({
                    "id": str(uuid.uuid4()),
                    "title": f"Slide {idx + 1} (Error)",
                    "components": []
                })
                
        deck = {
            "uuid": str(uuid.uuid4()),
            "name": "Imported Presentation",
            "slides": slides,
            "size": {"width": 1920, "height": 1080},
            "metadata": {
                "source": "pptx",
                "original_size": {"width": width, "height": height},
                "import_stats": self.stats,
                "theme_colors": self.theme_colors,
                "theme_fonts": getattr(self, 'theme_fonts', {}),
                "slide_count": len(prs.slides)
            }
        }
        
        logger.info(f"Import complete. Stats: {self.stats}")
        return deck
    def _extract_theme_fonts(self, file_path: Optional[str]) -> Dict[str, Any]:
        """Extract theme font scheme (major/minor fonts) from theme XML if available."""
        fonts: Dict[str, Any] = {}
        if not file_path:
            return fonts
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            with zipfile.ZipFile(file_path, 'r') as pptx_zip:
                theme_files = [f for f in pptx_zip.namelist() if 'theme/theme' in f and f.endswith('.xml')]
                if not theme_files:
                    return fonts
                with pptx_zip.open(theme_files[0]) as theme_xml:
                    root = ET.parse(theme_xml).getroot()
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    font_scheme = root.find('.//a:fontScheme', ns)
                    if font_scheme is not None:
                        major = font_scheme.find('.//a:majorFont', ns)
                        minor = font_scheme.find('.//a:minorFont', ns)
                        def _font_from(elem):
                            if elem is None:
                                return None
                            latin = elem.find('./a:latin', ns)
                            if latin is not None and latin.get('typeface'):
                                return latin.get('typeface')
                            return None
                        fonts['major'] = _font_from(major)
                        fonts['minor'] = _font_from(minor)
        except Exception as e:
            logger.debug(f"Failed to extract theme fonts: {e}")
        return fonts
        
    async def _process_slide(self, slide, idx: int, scale: float) -> Dict[str, Any]:
        """Process a single slide"""
        components = []
        background_shape_id = None  # Track if we found a shape-based background
        
        # Extract slide title
        title = self._extract_slide_title(slide) or f"Slide {idx + 1}"
        
        # Process background
        background_data = self._extract_background(slide, scale)
        if background_data:
            background, shape_id = background_data if isinstance(background_data, tuple) else (background_data, None)
            components.append(background)
            background_shape_id = shape_id
            
        # Process all shapes
        z_index = 1
        for shape in slide.shapes:
            try:
                # Skip shape if it was used as background
                if background_shape_id and hasattr(shape, '_element') and shape._element is background_shape_id:
                    continue
                    
                component = self._process_shape(shape, scale, z_index)
                if component:
                    # Handle groups that return lists of components
                    if isinstance(component, list):
                        components.extend(component)
                        z_index += len(component)
                        self.stats["components"] += len(component)
                    else:
                        components.append(component)
                        z_index += 1
                        self.stats["components"] += 1
            except Exception as e:
                logger.warning(f"Failed to process shape: {e}")
                self.stats["errors"] += 1
        # Enforce a stable render ordering when frontend ignores zIndex
        # Order: Background -> Image -> Shape -> Text -> everything else
        def _order_key(c: Dict[str, Any]) -> int:
            t = c.get("type")
            if t == "Background":
                return 0
            if t == "Image":
                return 1
            if t == "Shape":
                return 2
            if t == "TiptapTextBlock":
                return 3
            return 4

        components.sort(key=_order_key)

        # Fallback: if no Image component exists but Background has an image, also emit an Image
        try:
            has_image = any(c.get("type") == "Image" for c in components)
            if not has_image:
                for c in components:
                    if c.get("type") == "Background":
                        p = c.get("props", {})
                        if p.get("backgroundType") == "image" and p.get("backgroundImageUrl"):
                            components.insert(1, {
                                "id": str(uuid.uuid4()),
                                "type": "Image",
                                "props": {
                                    "position": {"x": 0, "y": 0},
                                    "width": 1920,
                                    "height": 1080,
                                    "opacity": 1,
                                    "rotation": 0,
                                    "zIndex": 1,
                                    "src": p.get("backgroundImageUrl"),
                                    "objectFit": "cover"
                                }
                            })
                            break
        except Exception:
            pass

        return {
            "id": str(uuid.uuid4()),
            "title": title,
            "components": components
        }
        
    def _process_shape(self, shape, scale: float, z_index: int) -> Optional[Union[Dict[str, Any], List[Dict[str, Any]]]]:
        """Process a single shape into a component"""
        # Get common properties
        bounds = self._get_shape_bounds(shape, scale)
        
        # Route to specific handlers based on shape type
        # 0) Picture placeholder (Google Slides often exports images as PLACEHOLDER with .image)
        try:
            if hasattr(shape, 'image') and getattr(shape, 'image') is not None:
                return self._create_image_component(shape, bounds, z_index)
        except Exception:
            pass
        # 0b) Shapes with picture fill (blipFill) should become masked Shape with image fill
        try:
            blob = self._try_extract_picture_fill_blob(shape)
            if blob:
                # Extract crop from a:srcRect if present
                crop_rect = self._extract_blipfill_crop_rect(shape)
                # Detect flips from xfrm
                flip_x = False
                flip_y = False
                try:
                    el = getattr(shape, '_element', None)
                    if el is not None and hasattr(el, 'xfrm') and el.xfrm is not None:
                        flip_x = bool(getattr(el.xfrm, 'flipH', False))
                        flip_y = bool(getattr(el.xfrm, 'flipV', False))
                except Exception:
                    pass
                # Prefer preserving original shape mask; build a Shape with embedded image fill
                return self._create_shape_with_picture_image(blob, shape, bounds, z_index, crop_rect, flip_x, flip_y)
        except Exception:
            pass
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._create_image_component(shape, bounds, z_index)
        # Recognize charts
        if hasattr(shape, 'chart') and getattr(shape, 'chart') is not None:
            comp = self._create_chart_component(shape, bounds, z_index)
            if comp:
                return comp
        # Recognize tables
        if getattr(shape, 'has_table', False):
            comp = self._create_table_component(shape, bounds, z_index)
            if comp:
                return comp
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Process groups by creating a Group container + children (schema-compliant)
            return self._create_group_with_children(shape, scale, z_index)
        elif (
            shape.shape_type == MSO_SHAPE_TYPE.LINE
            or (hasattr(MSO_SHAPE_TYPE, 'CONNECTOR') and shape.shape_type == getattr(MSO_SHAPE_TYPE, 'CONNECTOR'))
        ):
            # Lines and connectors
            return self._create_line_component(shape, bounds, z_index)
        elif shape.has_text_frame and shape.text_frame and shape.text_frame.text.strip():
            # Import placeholders (title, subtitle, etc.) as text blocks instead of skipping
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                return self._create_text_component(shape, bounds, z_index, scale)
            # Check if this is a shape with text (not just a text box)
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # This is a shape with text - create ShapeWithText component
                return self._create_shape_with_text_component(shape, bounds, z_index, scale)
            else:
                # This is a text box or text placeholder - create TiptapTextBlock
                return self._create_text_component(shape, bounds, z_index, scale)
        else:
            # All other shapes (AUTO_SHAPE, FREEFORM, CONNECTOR, PLACEHOLDER without text, etc.)
            return self._create_shape_component(shape, bounds, z_index)
            
    def _create_text_component(self, shape, bounds: Dict[str, float], z_index: int, scale: float = 1.0) -> Dict[str, Any]:
        """Create a text component from a shape with text"""
        self.stats["text_blocks"] += 1
        logger.debug(f"Creating text component for shape: {shape.name}")
        
        # Extract text segments with formatting
        segments = []
        paragraphs_segments: List[List[Dict[str, Any]]] = []
        alignment = "left"
        line_height = 1.2  # More natural default
        
        # Get text frame properties
        text_frame = shape.text_frame
        
        # Extract actual padding from text frame margins
        # Set padding to 0 to prevent text cropping in imports
        padding = 0
        
        for paragraph in text_frame.paragraphs:
            # Get paragraph alignment
            if paragraph.alignment == PP_ALIGN.CENTER:
                alignment = "center"
            elif paragraph.alignment == PP_ALIGN.RIGHT:
                alignment = "right"
            elif paragraph.alignment == PP_ALIGN.JUSTIFY:
                alignment = "justify"
                
            # Get line spacing from first paragraph
            if hasattr(paragraph, 'line_spacing') and paragraph.line_spacing is not None:
                line_height = float(paragraph.line_spacing)
                
            para_runs: List[Dict[str, Any]] = []
            for run in paragraph.runs:
                if not run.text:
                    continue
                    
                # Extract font properties safely
                font = run.font
                # For Google Slides exports, we might want to map fonts
                # If the presentation uses custom Google Fonts, they'll be lost in PPTX
                font_name = "Calibri"  # Default Office font
                
                # Optional: Add font mapping for known Google Slides templates
                # e.g., if this is a known template that uses Newsreader:
                # font_name = "Newsreader"
                font_size = 12  # Default 12pt
                # Default to black text (as defined in slide masters)
                font_color = "#000000ff"  # Black text by default
                
                # Check for layout-defined color first (title/body placeholders)
                shape = getattr(text_frame, '_parent', None)
                if shape:
                    layout_color = self._get_layout_placeholder_color(shape)
                    if layout_color:
                        if layout_color.startswith('#'):
                            font_color = layout_color
                        else:
                            # layout_color may be theme key (e.g., 'dk1', 'lt1', 'accent1')
                            color_val = self.theme_colors.get(layout_color)
                            if color_val:
                                font_color = color_val
                
                # Try to get font name from multiple sources
                try:
                    if font.name:
                        font_name = font.name
                    elif hasattr(paragraph, 'font') and paragraph.font and paragraph.font.name:
                        font_name = paragraph.font.name
                    elif hasattr(text_frame, 'font') and text_frame.font and text_frame.font.name:
                        font_name = text_frame.font.name
                    else:
                        # First try to get from layout placeholder
                        shape = getattr(text_frame, '_parent', None)
                        if shape:
                            layout_font = self._get_layout_placeholder_style(shape)
                            if layout_font:
                                font_name = layout_font
                            else:
                                # Theme font fallback based on placeholder type
                                try:
                                    ph_type = shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else None
                                except Exception:
                                    ph_type = None
                                theme_major = getattr(self, 'theme_fonts', {}).get('major')
                                theme_minor = getattr(self, 'theme_fonts', {}).get('minor')
                                if ph_type in [1, 3] and theme_major:
                                    font_name = theme_major
                                elif theme_minor:
                                    font_name = theme_minor
                        else:
                            # Fall back to default extraction
                            extracted_font = self._extract_default_font(text_frame, paragraph)
                            if extracted_font:
                                font_name = extracted_font
                except:
                    pass
                    
                # Try to get font size from multiple sources
                try:
                    if font.size is not None:
                        font_size = self._pt_to_pt(font.size)
                        logger.debug(f"Font size from run: {font.size.pt}pt -> {font_size}pt")
                    elif hasattr(paragraph, 'font') and paragraph.font and paragraph.font.size is not None:
                        font_size = self._pt_to_pt(paragraph.font.size)
                        logger.debug(f"Font size from paragraph: {paragraph.font.size.pt}pt -> {font_size}pt")
                    elif hasattr(text_frame, 'font') and text_frame.font and text_frame.font.size is not None:
                        font_size = self._pt_to_pt(text_frame.font.size)
                        logger.debug(f"Font size from text_frame: {text_frame.font.size.pt}pt -> {font_size}pt")
                    else:
                        # First try to get from layout placeholder
                        shape = getattr(text_frame, '_parent', None)
                        if shape:
                            logger.debug(f"Checking layout for shape: {getattr(shape, 'name', 'unknown')}")
                            layout_size = self._get_layout_placeholder_size(shape)
                            if layout_size:
                                # Convert layout size from px to points
                                font_size = layout_size * 72 / 96  # Convert px to pt
                                logger.debug(f"Font size from layout placeholder: {font_size}pt")

                            # For placeholders without explicit layout size, use defaults based on type
                            elif shape.is_placeholder:
                                logger.debug(f"Using default for placeholder type: {shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else 'unknown'}")
                                # Title placeholders are typically larger
                                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type in [1, 3]:  # TITLE
                                    font_size = 42  # 42pt (from layout analysis)
                                else:
                                    font_size = 14  # 14pt default
                        # Try to infer from paragraph level
                        elif hasattr(paragraph, 'level') and paragraph.level is not None:
                            # Common size hierarchy for levels
                            level_sizes = {0: 28, 1: 24, 2: 20, 3: 18, 4: 16}
                            font_size = level_sizes.get(paragraph.level, 16)
                    
                    # Apply minimum but preserve small sizes where reasonable  
                    # Apply scale factor to font size
                    font_size = font_size * scale
                    
                    if font_size < 6:
                        font_size = 6  # 6pt minimum

                except:
                    pass
                    
                # Get color safely - override layout color when run specifies a color (RGB or theme)
                try:
                    if getattr(font, 'color', None) is not None:
                        col = font.color
                        # Use generic extractor to handle RGB, theme indices, and transforms
                        # Some PPTX use scheme colors (no rgb) that we must honor
                        if getattr(col, 'type', None) is not None or getattr(col, 'rgb', None) is not None or getattr(col, 'theme_color', None) is not None:
                            extracted_color = self._get_color_hex(col, is_text_color=True, theme_colors=self.theme_colors)
                            if extracted_color:
                                font_color = extracted_color
                except Exception:
                    pass
                
                # Extract text style - only inline properties
                logger.debug(f"Final font color for text: {font_color}")
                style = {
                    "textColor": font_color,
                    "backgroundColor": "#00000000",
                    "bold": bool(font.bold) if font.bold is not None else False,
                    "italic": bool(font.italic) if font.italic is not None else False,
                    "underline": bool(font.underline) if font.underline is not None else False,
                    "strike": False,  # Not directly available in python-pptx
                    # Include font properties in style for proper rendering
                    "fontSize": font_size,
                    "fontFamily": font_name
                }
                
                para_runs.append({
                    "text": run.text,
                    "style": style
                })

            if para_runs:
                paragraphs_segments.append(para_runs)
                segments.extend(para_runs)
                
        if not segments:
            return None
            
        # Extract the most common font properties from segments for block-level styling
        font_family = "Calibri"  # Default Office font
        font_size = 16
        if segments:
            # Use the first segment's extracted font info
            # These were extracted above but not stored - need to re-extract
            first_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
            if first_paragraph and first_paragraph.runs:
                first_run = first_paragraph.runs[0]
                font = first_run.font
                
                # Get font family
                if font.name:
                    font_family = font.name
                elif hasattr(first_paragraph, 'font') and first_paragraph.font and first_paragraph.font.name:
                    font_family = first_paragraph.font.name
                else:
                    # Try to get from default properties
                    extracted_font = self._extract_default_font(text_frame, first_paragraph)
                    if extracted_font:
                        font_family = extracted_font
                    
                # Get font size
                if font.size is not None:
                    font_size = self._pt_to_pt(font.size)
                    logger.debug(f"Font size from connector run: {font.size.pt}pt -> {font_size}pt")
                elif hasattr(first_paragraph, 'font') and first_paragraph.font and first_paragraph.font.size is not None:
                    font_size = self._pt_to_pt(first_paragraph.font.size)
                    logger.debug(f"Font size from connector paragraph: {first_paragraph.font.size.pt}pt -> {font_size}pt")
                else:
                    # Check placeholder defaults
                    if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.type in [1, 3]:  # TITLE
                            font_size = 24  # 24pt for title connectors
                        else:
                            font_size = 14  # 14pt default
                
                # Allow smaller font sizes - only enforce a very small minimum
                if font_size < 8:
                    logger.warning(f"Connector font size {font_size}pt is too small, setting to minimum 6pt")
                    font_size = 6
        
        # Create Tiptap content structure from segments preserving paragraphs
        tiptap_content = {
            "type": "doc",
            "content": []
        }
        
        if paragraphs_segments:
            for para in paragraphs_segments:
                paragraph_node = {
                    "type": "paragraph",
                    "content": [],
                    "attrs": {
                        "textAlign": alignment,
                        "lineHeight": line_height
                    }
                }
            
                for segment in para:
                    text_node = {
                        "type": "text",
                        "text": segment["text"]
                    }
                    
                    marks = []
                    style = segment.get("style", {})
                    if style.get("bold"):
                        marks.append({"type": "bold"})
                    if style.get("italic"):
                        marks.append({"type": "italic"})
                    if style.get("underline"):
                        marks.append({"type": "underline"})
                    if style.get("strike"):
                        marks.append({"type": "strike"})
                    
                    text_style_mark = {
                        "type": "textStyle",
                        "attrs": {
                            "color": style.get("textColor", "#000000ff"),
                            "fontSize": style.get("fontSize", font_size),
                            "fontFamily": style.get("fontFamily", font_family)
                        }
                    }
                    marks.append(text_style_mark)
                    
                    if marks:
                        text_node["marks"] = marks
                    
                    paragraph_node["content"].append(text_node)
            
                tiptap_content["content"].append(paragraph_node)
        
        # Use font properties from the first segment for the component
        component_font_family = segments[0]["style"]["fontFamily"] if segments else "Arial"
        component_font_size = segments[0]["style"]["fontSize"] if segments else 16
        
        return {
            "id": str(uuid.uuid4()),
            "type": "TiptapTextBlock",
            "props": {
                **bounds,
                # Compatibility: duplicate x/y at top-level of props
                "x": int(bounds["position"]["x"]),
                "y": int(bounds["position"]["y"]),
                "texts": segments,
                "content": tiptap_content,  # Add Tiptap content structure
                "fontFamily": component_font_family,
                "fontSize": component_font_size,
                "fontSizePx": int(round(component_font_size * 96 / 72)) if isinstance(component_font_size, (int, float)) else 48,
                "fontWeight": "normal",
                "fontStyle": "normal",
                "textColor": segments[0]["style"]["textColor"] if segments else "#000000ff",
                "backgroundColor": "#00000000",
                "letterSpacing": 0,
                "lineHeight": line_height,
                                    "alignment": alignment,
                    "verticalAlignment": ("middle" if getattr(text_frame, 'vertical_anchor', None) in (getattr(text_frame, 'ANCHOR_MIDDLE', object()), getattr(text_frame, 'ANCHOR_CENTER', object())) else "top"),
                    "padding": padding,  # Use extracted padding from text frame
                    "textPadding": padding,  # Frontend schema name
                    "opacity": 1,
                    "zIndex": z_index
            }
        }
        
    def _create_image_component(self, shape, bounds: Dict[str, float], z_index: int) -> Optional[Dict[str, Any]]:
        """Create an image component from a picture shape"""
        self.stats["images"] += 1
        logger.debug(f"Extracting image from shape: {shape.name}")
        
        try:
            # Extract image data
            image = shape.image
            image_bytes = image.blob
            logger.debug(f"Image size: {len(image_bytes)} bytes")
            
            # Determine format from bytes (more reliable than extension)
            mime_type = self._detect_image_mime_from_bytes(image_bytes)
            
            # Convert to base64 data URL
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            data_url = f"data:{mime_type};base64,{image_base64}"
            
            # Check for cropping and flips
            crop_props = {}
            crop_rect = None
            try:
                cl = getattr(shape, 'crop_left', None)
                cr = getattr(shape, 'crop_right', None)
                ct = getattr(shape, 'crop_top', None)
                cb = getattr(shape, 'crop_bottom', None)
                if any(v not in (None, 0) for v in [cl, cr, ct, cb]):
                    crop_rect = {
                        "left": float(cl or 0),
                        "right": float(cr or 0),
                        "top": float(ct or 0),
                        "bottom": float(cb or 0),
                    }
                    # When cropping is present, cover usually matches authoring intent
                crop_props["objectFit"] = "cover"
            except Exception:
                pass

            if "objectFit" not in crop_props:
                crop_props["objectFit"] = "contain"

            flip_x = False
            flip_y = False
            try:
                el = getattr(shape, '_element', None)
                if el is not None and hasattr(el, 'xfrm') and el.xfrm is not None:
                    xfrm = el.xfrm
                    # flipH / flipV may be present on transform
                    flip_x = bool(getattr(xfrm, 'flipH', False))
                    flip_y = bool(getattr(xfrm, 'flipV', False))
            except Exception:
                pass
            
            return {
                "id": str(uuid.uuid4()),
                "type": "Image",
                "props": {
                    **bounds,
                    "src": data_url,
                    "alt": shape.name or "Image",
                    "objectFit": crop_props.get("objectFit", "contain"),
                    **({"cropRect": crop_rect} if crop_rect else {}),
                    "borderRadius": 0,
                    "borderWidth": 0,
                    "borderColor": "#000000ff",
                    "shadow": False,
                    "shadowBlur": 10,
                    "shadowColor": "#0000004D",
                    "shadowOffsetX": 0,
                    "shadowOffsetY": 4,
                    "shadowSpread": 0,
                    "opacity": 1,
                    "zIndex": z_index,
                    # Filter defaults
                    "filterPreset": "none",
                    "brightness": 100,
                    "contrast": 100,
                    "saturation": 100,
                    "grayscale": 0,
                    "sepia": 0,
                    "hueRotate": 0,
                    "blur": 0,
                    "invert": 0,
                    **({"flipX": True} if flip_x else {}),
                    **({"flipY": True} if flip_y else {}),
                }
            }
        except Exception as e:
            logger.warning(f"Failed to extract image from {shape.name}: {e}")
            self.stats["images"] -= 1  # Don't count failed images
            return None
            
    def _create_shape_component(self, shape, bounds: Dict[str, float], z_index: int) -> Dict[str, Any]:
        """Create a shape component"""
        self.stats["shapes"] += 1
        
        # Get the auto shape type safely
        shape_type = MSO_SHAPE.RECTANGLE  # Default
        
        # Try to get auto shape type for AUTO_SHAPE types
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                shape_type = shape.auto_shape_type
            except Exception as e:
                pass
        # For PLACEHOLDER shapes, try to infer from placeholder type
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            # Placeholders are usually rectangles unless specified otherwise
            shape_type = MSO_SHAPE.RECTANGLE
        # For FREEFORM shapes, try to detect the actual shape
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            # Google Slides exports shapes as FREEFORM, try to detect circle/ellipse by path
            shape_type = self._detect_freeform_shape_type(shape)
        
        # Shape mapping to our supported types: rectangle, circle, ellipse, triangle, star, hexagon, pentagon, diamond, arrow, heart
        shape_type_map = {
            MSO_SHAPE.RECTANGLE: "rectangle",
            MSO_SHAPE.ROUNDED_RECTANGLE: "rectangle",
            MSO_SHAPE.ROUND_2_SAME_RECTANGLE: "rectangle",
            MSO_SHAPE.OVAL: "ellipse",  # Will check for circle below
            MSO_SHAPE.DIAMOND: "diamond",
            MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
            MSO_SHAPE.RIGHT_TRIANGLE: "triangle",
            MSO_SHAPE.HEXAGON: "hexagon",
            MSO_SHAPE.REGULAR_PENTAGON: "pentagon",
            MSO_SHAPE.STAR_4_POINT: "star",
            MSO_SHAPE.STAR_5_POINT: "star",
            MSO_SHAPE.STAR_6_POINT: "star",
            MSO_SHAPE.STAR_7_POINT: "star",
            MSO_SHAPE.STAR_8_POINT: "star",
            MSO_SHAPE.STAR_10_POINT: "star",
            MSO_SHAPE.STAR_12_POINT: "star",
            MSO_SHAPE.STAR_16_POINT: "star",
            MSO_SHAPE.STAR_24_POINT: "star",
            MSO_SHAPE.STAR_32_POINT: "star",
            MSO_SHAPE.HEART: "heart",
            MSO_SHAPE.RIGHT_ARROW: "arrow",
            MSO_SHAPE.LEFT_ARROW: "arrow",
            MSO_SHAPE.UP_ARROW: "arrow",
            MSO_SHAPE.DOWN_ARROW: "arrow",
            MSO_SHAPE.LEFT_RIGHT_ARROW: "arrow",
            MSO_SHAPE.UP_DOWN_ARROW: "arrow",
            # Map all other shapes to rectangle as fallback
            MSO_SHAPE.ROUND_1_RECTANGLE: "rectangle",
            MSO_SHAPE.ROUND_2_SAME_RECTANGLE: "rectangle",  # Rounded rectangle
            MSO_SHAPE.ROUND_2_DIAG_RECTANGLE: "rectangle",
            MSO_SHAPE.SNIP_1_RECTANGLE: "rectangle",
            MSO_SHAPE.SNIP_2_SAME_RECTANGLE: "rectangle",
            MSO_SHAPE.SNIP_2_DIAG_RECTANGLE: "rectangle",
            MSO_SHAPE.FRAME: "rectangle",
            MSO_SHAPE.HALF_FRAME: "rectangle",
            MSO_SHAPE.TEAR: "ellipse",
            MSO_SHAPE.CLOUD: "ellipse",
            MSO_SHAPE.LIGHTNING_BOLT: "triangle",
            MSO_SHAPE.SUN: "star",
            MSO_SHAPE.MOON: "circle",
            MSO_SHAPE.ARC: "ellipse",
            MSO_SHAPE.DONUT: "circle",
            MSO_SHAPE.NO_SYMBOL: "circle",
            MSO_SHAPE.BLOCK_ARC: "ellipse",
            MSO_SHAPE.SMILEY_FACE: "circle",
            MSO_SHAPE.CUBE: "rectangle",
            MSO_SHAPE.BEVEL: "rectangle",
            MSO_SHAPE.FOLDED_CORNER: "rectangle",
            MSO_SHAPE.ACTION_BUTTON_HOME: "rectangle",
            MSO_SHAPE.ACTION_BUTTON_HELP: "circle",
            MSO_SHAPE.ACTION_BUTTON_INFORMATION: "circle",
            # MSO_SHAPE.CIRCLE_ARROW: "circle",  # Not available in python-pptx
            # MSO_SHAPE.EXPLOSION_1: "star",  # Not available in python-pptx
            # MSO_SHAPE.EXPLOSION_2: "star",  # Not available in python-pptx
            # Additional shapes
            MSO_SHAPE.PARALLELOGRAM: "rectangle",
            MSO_SHAPE.TRAPEZOID: "rectangle",
            MSO_SHAPE.OCTAGON: "rectangle",
        }
        
        our_shape_type = shape_type_map.get(shape_type, "rectangle")
        
        # Check if an oval is actually a circle (equal width and height)
        if shape_type == MSO_SHAPE.OVAL and hasattr(shape, 'width') and hasattr(shape, 'height'):
            # Compare dimensions with tolerance (5%) for near-circles
            width = shape.width
            height = shape.height
            if abs(width - height) <= min(width, height) * 0.05:
                our_shape_type = "circle"
        
        # Extract fill color safely and detect picture fill
        fill_color = "#00000000"  # Default transparent
        try:
            if hasattr(shape, 'fill') and shape.fill is not None:
                fill_type = shape.fill.type

                # Solid fill
                if fill_type == 1:
                    try:
                        if shape.fill.fore_color:
                            fill_color = self._get_color_hex(shape.fill.fore_color, theme_colors=self.theme_colors)
                    except Exception:
                        pass
                # No fill
                elif fill_type == 0 or fill_type is None:
                    fill_color = "#00000000"
                else:
                    # Picture/gradient/pattern fills
                    # If picture fill, extract the embedded image and return an Image component
                    picture_blob = self._try_extract_picture_fill_blob(shape)
                    if picture_blob is not None:
                        try:
                            # Build image data URL
                            image_base64 = base64.b64encode(picture_blob).decode('utf-8')
                            data_url = f"data:image/png;base64,{image_base64}"
                            # If the original shape was a circle, apply circular mask via borderRadius
                            border_radius = 0
                            if our_shape_type == "circle":
                                try:
                                    border_radius = int(round(min(bounds["width"], bounds["height"]) / 2))
                                except Exception:
                                    border_radius = 0
                            return {
                                "id": str(uuid.uuid4()),
                                "type": "Image",
                                "props": {
                                    **bounds,
                                    "src": data_url,
                                    "alt": shape.name or "Image",
                                    "objectFit": "cover",
                                    "borderRadius": border_radius,
                                    "borderWidth": 0,
                                    "borderColor": "#000000ff",
                                    "opacity": 1,
                                    "zIndex": z_index,
                                }
                            }
                        except Exception:
                            # Fall back to transparent fill if anything fails
                            fill_color = "#00000000"
                    # For other non-solid fills, keep transparent and rely on stroke if any
        except Exception as e:
            pass
            
        # Don't render completely invisible shapes (but allow transparent shapes with borders)
        # Comment out for now to see all shapes
        # if fill_color == "#00000000" and stroke_width == 0:
        #     return None
            
        # Extract line properties safely
        stroke_color = "#00000000"  # Default to transparent
        stroke_width = 0
        try:
            if hasattr(shape, 'line') and shape.line:
                # Check if line has any properties
                if hasattr(shape.line, 'width') and shape.line.width is not None:
                    # Use exact conversion for line widths (0.8pt -> 1px)
                    stroke_width = self._pt_to_px(shape.line.width)
                    if stroke_width == 0:  # Ensure minimum visibility
                        stroke_width = 1
                # Check line color/fill
                if hasattr(shape.line, 'fill') and shape.line.fill:
                    if shape.line.fill.type == 1:  # Solid fill
                        stroke_color = self._get_color_hex(shape.line.fill.fore_color, theme_colors=self.theme_colors)
                elif hasattr(shape.line, 'color') and shape.line.color:
                    stroke_color = self._get_color_hex(shape.line.color, theme_colors=self.theme_colors)
                # If no explicit stroke width/color, keep stroke transparent/0
        except Exception as e:
            pass
        
        # Determine border radius
        border_radius = 0
        if our_shape_type == "circle":
            try:
                border_radius = int(round(min(bounds["width"], bounds["height"]) / 2))
            except Exception:
                border_radius = 0
        elif shape_type in [MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.ROUND_1_RECTANGLE, 
                            MSO_SHAPE.ROUND_2_SAME_RECTANGLE, MSO_SHAPE.ROUND_2_DIAG_RECTANGLE]:
            border_radius = 8
                
        # Do NOT enforce minimum sizes; preserve exact authoring sizes
        
        # Skip placeholder full-bleed transparent rectangles at origin (often artifact when an Image exists)
        try:
            if (fill_color == "#00000000" and (stroke_width or 0) == 0 and 
                int(bounds["position"]["x"]) == 0 and int(bounds["position"]["y"]) == 0 and 
                int(bounds["width"]) >= 1918 and int(bounds["height"]) >= 1078):
                return None
        except Exception:
            pass
        
        # Build minimal schema-compliant props (omit optional fields entirely when not used)
        props: Dict[str, Any] = {
            "position": {"x": int(bounds["position"]["x"]), "y": int(bounds["position"]["y"])},
            "width": int(bounds["width"]),
            "height": int(bounds["height"]),
            "opacity": 1,
            "rotation": int(bounds.get("rotation", 0)),
            "zIndex": int(z_index),
            "shapeType": our_shape_type,
            # Provide both keys for broader renderer/frontend compatibility
            "shape": our_shape_type,
            "fill": fill_color,
            "backgroundColor": fill_color,
        }
        if border_radius > 0:
            props["borderRadius"] = int(border_radius)
        if stroke_width and stroke_width > 0:
            props["strokeWidth"] = int(stroke_width)
            props["stroke"] = stroke_color
            props["borderWidth"] = int(stroke_width)
            props["borderColor"] = stroke_color

        # Compatibility: duplicate x/y at top-level for Shape props
        props["x"] = props["position"]["x"]
        props["y"] = props["position"]["y"]

        result = {
            "id": str(uuid.uuid4()),
            "type": "Shape",
            "props": props,
        }
        return result
    
    def _detect_freeform_shape_type(self, shape):
        """Try to detect if a freeform shape is actually a circle/ellipse based on its path or properties"""
        try:
            # First check if it has equal width/height (likely a circle)
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                width = shape.width
                height = shape.height
                # Within 5% tolerance
                if abs(width - height) <= min(width, height) * 0.05:
                    # Check if it has a curved path (not a square)
                    # For now, assume it's a circle if dimensions are equal
                    # TODO: Could analyze the actual path data if available
                    return MSO_SHAPE.OVAL
            
            # Try to access the path data
            if hasattr(shape, 'path'):
                path_str = str(shape.path)
                # Simple heuristics - if path contains curves (C commands) it might be circular
                if 'C' in path_str or 'c' in path_str:
                    return MSO_SHAPE.OVAL
            
            # Check element tree for hints
            if hasattr(shape, '_element') and shape._element is not None:
                # Look for geometry presets in the XML
                from lxml import etree
                # Try to find any shape geometry hints
                for elem in shape._element.iter():
                    if elem.tag.endswith('prstGeom'):
                        prst = elem.get('prst')
                        if prst:
                            if prst in ('ellipse', 'circle'):
                                return MSO_SHAPE.OVAL
                            elif prst == 'rect':
                                return MSO_SHAPE.RECTANGLE
                
        except Exception as e:
            pass
        
        # Default to rectangle if we can't determine
        return MSO_SHAPE.RECTANGLE
        
    def _create_shape_with_text_component(self, shape, bounds: Dict[str, float], z_index: int, scale: float = 1.0) -> Dict[str, Any]:
        """Create a Shape component with inline texts for shapes that contain text"""
        self.stats["shapes"] += 1
        
        # First get the base shape properties
        shape_component = self._create_shape_component(shape, bounds, z_index)
        
        # Extract text content as inline segments compatible with Shape schema
        texts: List[Dict[str, Any]] = []
        root_text_color: Optional[str] = None
        text_frame = shape.text_frame
        
        for paragraph in getattr(text_frame, 'paragraphs', []) or []:
            for run in getattr(paragraph, 'runs', []) or []:
                if not run.text:
                    continue
                    
                font = run.font
                # Defaults
                font_color = "#000000ff"  # Default to black text
                
                # Check for layout-defined color
                layout_color = self._get_layout_placeholder_color(shape)
                if layout_color:
                    if layout_color.startswith('#'):
                        # RGB color
                        font_color = layout_color
                    else:
                        # Scheme color name (e.g., 'dk1')
                        # Map to theme color
                        color_val = self.theme_colors.get(layout_color)
                        if color_val:
                            font_color = color_val
                try:
                    if hasattr(font, 'color') and font.color is not None:
                        extracted_color = self._get_color_hex(font.color, is_text_color=True, theme_colors=self.theme_colors)
                        if extracted_color:
                            font_color = extracted_color
                except Exception:
                    pass
                if not root_text_color:
                    root_text_color = font_color

                # Extract font properties for this run
                # For Google Slides exports, we might want to map fonts
                # If the presentation uses custom Google Fonts, they'll be lost in PPTX
                font_name = "Calibri"  # Default Office font
                
                # Optional: Add font mapping for known Google Slides templates
                # e.g., if this is a known template that uses Newsreader:
                # font_name = "Newsreader"
                font_size = 12  # Default 12pt
                
                try:
                    if font.name:
                        font_name = font.name
                    else:
                        # First try to get from layout placeholder
                        layout_font = self._get_layout_placeholder_style(shape)
                        if layout_font:
                            font_name = layout_font
                        else:
                            # Try to get from default properties
                            extracted_font = self._extract_default_font(text_frame, paragraph)
                            if extracted_font:
                                font_name = extracted_font
                    
                    if font.size is not None:
                        font_size = self._pt_to_pt(font.size)
                    else:
                        # Try to get from layout placeholder
                        layout_size = self._get_layout_placeholder_size(shape)
                        if layout_size:
                            # Convert layout size from px to points
                            font_size = layout_size * 72 / 96  # Convert px to pt
                        elif shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                            # Use placeholder type to determine default size
                            ph_type = shape.placeholder_format.type
                            if ph_type in [1, 3]:  # TITLE or CENTER_TITLE
                                font_size = 42  # 42pt
                            elif ph_type == 2:  # SUBTITLE
                                font_size = 9  # 9pt  
                            else:
                                font_size = 10  # 10pt
                except:
                    pass
                    
                # Apply scale factor to font size
                font_size = font_size * scale
                    
                seg_style: Dict[str, Any] = {
                    "textColor": font_color,
                    "backgroundColor": "#00000000",
                    "bold": bool(font.bold) if font.bold is not None else False,
                    "italic": bool(font.italic) if font.italic is not None else False,
                    "underline": bool(font.underline) if font.underline is not None else False,
                    "strike": False,
                    # Additional keys for broader compatibility
                    "highlight": False,
                    "subscript": False,
                    "superscript": False,
                    "color": font_color,
                    "link": False,
                    # Add font properties
                    "fontSize": font_size,
                    "fontFamily": font_name,
                    "href": "",
                }
                # Best-effort font size extraction
                try:
                    if getattr(font, 'size', None):
                        seg_style["fontSize"] = self._pt_to_pt(font.size) * scale
                except Exception:
                    pass

                texts.append({
                    "text": run.text,
                    "style": seg_style,
                })

        # Merge texts into shape props
        shape_props = shape_component.get("props", {})
        if texts:
            shape_props.update({
                "hasText": True,
                "texts": texts,
                "textColor": root_text_color or "#000000ff",
            })
        
        return {
            "id": shape_component["id"],
            "type": "Shape",
            "props": shape_props,
        }

    def _create_table_component(self, shape, bounds: Dict[str, float], z_index: int) -> Optional[Dict[str, Any]]:
        """Create a Table component from a PPTX table shape"""
        try:
            table = getattr(shape, 'table', None)
            if table is None:
                return None
            data: List[List[str]] = []
            cell_styles: List[Dict[str, Any]] = []
            for row in table.rows:
                row_vals: List[str] = []
                row_style: List[Dict[str, Any]] = []
                for cell in row.cells:
                    try:
                        row_vals.append(cell.text or "")
                    except Exception:
                        row_vals.append("")
                    # Extract per-cell styles best-effort
                    c_style: Dict[str, Any] = {}
                    try:
                        # Background fill
                        if hasattr(cell, 'fill') and cell.fill is not None:
                            if getattr(cell.fill, 'type', None) == 1:
                                c_style['backgroundColor'] = self._get_color_hex(cell.fill.fore_color)
                        # Borders (python-pptx lacks full table border API; approximate via cell.text_frame.paragraphs)
                        # We keep defaults; renderer can use tableStyles
                        # Vertical alignment
                        va = getattr(cell, 'vertical_anchor', None)
                        if va is not None:
                            try:
                                # Map to top/middle/bottom if available
                                c_style['verticalAlignment'] = 'middle' if str(va).lower().endswith('middle') or str(va).lower().endswith('center') else ('bottom' if str(va).lower().endswith('bottom') else 'top')
                            except Exception:
                                pass
                        # Padding via margins if present
                        tf = getattr(cell, 'text_frame', None)
                        if tf is not None:
                            pad = 0
                            try:
                                # Use left margin as proxy; PPTX may not expose all
                                if getattr(tf, 'margin_left', None):
                                    pad = int(self._pt_to_px(tf.margin_left))
                            except Exception:
                                pass
                            if pad:
                                c_style['padding'] = pad
                    except Exception:
                        pass
                    row_style.append(c_style)
                data.append(row_vals)
                cell_styles.append({"cells": row_style})
            headers: List[str] = []
            show_header = False
            try:
                if hasattr(table, 'first_row') and table.first_row and len(data) > 0:
                    headers = data[0]
                    data = data[1:]
                    show_header = True
            except Exception:
                pass
            return {
                "id": str(uuid.uuid4()),
                "type": "Table",
                "props": {
                    "position": {"x": int(bounds["position"]["x"]), "y": int(bounds["position"]["y"])},
                    "width": int(bounds["width"]),
                    "height": int(bounds["height"]),
                    "opacity": 1,
                    "rotation": int(bounds.get("rotation", 0)),
                    "zIndex": int(z_index),
                    "textColor": "#000000ff",
                    **({"headers": headers} if headers else {}),
                    **({"showHeader": True} if show_header else {}),
                    "data": data,
                    "tableStyles": {
                        "fontFamily": "Inter",
                        "fontSize": 14,
                        "borderColor": "#e2e8f0",
                        "borderWidth": 1,
                        "cellPadding": 8,
                        "headerBackgroundColor": "#f8fafc",
                        "headerTextColor": "#334155",
                        "cellBackgroundColor": "#ffffff",
                        "textColor": "#334155",
                        "alignment": "left",
                        "alternatingRowColor": False,
                        "hoverEffect": False
                    },
                    "cellStyles": cell_styles
                }
            }
        except Exception:
            return None

    def _create_chart_component(self, shape, bounds: Dict[str, float], z_index: int) -> Optional[Dict[str, Any]]:
        """Create a Chart component placeholder from a PPTX chart shape"""
        try:
            chart = getattr(shape, 'chart', None)
            chart_type = None
            try:
                chart_type = str(getattr(chart, 'chart_type', '') or '')
            except Exception:
                chart_type = ""
            return {
                "id": str(uuid.uuid4()),
                "type": "Chart",
                "props": {
                    "position": {"x": int(bounds["position"]["x"]), "y": int(bounds["position"]["y"])},
                    "width": int(bounds["width"]),
                    "height": int(bounds["height"]),
                    "opacity": 1,
                    "rotation": int(bounds.get("rotation", 0)),
                    "zIndex": int(z_index),
                    **({"chartType": chart_type} if chart_type else {})
                }
            }
        except Exception:
            return None
        
    def _create_line_component(self, shape, bounds: Dict[str, float], z_index: int) -> Dict[str, Any]:
        """Create a line component"""
        self.stats["shapes"] += 1
        
        # Extract line color and width
        stroke_color = "#00000000"  # Default to transparent
        stroke_width = 2
        stroke_dash = "none"
        start_shape = "none"
        end_shape = "none"
        
        try:
            if hasattr(shape, 'line') and shape.line:
                if shape.line.width:
                    stroke_width = self._pt_to_px(shape.line.width)
                    if stroke_width == 0:
                        stroke_width = 1
                if hasattr(shape.line, 'fill') and shape.line.fill.type == 1:
                    stroke_color = self._get_color_hex(shape.line.fill.fore_color)
                elif hasattr(shape.line, 'color'):
                    stroke_color = self._get_color_hex(shape.line.color, theme_colors=self.theme_colors)
                # Map dash styles if available
                try:
                    dash_style = getattr(shape.line, 'dash_style', None)
                    if dash_style:
                        ds = str(dash_style).upper()
                        dash_map = {
                            'DASH': '6,6',
                            'DOT': '2,6',
                            'DASH_DOT': '6,4,2,4',
                            'LONG_DASH': '10,6',
                            'LONG_DASH_DOT': '10,6,2,6',
                        }
                        stroke_dash = dash_map.get(ds, 'none')
                except Exception:
                    pass
        except:
            pass
            
        # Arrowhead detection via underlying XML when available
        try:
            el = getattr(shape, '_element', None)
            if el is not None:
                from lxml import etree
                ns = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
                }
                ln = None
                # Try line props under shape
                for tag in ('.//a:ln', './/p:spPr/a:ln'):
                    ln = el.find(tag, namespaces=ns)
                    if ln is not None:
                        break
                if ln is not None:
                    head_end = ln.find('./a:headEnd', namespaces=ns)
                    tail_end = ln.find('./a:tailEnd', namespaces=ns)
                    def _arrow(val):
                        if val is None:
                            return 'none'
                        typ = val.get('type') or val.get('w')  # type attr indicates arrow type
                        return 'arrow' if typ and str(typ).lower() not in ('none',) else 'none'
                    start_shape = _arrow(head_end)
                    end_shape = _arrow(tail_end)
        except Exception:
            pass

        # Calculate start and end points from bounds
        x = int(bounds["position"]["x"])
        y = int(bounds["position"]["y"])
        width = int(bounds["width"])
        height = int(bounds["height"])
        
        # For lines, we need to determine the actual endpoints
        # PowerPoint lines are typically diagonal within their bounds
        startPoint = {"x": x, "y": y}
        endPoint = {"x": x + width, "y": y + height}
        
        # Handle rotation if present
        if bounds.get("rotation", 0) != 0:
            # For now, just use the bounds as-is
            # Proper rotation would require more complex math
            pass
            
        return {
            "id": str(uuid.uuid4()),
            "type": "Lines",
            "props": {
                "startPoint": startPoint,
                "endPoint": endPoint,
                "connectionType": "straight",
                "startShape": start_shape,
                "endShape": end_shape,
                "stroke": stroke_color,
                "strokeWidth": stroke_width,
                "strokeDasharray": stroke_dash,
                "opacity": 1,
                "rotation": 0,  # Lines handle their own rotation via endpoints
                "zIndex": z_index
            }
        }
        
    def _create_group_with_children(self, group_shape, scale: float, z_index: int) -> List[Dict[str, Any]]:
        """Create a Group container and child components with correct absolute group position and relative children."""
        results: List[Dict[str, Any]] = []
        try:
            group_id = str(uuid.uuid4())

            # Get group's absolute bounds on the slide
            gb = self._get_shape_bounds(group_shape, scale)
            gx = int(gb['position']['x'])
            gy = int(gb['position']['y'])
            gw = int(gb['width'])
            gh = int(gb['height'])

            # Attempt to extract group child scaling (when group is resized)
            # Default to 1.0 (no scaling) if unavailable
            sX, sY = 1.0, 1.0
            try:
                # Access underlying oxml to get group transform
                el = getattr(group_shape, '_element', None)
                if el is not None:
                    # Namespaces
                    from lxml import etree
                    ns = {
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    }
                    # grpSpPr/a:xfrm
                    xfrm = el.find('.//p:grpSpPr/a:xfrm', namespaces=ns)
                    if xfrm is not None:
                        ext = xfrm.find('./a:ext', namespaces=ns)
                        chExt = xfrm.find('./a:chExt', namespaces=ns)
                        if ext is not None and chExt is not None:
                            ext_cx = int(ext.get('cx') or '0')
                            ext_cy = int(ext.get('cy') or '0')
                            ch_cx = int(chExt.get('cx') or '0')
                            ch_cy = int(chExt.get('cy') or '0')
                            if ch_cx and ch_cy and ext_cx and ext_cy:
                                sX = float(ext_cx) / float(ch_cx)
                                sY = float(ext_cy) / float(ch_cy)
            except Exception:
                # If any error, keep sX/sY as 1.0
                sX, sY = 1.0, 1.0

            child_ids: List[str] = []

            if hasattr(group_shape, 'shapes'):
                child_z = z_index + 1
                for child in group_shape.shapes:
                    child_comp = self._process_shape(child, scale, child_z)
                    if not child_comp:
                        continue
                    if isinstance(child_comp, list):
                        # Nested group comes back with absolute positions; make it relative to this group
                        nested_group = child_comp[0]
                        np = nested_group.get('props', {})
                        npos = np.get('position', {})
                        # Convert nested group container position to ABSOLUTE by adding parent group offset
                        rx_rel, ry_rel = int(npos.get('x', 0)), int(npos.get('y', 0))
                        # Apply group child scaling if present
                        rx, ry = gx + int(round(rx_rel * sX)), gy + int(round(ry_rel * sY))
                        npos['x'], npos['y'] = rx, ry
                        np['position'] = npos
                        np['x'], np['y'] = rx, ry
                        # Scale nested group dimensions as well
                        try:
                            nw, nh = int(np.get('width', 0)), int(np.get('height', 0))
                            if nw and nh:
                                np['width'] = int(round(nw * sX))
                                np['height'] = int(round(nh * sY))
                        except Exception:
                            pass
                        nested_group['props'] = np
                        results.append(nested_group)
                        child_ids.append(nested_group['id'])
                        # Shift all nested children by the same parent group offset to make ABSOLUTE
                        for cc in child_comp[1:]:
                            ccp = cc.get('props', {})
                            cpos = ccp.get('position', {})
                            cx_rel, cy_rel = int(cpos.get('x', 0)), int(cpos.get('y', 0))
                            cx, cy = gx + int(round(cx_rel * sX)), gy + int(round(cy_rel * sY))
                            cpos['x'], cpos['y'] = cx, cy
                            ccp['position'] = cpos
                            ccp['x'], ccp['y'] = cx, cy
                            # Scale child dimensions
                            try:
                                cw, ch = int(ccp.get('width', 0)), int(ccp.get('height', 0))
                                if cw and ch:
                                    ccp['width'] = int(round(cw * sX))
                                    ccp['height'] = int(round(ch * sY))
                            except Exception:
                                pass
                            cc['props'] = ccp
                            results.append(cc)
                        child_z += len(child_comp)
                    else:
                        # Child shapes inside group are RELATIVE; convert to ABSOLUTE
                        cp = child_comp.get('props', {})
                        pos = cp.get('position', {})
                        rx_rel, ry_rel = int(pos.get('x', 0)), int(pos.get('y', 0))
                        rx, ry = gx + int(round(rx_rel * sX)), gy + int(round(ry_rel * sY))
                        pos['x'], pos['y'] = rx, ry
                        cp['position'] = pos
                        cp['x'], cp['y'] = rx, ry
                        # Scale child dimensions by group scaling
                        try:
                            cw, ch = int(cp.get('width', 0)), int(cp.get('height', 0))
                            if cw and ch:
                                cp['width'] = int(round(cw * sX))
                                cp['height'] = int(round(ch * sY))
                        except Exception:
                            pass
                        child_comp['props'] = cp
                        results.append(child_comp)
                        child_ids.append(child_comp['id'])
                        child_z += 1

            group_comp = {
                "id": group_id,
                "type": "Group",
                "props": {
                    "position": {"x": gx, "y": gy},
                    "x": gx,
                    "y": gy,
                    "width": gw,
                    "height": gh,
                    "opacity": 1,
                    "rotation": 0,
                    "zIndex": z_index,
                    "children": child_ids,
                    "locked": False,
                }
            }
            return [group_comp] + results
        except Exception as e:
            logger.warning(f"Failed to create group component: {e}")
            return []

    def _try_extract_picture_fill_blob_from_background(self, slide, fill) -> Optional[bytes]:
        """Extract picture data from a background fill object"""
        try:
            xFill = getattr(fill, '_xFill', None)
            if xFill is None or not hasattr(xFill, 'blipFill'):
                return None
            blipFill = xFill.blipFill
            if not hasattr(blipFill, 'blip') or not blipFill.blip:
                return None
            rEmbed = blipFill.blip.embed
            if not rEmbed:
                return None
            
            # Get the image part from slide
            slide_part = slide._part
            if slide_part:
                image_part = slide_part.rels[rEmbed].target_part
                return image_part.blob
        except Exception as e:
            logger.debug(f"Failed to extract picture from background fill: {e}")
            return None

    def _try_extract_picture_fill_blob_for_part(self, owner_part, fill) -> Optional[bytes]:
        """Extract picture data from a background fill using a given owning part (layout/master)."""
        try:
            xFill = getattr(fill, '_xFill', None)
            if xFill is None or not hasattr(xFill, 'blipFill'):
                return None
            blipFill = xFill.blipFill
            if not hasattr(blipFill, 'blip') or not blipFill.blip:
                return None
            rEmbed = blipFill.blip.embed
            if not rEmbed or not hasattr(owner_part, 'rels'):
                return None
            rel = owner_part.rels.get(rEmbed)
            if not rel:
                return None
            image_part = rel.target_part
            if hasattr(image_part, 'blob'):
                return image_part.blob
        except Exception as e:
            logger.debug(f"Failed to extract picture from owner part background: {e}")
            return None
    
    def _try_extract_picture_fill_blob(self, shape) -> Optional[bytes]:
        """Attempt to extract embedded image bytes from a shape's picture fill (blipFill)."""
        try:
            if not hasattr(shape, 'fill') or shape.fill is None:
                return None
            xFill = getattr(shape.fill, '_xFill', None)
            if xFill is None or not hasattr(xFill, 'blipFill'):
                # Fallback to XML-based extraction when fill API doesn't expose blipFill
                try:
                    el = getattr(shape, '_element', None)
                    if el is None:
                        return None
                    from lxml import etree
                    ns = {
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                    }
                    # Look for any blip under a blipFill in this shape
                    blip = el.find('.//a:blipFill/a:blip', namespaces=ns)
                    if blip is None:
                        # Some picture shapes use p:pic/a:blipFill/a:blip
                        blip = el.find('.//p:pic//a:blip', namespaces=ns)
                    if blip is None:
                        return None
                    rid = blip.get('{%s}embed' % ns['r']) or blip.get('r:embed') or blip.get('{%s}link' % ns['r'])
                    if not rid:
                        return None
                    # Try both related_parts and rels maps
                    part = getattr(shape, 'part', None)
                    if part is not None:
                        image_part = None
                        try:
                            image_part = part.related_parts.get(rid)
                        except Exception:
                            image_part = None
                        if image_part is None:
                            try:
                                rel = getattr(part, 'rels', {})
                                if rid in rel:
                                    image_part = rel[rid].target_part
                            except Exception:
                                image_part = None
                        if image_part is not None and hasattr(image_part, 'blob'):
                            return image_part.blob
                except Exception:
                    return None
                return None
            blipFill = xFill.blipFill
            if blipFill is None or not hasattr(blipFill, 'blip'):
                return None
            blip = blipFill.blip
            rId = getattr(blip, 'rEmbed', None)
            if not rId:
                # Try XML fallback (embed/link attr)
                try:
                    el = getattr(shape, '_element', None)
                    if el is not None:
                        from lxml import etree
                        ns = {
                            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                        }
                        xml_blip = el.find('.//a:blipFill/a:blip', namespaces=ns)
                        if xml_blip is not None:
                            rId = xml_blip.get('{%s}embed' % ns['r']) or xml_blip.get('r:embed') or xml_blip.get('{%s}link' % ns['r'])
                except Exception:
                    rId = None
                if not rId:
                    return None
            image_part = None
            try:
                image_part = shape.part.related_parts.get(rId)
            except Exception:
                image_part = None
            if image_part is None:
                try:
                    rels = getattr(shape.part, 'rels', {})
                    if rId in rels:
                        image_part = rels[rId].target_part
                except Exception:
                    image_part = None
            if image_part is None or not hasattr(image_part, 'blob'):
                return None
            return image_part.blob
        except Exception:
            return None

    def _extract_blipfill_crop_rect(self, shape) -> Optional[Dict[str, float]]:
        """Extract cropRect fractions from a shape's blipFill (a:srcRect l/t/r/b in 1/100000 units)."""
        try:
            el = getattr(shape, '_element', None)
            if el is None:
                return None
            # Namespaces
            ns = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            }
            src_rect = el.find('.//a:blipFill/a:srcRect', namespaces=ns)
            if src_rect is None:
                return None
            def _frac(attr: str) -> float:
                val = src_rect.get(attr)
                if val is None:
                    return 0.0
                try:
                    return max(0.0, min(1.0, float(int(val)) / 100000.0))
                except Exception:
                    return 0.0
            left = _frac('l')
            right = _frac('r')
            top = _frac('t')
            bottom = _frac('b')
            if any(v > 0 for v in (left, right, top, bottom)):
                return {
                    'left': left,
                    'right': right,
                    'top': top,
                    'bottom': bottom,
                }
            return None
        except Exception:
            return None

    def _detect_image_mime_from_bytes(self, blob: bytes) -> str:
        """Best-effort detection of common image mime types from raw bytes."""
        try:
            if blob.startswith(b"\x89PNG\r\n\x1a\n"):
                return 'image/png'
            if blob.startswith(b"\xff\xd8"):
                return 'image/jpeg'
            if blob.startswith(b"GIF8"):
                return 'image/gif'
        except Exception:
            pass
        return 'image/png'

    def _create_image_component_from_blob(
        self,
        image_bytes: bytes,
        shape,
        bounds: Dict[str, float],
        z_index: int,
        crop_rect: Optional[Dict[str, float]] = None,
        flip_x: bool = False,
        flip_y: bool = False,
    ) -> Dict[str, Any]:
        """Create an Image component from raw image bytes extracted via blipFill."""
        try:
            mime_type = self._detect_image_mime_from_bytes(image_bytes)
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            data_url = f"data:{mime_type};base64,{image_base64}"
            props: Dict[str, Any] = {
                **bounds,
                'src': data_url,
                'alt': getattr(shape, 'name', None) or 'Image',
                'objectFit': 'cover' if crop_rect else 'contain',
                'borderRadius': 0,
                'borderWidth': 0,
                'borderColor': '#000000ff',
                'shadow': False,
                'shadowBlur': 10,
                'shadowColor': '#0000004D',
                'shadowOffsetX': 0,
                'shadowOffsetY': 4,
                'shadowSpread': 0,
                'opacity': 1,
                'zIndex': z_index,
                'filterPreset': 'none',
                'brightness': 100,
                'contrast': 100,
                'saturation': 100,
                'grayscale': 0,
                'sepia': 0,
                'hueRotate': 0,
                'blur': 0,
            }
            if crop_rect:
                props['cropRect'] = crop_rect
            if flip_x:
                props['flipX'] = True
            if flip_y:
                props['flipY'] = True
            return {
                'id': str(uuid.uuid4()),
                'type': 'Image',
                'props': props,
            }
        except Exception as e:
            logger.warning(f"Failed to build image component from blob: {e}")
            raise

    def _create_shape_with_picture_image(
        self,
        image_bytes: bytes,
        shape,
        bounds: Dict[str, float],
        z_index: int,
        crop_rect: Optional[Dict[str, float]] = None,
        flip_x: bool = False,
        flip_y: bool = False,
    ) -> Dict[str, Any]:
        """Create a Shape component that preserves the original shape mask and applies the picture as fill via props.imageFill.

        Frontend can render this by clipping to the shape and drawing the image with objectFit and cropRect.
        Falls back to plain Image if frontend lacks masked image support.
        """
        try:
            # Determine original mapped shape type to hint frontend masking
            shape_type = 'rectangle'
            try:
                if getattr(shape, 'auto_shape_type', None) is not None:
                    shape_type = str(getattr(shape, 'auto_shape_type')).lower()
            except Exception:
                pass

            # Build data URL
            mime_type = self._detect_image_mime_from_bytes(image_bytes)
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            data_url = f"data:{mime_type};base64,{image_base64}"

            # If the shape is circular, compute border radius hint
            border_radius = 0
            try:
                if hasattr(shape, 'auto_shape_type') and str(shape.auto_shape_type).endswith('OVAL'):
                    border_radius = int(round(min(bounds["width"], bounds["height"]) / 2))
            except Exception:
                pass

            return {
                'id': str(uuid.uuid4()),
                'type': 'Shape',
                'props': {
                    **bounds,
                    'shape': 'rectangle',  # Frontend will clip via imageFill.mask if provided
                    **({'borderRadius': border_radius} if border_radius else {}),
                    'fill': '#00000000',
                    'stroke': '#00000000',
                    'strokeWidth': 0,
                    'opacity': 1,
                    'zIndex': z_index,
                    # Extended prop to preserve original picture fill semantics
                    'imageFill': {
                        'src': data_url,
                        'objectFit': 'cover' if crop_rect else 'contain',
                        **({'cropRect': crop_rect} if crop_rect else {}),
                        **({'flipX': True} if flip_x else {}),
                        **({'flipY': True} if flip_y else {}),
                        # Optional: mask type for future advanced shapes
                        'mask': 'auto',
                    },
                },
            }
        except Exception:
            # Fallback to Image component if anything fails
            return self._create_image_component_from_blob(image_bytes, shape, bounds, z_index, crop_rect, flip_x, flip_y)
        
    def _extract_background(self, slide, scale: float) -> Optional[Dict[str, Any]]:
        """Extract slide background"""
        try:
            # Helper: direct XML lookup for slide bg blip
            def _extract_slide_bg_blip(sl):
                try:
                    from lxml import etree
                    ns = {
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                    }
                    el = getattr(sl, '_element', None)
                    if el is None:
                        return None
                    blip = el.find('.//p:bg/p:bgPr/a:blipFill/a:blip', namespaces=ns)
                    if blip is not None:
                        rid = blip.get('{%s}embed' % ns['r']) or blip.get('r:embed')
                        if rid:
                            part = getattr(sl, '_part', None)
                            if part is not None and hasattr(part, 'rels') and rid in part.rels:
                                image_part = part.rels[rid].target_part
                                if hasattr(image_part, 'blob'):
                                    return image_part.blob
                except Exception:
                    return None
                return None
            # First check for shape-based backgrounds (common in Google Slides exports)
            # Look for a shape that covers the entire slide
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'width') and hasattr(shape, 'height'):
                        # Do NOT convert actual picture shapes to background; keep them as Image components
                        if getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.PICTURE:
                            continue
                        # Compute scaled shape bounds (to compare with normalized 1920x1080)
                        shape_width = self._emu_to_px(shape.width) * scale
                        shape_height = self._emu_to_px(shape.height) * scale
                        shape_left = (self._emu_to_px(shape.left) * scale) if shape.left else 0
                        shape_top = (self._emu_to_px(shape.top) * scale) if shape.top else 0
                        
                        # If shape is at/near origin and covers most of slide, it's likely a background
                        slide_width_px = 1920
                        slide_height_px = 1080
                        if (shape_left <= 50 and shape_top <= 50 and 
                            shape_width >= slide_width_px * 0.9 and 
                            shape_height >= slide_height_px * 0.9):
                            
                            # Extract fill if available
                            if hasattr(shape, 'fill') and shape.fill:
                                fill = shape.fill
                                # Solid color
                                if fill.type == 1:
                                    color = self._get_color_hex(fill.fore_color, is_background=True, theme_colors=self.theme_colors) if fill.fore_color else "#FFFFFFff"
                                    return ({
                                        "id": str(uuid.uuid4()),
                                        "type": "Background",
                                        "props": {
                                            "position": {"x": 0, "y": 0},
                                            "width": 1920,
                                            "height": 1080,
                                            "opacity": 1,
                                            "rotation": 0,
                                            "zIndex": 0,
                                            "backgroundType": "color",
                                            "backgroundColor": color,
                                            "gradient": None,
                                            "isAnimated": False,
                                            "animationSpeed": 1,
                                            "kind": "background"
                                        }
                                    }, shape._element)
                                else:
                                    # Picture fill → convert to Background image
                                    picture_blob = self._try_extract_picture_fill_blob(shape)
                                    if picture_blob is not None:
                                        try:
                                            image_base64 = base64.b64encode(picture_blob).decode('utf-8')
                                            data_url = f"data:image/png;base64,{image_base64}"
                                            return ({
                                                "id": str(uuid.uuid4()),
                                                "type": "Background",
                                                "props": {
                                                    "position": {"x": 0, "y": 0},
                                                    "width": 1920,
                                                    "height": 1080,
                                                    "opacity": 1,
                                                    "rotation": 0,
                                                    "zIndex": 0,
                                                    "backgroundType": "image",
                                                    "backgroundImageUrl": data_url,
                                                    "backgroundColor": "#FFFFFFff",
                                                    "gradient": None,
                                                    "isAnimated": False,
                                                    "animationSpeed": 1,
                                                    "kind": "background"
                                                }
                                            }, shape._element)
                                        except Exception:
                                            pass
                                    # Gradient fill on shape treated as gradient background
                                    try:
                                        if hasattr(fill, 'gradient_stops') and fill.gradient_stops:
                                            stops_list = list(fill.gradient_stops)
                                            if len(stops_list) >= 2:
                                                gradient_stops = []
                                                total = max(1, len(stops_list) - 1)
                                                for i, stop in enumerate(stops_list):
                                                    c = self._get_color_hex(stop.color, is_background=True, theme_colors=self.theme_colors)
                                                    pos = 0 if total == 0 else int(round((i / total) * 100))
                                                    gradient_stops.append({"color": c, "position": pos})
                                                return ({
                                                    "id": str(uuid.uuid4()),
                                                    "type": "Background",
                                                    "props": {
                                                        "position": {"x": 0, "y": 0},
                                                        "width": 1920,
                                                        "height": 1080,
                                                        "opacity": 1,
                                                        "rotation": 0,
                                                        "zIndex": 0,
                                                        "backgroundType": "gradient",
                                                        "backgroundColor": gradient_stops[0]["color"],
                                                        "gradient": {"type": "linear", "angle": 0, "stops": gradient_stops},
                                                        "isAnimated": False,
                                                        "animationSpeed": 1,
                                                        "kind": "background"
                                                    }
                                                }, shape._element)
                                    except Exception:
                                        pass
                except:
                    continue
            
            # Then check slide background property
            background = slide.background
            
            # If slide doesn't have a background, check layout and master
            if not background or not hasattr(background, 'fill'):
                # Try layout background
                if hasattr(slide, 'slide_layout') and slide.slide_layout:
                    layout_bg = slide.slide_layout.background
                    if layout_bg and hasattr(layout_bg, 'fill'):
                        background = layout_bg
                
                # If still no background, try master
                if (not background or not hasattr(background, 'fill')) and hasattr(slide, 'slide_layout'):
                    master = slide.slide_layout.slide_master
                    if master and hasattr(master, 'background'):
                        master_bg = master.background
                        if master_bg and hasattr(master_bg, 'fill'):
                            background = master_bg
                
                # If still no background, return a default based on theme
                if not background or not hasattr(background, 'fill'):
                    # Use white as default background (most common)
                    # Don't use dk1 as it might be black in some themes
                    default_color = '#FFFFFFff'
                    return {
                        "id": str(uuid.uuid4()),
                        "type": "Background",
                        "props": {
                            "position": {"x": 0, "y": 0},
                            "width": 1920,
                            "height": 1080,
                            "opacity": 1,
                            "rotation": 0,
                            "zIndex": 0,
                            "backgroundType": "color",
                            "backgroundColor": default_color,
                            "gradient": None,
                            "isAnimated": False,
                            "animationSpeed": 1,
                            "kind": "background"
                        }
                    }
                
            fill = background.fill
            
            # Import MSO_FILL_TYPE for proper comparison
            from pptx.enum.dml import MSO_FILL_TYPE
            
            logger.debug(f"Slide background fill type: {fill.type}")
            
            # Slide-level picture background via XML (handles cases python-pptx misses)
            direct_blob = _extract_slide_bg_blip(slide)
            if direct_blob:
                try:
                    image_base64 = base64.b64encode(direct_blob).decode('utf-8')
                    data_url = f"data:image/png;base64,{image_base64}"
                    return {
                        "id": str(uuid.uuid4()),
                        "type": "Background",
                        "props": {
                            "position": {"x": 0, "y": 0},
                            "width": 1920,
                            "height": 1080,
                            "opacity": 1,
                            "rotation": 0,
                            "zIndex": 0,
                            "backgroundType": "image",
                            "backgroundImageUrl": data_url,
                            "backgroundColor": "#FFFFFFff",
                            "gradient": None,
                            "isAnimated": False,
                            "animationSpeed": 1,
                            "kind": "background"
                        }
                    }
                except Exception:
                    pass
            
            # Solid color background
            if fill.type == MSO_FILL_TYPE.SOLID:
                color = self._get_color_hex(fill.fore_color, is_background=True, theme_colors=self.theme_colors) if fill.fore_color else "#FFFFFFff"
                return {
                    "id": str(uuid.uuid4()),
                    "type": "Background",
                    "props": {
                        "position": {"x": 0, "y": 0},
                        "width": 1920,
                        "height": 1080,
                        "opacity": 1,
                        "rotation": 0,
                        "zIndex": 0,
                        "backgroundType": "color",
                        "backgroundColor": color,
                        "gradient": None,
                        "isAnimated": False,
                        "animationSpeed": 1,
                        "kind": "background"
                    }
                }
            elif fill.type == MSO_FILL_TYPE.GRADIENT:
                # Extract gradient stops {color, position}
                gradient_stops: List[Dict[str, Any]] = []
                try:
                    if hasattr(fill, 'gradient_stops'):
                        stops_list = list(fill.gradient_stops)
                        total = max(1, len(stops_list) - 1)
                        for i, stop in enumerate(stops_list):
                            color = self._get_color_hex(stop.color, is_background=True, theme_colors=self.theme_colors)
                            position = 0 if total == 0 else int(round((i / total) * 100))
                            gradient_stops.append({"color": color, "position": position})
                except Exception:
                    gradient_stops = []

                if len(gradient_stops) < 2:
                    gradient_stops = [
                        {"color": "#FFFFFFff", "position": 0},
                        {"color": "#000000ff", "position": 100},
                    ]
                    
                return {
                    "id": str(uuid.uuid4()),
                    "type": "Background",
                    "props": {
                        "position": {"x": 0, "y": 0},
                        "width": 1920,
                        "height": 1080,
                        "opacity": 1,
                        "rotation": 0,
                        "zIndex": 0,
                        "backgroundType": "gradient",
                        "backgroundColor": gradient_stops[0]["color"],
                        "gradient": {
                            "type": "linear",
                            "angle": 0,
                            "stops": gradient_stops
                        },
                        "isAnimated": False,
                        "animationSpeed": 1,
                        "kind": "background"
                    }
                }
            elif fill.type == MSO_FILL_TYPE.PICTURE:
                # Try to extract the image
                # For slide backgrounds, we need to pass the slide object
                picture_blob = self._try_extract_picture_fill_blob_from_background(slide, fill)
                if picture_blob:
                    try:
                        image_base64 = base64.b64encode(picture_blob).decode('utf-8')
                        data_url = f"data:image/png;base64,{image_base64}"
                        
                        return {
                            "id": str(uuid.uuid4()),
                            "type": "Background",
                            "props": {
                                "position": {"x": 0, "y": 0},
                                "width": 1920,
                                "height": 1080,
                                "opacity": 1,
                                "rotation": 0,
                                "zIndex": 0,
                                "backgroundType": "image",
                                "backgroundImageUrl": data_url,
                                "backgroundColor": "#FFFFFFff",
                                "gradient": None,
                                "isAnimated": False,
                                "animationSpeed": 1,
                                "kind": "background"
                            }
                        }
                    except Exception:
                        pass
            elif fill.type == MSO_FILL_TYPE.BACKGROUND:
                # BACKGROUND type means inherit from parent (layout or master)
                bg_color = "#FFFFFFff"  # Default white
                try:
                    if hasattr(slide, 'slide_layout') and slide.slide_layout:
                        layout = slide.slide_layout
                        
                        # First check if layout has a background
                        if hasattr(layout, 'background') and layout.background:
                            layout_bg = layout.background
                            if hasattr(layout_bg, 'fill') and layout_bg.fill:
                                layout_fill = layout_bg.fill
                                if layout_fill.type == MSO_FILL_TYPE.SOLID:
                                    # Use layout's solid fill
                                    bg_color = self._get_color_hex(layout_fill.fore_color, is_text_color=False, is_background=True, theme_colors=self.theme_colors)
                                elif layout_fill.type == MSO_FILL_TYPE.PICTURE:
                                    # Extract picture from layout background
                                    blob = self._try_extract_picture_fill_blob_for_part(getattr(layout, 'part', None), layout_fill)
                                    if blob:
                                        image_base64 = base64.b64encode(blob).decode('utf-8')
                                        data_url = f"data:image/png;base64,{image_base64}"
                                        return {
                                            "id": str(uuid.uuid4()),
                                            "type": "Background",
                                            "props": {
                                                "position": {"x": 0, "y": 0},
                                                "width": 1920,
                                                "height": 1080,
                                                "opacity": 1,
                                                "rotation": 0,
                                                "zIndex": 0,
                                                "backgroundType": "image",
                                                "backgroundImageUrl": data_url,
                                                "backgroundColor": "#FFFFFFff",
                                                "gradient": None,
                                                "isAnimated": False,
                                                "animationSpeed": 1,
                                                "kind": "background"
                                            }
                                        }
                                elif layout_fill.type == MSO_FILL_TYPE.BACKGROUND:
                                    # Layout also inherits, check master
                                    if hasattr(layout, 'slide_master') and layout.slide_master:
                                        master = layout.slide_master
                                        if hasattr(master, 'background') and master.background:
                                            master_bg = master.background
                                            if hasattr(master_bg, 'fill') and master_bg.fill:
                                                master_fill = master_bg.fill
                                                if master_fill.type == MSO_FILL_TYPE.SOLID:
                                                    bg_color = self._get_color_hex(master_fill.fore_color, is_text_color=False, is_background=True, theme_colors=self.theme_colors)
                                                elif master_fill.type == MSO_FILL_TYPE.PICTURE:
                                                    blob = self._try_extract_picture_fill_blob_for_part(getattr(master, 'part', None), master_fill)
                                                    if blob:
                                                        image_base64 = base64.b64encode(blob).decode('utf-8')
                                                        data_url = f"data:image/png;base64,{image_base64}"
                                                        return {
                                                            "id": str(uuid.uuid4()),
                                                            "type": "Background",
                                                            "props": {
                                                                "position": {"x": 0, "y": 0},
                                                                "width": 1920,
                                                                "height": 1080,
                                                                "opacity": 1,
                                                                "rotation": 0,
                                                                "zIndex": 0,
                                                                "backgroundType": "image",
                                                                "backgroundImageUrl": data_url,
                                                                "backgroundColor": "#FFFFFFff",
                                                                "gradient": None,
                                                                "isAnimated": False,
                                                                "animationSpeed": 1,
                                                                "kind": "background"
                                                            }
                                                        }
                                    elif master_fill.type == MSO_FILL_TYPE.GRADIENT:
                                        # Extract full gradient instead of just first color
                                        if hasattr(master_fill, 'gradient_stops'):
                                            stops_list = list(master_fill.gradient_stops)
                                            if len(stops_list) >= 2:
                                                # Return gradient background instead of solid
                                                gradient_stops = []
                                                total = max(1, len(stops_list) - 1)
                                                for i, stop in enumerate(stops_list):
                                                    color = self._get_color_hex(stop.color, is_background=True, theme_colors=self.theme_colors)
                                                    position = 0 if total == 0 else int(round((i / total) * 100))
                                                    gradient_stops.append({"color": color, "position": position})
                                                
                                                return {
                                                    "id": str(uuid.uuid4()),
                                                    "type": "Background",
                                                    "props": {
                                                        "position": {"x": 0, "y": 0},
                                                        "width": 1920,
                                                        "height": 1080,
                                                        "opacity": 1,
                                                        "rotation": 0,
                                                        "zIndex": 0,
                                                        "backgroundType": "gradient",
                                                        "backgroundColor": gradient_stops[0]["color"],
                                                        "gradient": {
                                                            "type": "linear",
                                                            "angle": 0,
                                                            "stops": gradient_stops
                                                        },
                                                        "isAnimated": False,
                                                        "animationSpeed": 1,
                                                        "kind": "background"
                                                    }
                                                }
                                            else:
                                                # Fallback to first color if not enough stops
                                                bg_color = self._get_color_hex(stops_list[0].color, is_background=True, theme_colors=self.theme_colors) if stops_list else bg_color
                except Exception as e:
                    logger.debug(f"Could not extract master background: {e}")
                
                return {
                    "id": str(uuid.uuid4()),
                    "type": "Background",
                    "props": {
                        "position": {"x": 0, "y": 0},
                        "width": 1920,
                        "height": 1080,
                        "opacity": 1,
                        "rotation": 0,
                        "zIndex": 0,
                        "backgroundType": "color",
                        "backgroundColor": bg_color,
                        "gradient": None,
                        "isAnimated": False,
                        "animationSpeed": 1,
                        "kind": "background"
                    }
                }
        except Exception as e:
            logger.warning(f"Failed to extract background: {e}")
            
        return None
        
    def _extract_slide_title(self, slide) -> Optional[str]:
        """Extract slide title from title placeholder"""
        try:
            for shape in slide.shapes:
                if shape.is_placeholder and hasattr(shape, 'text_frame'):
                    # Check if it's a title placeholder
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type in [1, 3]:  # TITLE or CENTER_TITLE
                        if shape.text_frame and shape.text_frame.text:
                            return shape.text_frame.text.strip()
        except:
            pass
        return None
        
    def _get_shape_bounds(self, shape, scale: float) -> Dict[str, float]:
        """Get shape position and size in pixels"""
        # Convert from EMU to pixels and apply scale
        # Use round instead of int for better accuracy
        x = round(self._emu_to_px(shape.left) * scale) if shape.left is not None else 0
        y = round(self._emu_to_px(shape.top) * scale) if shape.top is not None else 0
        width = round(self._emu_to_px(shape.width) * scale) if shape.width is not None else 100
        height = round(self._emu_to_px(shape.height) * scale) if shape.height is not None else 100
        
        # Keep original positions including negative ones
        # This preserves the exact layout from PowerPoint
        
        # Handle rotation if available; fall back to XML a:xfrm@rot (1/60000 deg)
        rotation = 0
        try:
            if hasattr(shape, 'rotation') and shape.rotation is not None:
                rotation = int(shape.rotation)
            if not rotation:
                # Fallback: read from underlying XML transform
                el = getattr(shape, '_element', None)
                xfrm = getattr(el, 'xfrm', None) if el is not None else None
                if xfrm is not None and getattr(xfrm, 'rot', None) is not None:
                    try:
                        xml_rot = float(xfrm.rot)
                        # a:xfrm@rot is in 1/60000 degrees
                        rotation = int(round(xml_rot / 60000.0))
                    except Exception:
                        pass
            rotation = rotation % 360
        except Exception:
            rotation = 0
            
        return {
            "position": {"x": x, "y": y},
            "width": width,
            "height": height,
            "rotation": rotation
        }

    def _normalize_color(self, color_value: Optional[str]) -> str:
        """Convert #RRGGBBAA to rgba(r,g,b,a). Return 'none' for transparent."""
        if not color_value:
            return "none"
        try:
            if color_value.startswith('#'):
                if len(color_value) == 9:  # #RRGGBBAA
                    r = int(color_value[1:3], 16)
                    g = int(color_value[3:5], 16)
                    b = int(color_value[5:7], 16)
                    a = int(color_value[7:9], 16) / 255.0
                    if a <= 0:
                        return "none"
                    return f"rgba({r},{g},{b},{a:.3f})"
                elif len(color_value) == 7:
                    return color_value
        except Exception:
            pass
        return color_value
        
    def _emu_to_px(self, emu: int) -> float:
        """Convert EMU to pixels at 96 DPI"""
        if emu is None:
            return 0
        # 1 inch = 914400 EMU, 1 inch = 96 pixels at 96 DPI
        return emu * 96 / 914400
        
    def _extract_theme_colors(self, presentation):
        """Extract color scheme from theme"""
        try:
            # First try to use stored file path if available
            if hasattr(self, '_file_path') and self._file_path:
                file_path = self._file_path
                
                # Open as zip file and extract theme directly
                import zipfile
                import xml.etree.ElementTree as ET
                
                with zipfile.ZipFile(file_path, 'r') as pptx_zip:
                    # Find theme files
                    theme_files = [f for f in pptx_zip.namelist() if 'theme/theme' in f and f.endswith('.xml')]
                    
                    if theme_files:
                        # Read first theme
                        with pptx_zip.open(theme_files[0]) as theme_xml:
                            tree = ET.parse(theme_xml)
                            root = tree.getroot()
                            
                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            
                            clr_scheme = root.find('.//a:clrScheme', ns)
                            if clr_scheme is not None:
                                colors = {}
                                
                                for child in clr_scheme:
                                    tag = child.tag.split('}')[-1]
                                    srgb_clr = child.find('.//a:srgbClr', ns)
                                    sys_clr = child.find('.//a:sysClr', ns)
                                    
                                    if srgb_clr is not None:
                                        color_val = srgb_clr.get('val')
                                        colors[tag] = f"#{color_val}ff"
                                    elif sys_clr is not None:
                                        last_clr = sys_clr.get('lastClr')
                                        if last_clr:
                                            colors[tag] = f"#{last_clr}ff"
                                
                                return colors
        except Exception as e:
            logger.debug(f"Error extracting theme colors via zipfile: {e}")
            
        try:
            # Try to access theme through the package
            if hasattr(presentation, 'package'):
                package = presentation.package
                
                # Look for theme parts
                for rel in package.relationships:
                    if 'theme' in rel.reltype:
                        theme_part = rel.target_part
                        if hasattr(theme_part, 'blob'):
                            # Parse theme XML
                            import xml.etree.ElementTree as ET
                            theme_xml = theme_part.blob
                            root = ET.fromstring(theme_xml)
                            
                            # Define namespace
                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            
                            # Find color scheme
                            clr_scheme = root.find('.//a:clrScheme', ns)
                            if clr_scheme is not None:
                                colors = {}
                                
                                # Map theme color names to their values
                                for child in clr_scheme:
                                    tag = child.tag.split('}')[-1]
                                    
                                    # Get color value
                                    srgb_clr = child.find('.//a:srgbClr', ns)
                                    sys_clr = child.find('.//a:sysClr', ns)
                                    
                                    if srgb_clr is not None:
                                        color_val = srgb_clr.get('val')
                                        colors[tag] = f"#{color_val}ff"
                                    elif sys_clr is not None:
                                        # System colors - use lastClr attribute
                                        last_clr = sys_clr.get('lastClr')
                                        if last_clr:
                                            colors[tag] = f"#{last_clr}ff"
                                
                                return colors
            
            # Alternative: Try through slide master
            if hasattr(presentation, 'slide_masters') and len(presentation.slide_masters) > 0:
                master = presentation.slide_masters[0]
                
                # Access through master part's package
                if hasattr(master, 'part') and hasattr(master.part, 'package'):
                    package = master.part.package
                    
                    # Look for theme in package parts
                    for partname, part in package.iter_parts():
                        if 'theme' in partname:
                            if hasattr(part, 'blob'):
                                # Parse theme XML
                                import xml.etree.ElementTree as ET
                                theme_xml = part.blob
                                root = ET.fromstring(theme_xml)
                                
                                # Define namespace
                                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                                
                                # Find color scheme
                                clr_scheme = root.find('.//a:clrScheme', ns)
                                if clr_scheme is not None:
                                    colors = {}
                                    
                                    # Map theme color names to their values
                                    for child in clr_scheme:
                                        tag = child.tag.split('}')[-1]
                                        
                                        # Get color value
                                        srgb_clr = child.find('.//a:srgbClr', ns)
                                        sys_clr = child.find('.//a:sysClr', ns)
                                        
                                        if srgb_clr is not None:
                                            color_val = srgb_clr.get('val')
                                            colors[tag] = f"#{color_val}ff"
                                        elif sys_clr is not None:
                                            # System colors - use lastClr attribute
                                            last_clr = sys_clr.get('lastClr')
                                            if last_clr:
                                                colors[tag] = f"#{last_clr}ff"
                                    
                                    return colors
        except Exception as e:
            logger.debug(f"Error extracting theme colors: {e}")
        
        # Return conservative default theme mapping (Office-like)
        return {
            'dk1': '#000000ff',
            'lt1': '#FFFFFFFF',
            'dk2': '#1F497Dff',
            'lt2': '#EEECECFF',
            'accent1': '#4472C4FF',
            'accent2': '#ED7D31FF',
            'accent3': '#A5A5A5FF',
            'accent4': '#FFC000FF',
            'accent5': '#5B9BD5FF',
            'accent6': '#70AD47FF',
            'hlink': '#0563C1FF',
            'folHlink': '#954F72FF'
        }
    
    def _get_layout_placeholder_style(self, shape):
        """Get style information from layout placeholder"""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                
                # Get the layout
                slide = shape.part.slide
                if hasattr(slide, 'slide_layout') and slide.slide_layout:
                    layout = slide.slide_layout
                    
                    # Try to find matching placeholder in layout
                    for layout_placeholder in layout.placeholders:
                        if (hasattr(layout_placeholder, 'placeholder_format') and 
                            layout_placeholder.placeholder_format.idx == ph_idx):
                            
                            # Extract font from layout placeholder
                            if hasattr(layout_placeholder, 'text_frame') and layout_placeholder.text_frame:
                                tf = layout_placeholder.text_frame
                                if hasattr(tf, '_element'):
                                    # Check for list style fonts
                                    return self._extract_layout_fonts(tf._element, ph_type)
        except:
            pass
        
        return None
    
    def _get_layout_placeholder_color(self, shape):
        """Get text color from layout placeholder"""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                
                # Get the layout
                slide = shape.part.slide
                if hasattr(slide, 'slide_layout') and slide.slide_layout:
                    layout = slide.slide_layout
                    
                    # Try to find matching placeholder in layout
                    for layout_placeholder in layout.placeholders:
                        if (hasattr(layout_placeholder, 'placeholder_format') and 
                            layout_placeholder.placeholder_format.idx == ph_idx):
                            
                            # Extract color from layout placeholder
                            if hasattr(layout_placeholder, 'text_frame') and layout_placeholder.text_frame:
                                tf = layout_placeholder.text_frame
                                if hasattr(tf, '_element'):
                                    # Extract color from list style
                                    return self._extract_layout_color(tf._element, ph_type)
        except:
            pass
        
        return None
    
    def _extract_layout_fonts(self, txBody_element, ph_type):
        """Extract fonts from layout text body element"""
        try:
            # Use XML find method to access lstStyle
            lstStyle = txBody_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
            if lstStyle is not None:
                # Check level 1 properties
                lvl1pPr = lstStyle.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr')
                if lvl1pPr is not None:
                    defRPr = lvl1pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
                    if defRPr is not None:
                        latin = defRPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                        if latin is not None:
                            typeface = latin.get('typeface')
                            if typeface:
                                # Map known layout fonts
                                if typeface == 'DM Sans Light':
                                    return 'DM Sans'
                                elif typeface and not typeface.startswith('+'):
                                    return typeface
        except:
            pass
        
        return None
    
    def _extract_layout_color(self, txBody_element, ph_type):
        """Extract text color from layout text body element"""
        try:
            # Use XML find method to access lstStyle
            lstStyle = txBody_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
            logger.debug(f"_extract_layout_color: lstStyle found: {lstStyle is not None}")
            if lstStyle is not None:
                # Check level 1 properties - use direct child search
                lvl1pPr = lstStyle.find('./{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr')
                logger.debug(f"_extract_layout_color: lvl1pPr found: {lvl1pPr is not None}")
                if lvl1pPr is not None:
                    defRPr = lvl1pPr.find('./{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
                    logger.debug(f"_extract_layout_color: defRPr found: {defRPr is not None}")
                    if defRPr is not None:
                        # Check for solid fill - use direct child search, not descendant
                        solidFill = defRPr.find('./{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                        logger.debug(f"_extract_layout_color: solidFill found: {solidFill is not None}")
                        if solidFill is not None:
                            # Check for scheme color
                            schemeClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
                            if schemeClr is not None:
                                val = schemeClr.get('val')
                                if val:
                                    # Map scheme color name to theme color
                                    return val  # Return the scheme color name (e.g., 'dk1')
                            
                            # Check for RGB color
                            srgbClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                            if srgbClr is not None:
                                val = srgbClr.get('val')
                                if val:
                                    return f"#{val}ff"  # Return RGB color
        except:
            pass
        
        return None
    
    def _get_layout_placeholder_size(self, shape):
        """Get font size from layout placeholder"""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                logger.debug(f"Looking for layout size for placeholder idx={ph_idx}, type={ph_type}")
                
                # Get the layout
                slide = shape.part.slide
                if hasattr(slide, 'slide_layout') and slide.slide_layout:
                    layout = slide.slide_layout
                    
                    # Try to find matching placeholder in layout
                    for layout_placeholder in layout.placeholders:
                        if (hasattr(layout_placeholder, 'placeholder_format') and 
                            layout_placeholder.placeholder_format.idx == ph_idx):
                            logger.debug(f"Found matching layout placeholder!")
                            
                            # Extract font size from layout placeholder
                            if hasattr(layout_placeholder, 'text_frame') and layout_placeholder.text_frame:
                                tf = layout_placeholder.text_frame
                                if hasattr(tf, '_element'):
                                    size = self._extract_layout_font_size(tf._element, ph_type)
                                    logger.debug(f"Extracted layout size: {size}")
                                    return size
                    logger.debug(f"No matching layout placeholder found for idx={ph_idx}")
        except Exception as e:
            logger.debug(f"Error in _get_layout_placeholder_size: {e}")
        
        return None
    
    def _extract_layout_font_size(self, txBody_element, ph_type):
        """Extract font size from layout text body element"""
        try:
            logger.debug(f"Extracting font size from layout element")
            # Use XML find method to access lstStyle
            lstStyle = txBody_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
            if lstStyle is not None:
                logger.debug(f"Found lstStyle")
                # Check level 1 properties
                lvl1pPr = lstStyle.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr')
                if lvl1pPr is not None:
                    logger.debug(f"Found lvl1pPr")
                    defRPr = lvl1pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
                    if defRPr is not None:
                        logger.debug(f"Found defRPr")
                        sz = defRPr.get('sz')
                        if sz:
                            # sz is in hundreds of a point (e.g. 4200 = 42pt)
                            pt_size = int(sz) / 100
                            # Convert to pixels (1pt = 4/3 px at 96 DPI)
                            px_size = int(pt_size * 4/3)
                            logger.debug(f"Found size: {pt_size}pt = {px_size}px")
                            return px_size
                        else:
                            logger.debug(f"No sz attribute in defRPr")
            else:
                logger.debug(f"No lstStyle found in layout element")
        except Exception as e:
            logger.debug(f"Error in _extract_layout_font_size: {e}")
        
        return None
    
    def _extract_default_font(self, text_frame, paragraph):
        """Extract default font from text frame or paragraph properties"""
        try:
            # Try paragraph default run properties
            if hasattr(paragraph, '_element') and paragraph._element is not None:
                pPr = paragraph._element.pPr
                if pPr is not None and hasattr(pPr, 'defRPr'):
                    defRPr = pPr.defRPr
                    if defRPr is not None and hasattr(defRPr, 'latin'):
                        latin = defRPr.latin
                        if latin is not None and hasattr(latin, 'typeface'):
                            return latin.typeface
            
            # Try text frame list style
            if hasattr(text_frame, '_element') and text_frame._element is not None:
                txBody = text_frame._element
                if hasattr(txBody, 'lstStyle') and txBody.lstStyle is not None:
                    lstStyle = txBody.lstStyle
                    # Check level 1 properties (most common)
                    if hasattr(lstStyle, 'lvl1pPr') and lstStyle.lvl1pPr is not None:
                        lvl1 = lstStyle.lvl1pPr
                        if hasattr(lvl1, 'defRPr') and lvl1.defRPr is not None:
                            defRPr = lvl1.defRPr
                            if hasattr(defRPr, 'latin') and defRPr.latin is not None:
                                latin = defRPr.latin
                                if hasattr(latin, 'typeface'):
                                    return latin.typeface
            
            # For embedded fonts from Google Slides, we may need to check the shape
            # In this case, we'll use heuristics based on placeholder type
            shape = getattr(text_frame, '_parent', None)
            if shape and hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    # Title placeholders often use display fonts like Newsreader
                    if ph_type in [1, 3]:  # TITLE or CENTER_TITLE
                        return "Newsreader"  # Known from the embedded fonts
                    else:
                        return "DM Sans"  # Body text font from the layout
        except:
            pass
        
        return None
        
    def _pt_to_px(self, pt: Optional[Length]) -> int:
        """Convert points/centipoints to pixels"""
        if pt is None:
            return 16  # Default font size
        
        # Handle different Length types from python-pptx
        try:
            if hasattr(pt, 'pt'):
                # This is a Length object with points
                points = pt.pt
            elif hasattr(pt, 'cm'):
                # Convert from cm to points (1 cm = 28.35 pt)
                points = pt.cm * 28.35
            elif hasattr(pt, 'inches'):
                # Convert from inches to points (1 inch = 72 pt)
                points = pt.inches * 72
            else:
                # Try direct conversion
                points = float(pt)
        except:
            return 16  # Default on any error
            
        # 1 pt = 1/72 inch, 1 inch = 96 px at 96 DPI
        # For line widths, use exact conversion without minimum
        return int(round(points * 96 / 72))
        
    def _pt_to_pt(self, pt: Optional[Length]) -> float:
        """Convert Length to points (matching Google Slides behavior)"""
        if pt is None:
            return 12.0  # Default font size in points
        
        # Handle different Length types from python-pptx
        try:
            if hasattr(pt, 'pt'):
                # This is a Length object with points
                return pt.pt
            elif hasattr(pt, 'cm'):
                # Convert from cm to points (1 cm = 28.35 pt)
                return pt.cm * 28.35
            elif hasattr(pt, 'inches'):
                # Convert from inches to points (1 inch = 72 pt)
                return pt.inches * 72
            else:
                # Try direct conversion
                return float(pt)
        except:
            return 12.0  # Default on any error
        
    def _get_color_hex(self, color, is_text_color: bool = False, is_background: bool = False, theme_colors: dict = None) -> str:
        """Convert a color object to hex string with alpha
        
        Args:
            color: The color object to convert
            is_text_color: Whether this color is being used for text
            is_background: Whether this color is being used for a slide background
            theme_colors: Optional dict of theme colors to use instead of hardcoded ones
        """
        if color is None:
            # For text, default to black; for other elements, transparent
            return "#000000ff" if is_text_color else "#00000000"
            
        try:
            # Check color type
            color_type = getattr(color, 'type', None)
            
            if hasattr(color, 'rgb') and color.rgb:
                # RGBColor object
                rgb = color.rgb
                if isinstance(rgb, tuple) and len(rgb) >= 3:
                    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}ff"
                else:
                    # Try to handle as an integer
                    try:
                        hex_str = f"{rgb:06x}"
                        return f"#{hex_str}ff"
                    except:
                        pass
            elif color_type == 2 or hasattr(color, 'theme_color'):  # Theme/Scheme color
                # python-pptx theme color constants
                from pptx.enum.dml import MSO_THEME_COLOR_INDEX
                
                # Use provided theme colors or defaults
                if not theme_colors:
                    theme_colors = {
                    MSO_THEME_COLOR_INDEX.DARK_1: "#F9F9F9ff",      # Dark 1 (light background from theme)
                    MSO_THEME_COLOR_INDEX.LIGHT_1: "#183C40ff",     # Light 1 (dark text from theme)
                    MSO_THEME_COLOR_INDEX.DARK_2: "#20494Dff",      # Dark 2
                    MSO_THEME_COLOR_INDEX.LIGHT_2: "#799DA7ff",     # Light 2
                    MSO_THEME_COLOR_INDEX.BACKGROUND_1: "#F9F9F9ff",  # Background 1 (uses dk1)
                    MSO_THEME_COLOR_INDEX.TEXT_1: "#183C40ff",      # Text 1 (uses lt1)
                    MSO_THEME_COLOR_INDEX.BACKGROUND_2: "#20494Dff",  # Background 2 (uses dk2)
                    MSO_THEME_COLOR_INDEX.TEXT_2: "#799DA7ff",      # Text 2 (uses lt2)
                    MSO_THEME_COLOR_INDEX.ACCENT_1: "#B0D5CDff",    # Accent 1
                    MSO_THEME_COLOR_INDEX.ACCENT_2: "#A6DAEAff",    # Accent 2
                    MSO_THEME_COLOR_INDEX.ACCENT_3: "#B78D43ff",    # Accent 3
                    MSO_THEME_COLOR_INDEX.ACCENT_4: "#9AA864ff",    # Accent 4
                    MSO_THEME_COLOR_INDEX.ACCENT_5: "#426345ff",    # Accent 5
                    MSO_THEME_COLOR_INDEX.ACCENT_6: "#DCDCDCff",    # Accent 6
                    # Additional theme colors
                    MSO_THEME_COLOR_INDEX.HYPERLINK: "#B0D5CDff",   # Uses accent1
                    MSO_THEME_COLOR_INDEX.FOLLOWED_HYPERLINK: "#0097A7ff",
                }
                
                theme_color_idx = getattr(color, 'theme_color', None)
                if theme_color_idx is not None:
                    # For text elements, always use the actual theme color (never transparent)
                    # For non-text, non-background elements, TEXT_1/DARK_1 might indicate transparency
                    if not is_text_color and not is_background and theme_color_idx in [MSO_THEME_COLOR_INDEX.TEXT_1, 
                                                                  getattr(MSO_THEME_COLOR_INDEX, 'DARK_1', None)]:
                        logger.debug(f"Converting theme color {theme_color_idx} to transparent for non-text element")
                        return "#00000000"  # Return transparent for non-text uses of TEXT_1/DARK_1
                    
                    # Map MSO theme color indices to theme color names
                    idx_to_name = {
                        MSO_THEME_COLOR_INDEX.DARK_1: 'dk1',
                        MSO_THEME_COLOR_INDEX.LIGHT_1: 'lt1',
                        MSO_THEME_COLOR_INDEX.DARK_2: 'dk2',
                        MSO_THEME_COLOR_INDEX.LIGHT_2: 'lt2',
                        MSO_THEME_COLOR_INDEX.ACCENT_1: 'accent1',
                        MSO_THEME_COLOR_INDEX.ACCENT_2: 'accent2',
                        MSO_THEME_COLOR_INDEX.ACCENT_3: 'accent3',
                        MSO_THEME_COLOR_INDEX.ACCENT_4: 'accent4',
                        MSO_THEME_COLOR_INDEX.ACCENT_5: 'accent5',
                        MSO_THEME_COLOR_INDEX.ACCENT_6: 'accent6',
                        MSO_THEME_COLOR_INDEX.HYPERLINK: 'hlink',
                        MSO_THEME_COLOR_INDEX.FOLLOWED_HYPERLINK: 'folHlink',
                        MSO_THEME_COLOR_INDEX.BACKGROUND_1: 'dk1',  # Maps to dk1
                        MSO_THEME_COLOR_INDEX.TEXT_1: 'lt1',        # Maps to lt1
                        MSO_THEME_COLOR_INDEX.BACKGROUND_2: 'dk2',  # Maps to dk2
                        MSO_THEME_COLOR_INDEX.TEXT_2: 'lt2',        # Maps to lt2
                    }
                    
                    color_name = idx_to_name.get(theme_color_idx)
                    if color_name and color_name in theme_colors:
                        return theme_colors[color_name]
                    
                    # If using old hardcoded mapping, use it
                    return theme_colors.get(theme_color_idx, "#4472C4ff")  # Default to blue accent
                    
            elif color_type == 1:  # RGB color or scheme color with transforms
                # This shouldn't happen if rgb was already checked above
                # Try to extract RGB value again
                if hasattr(color, 'rgb'):
                    try:
                        rgb = color.rgb
                        if hasattr(rgb, '__iter__'):
                            r, g, b = rgb
                        else:
                            hex_str = f"{rgb:06x}"
                            r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
                        # Apply transforms if present (tint/shade/alpha)
                        try:
                            # python-pptx exposes color._xClr for transforms in some cases
                            xClr = getattr(color, '_xClr', None)
                            if xClr is not None:
                                # Check tint ("tint" increases L, shade decreases L)
                                tint = getattr(xClr, 'tint', None)
                                shade = getattr(xClr, 'shade', None)
                                a = getattr(xClr, 'alpha', None)
                                # Apply tint/shade approximately in RGB
                                if tint is not None:
                                    factor = 1 + float(tint)
                                    r = min(255, int(r * factor))
                                    g = min(255, int(g * factor))
                                    b = min(255, int(b * factor))
                                if shade is not None:
                                    factor = 1 - float(shade)
                                    r = max(0, int(r * factor))
                                    g = max(0, int(g * factor))
                                    b = max(0, int(b * factor))
                                alpha_hex = 'ff'
                                if a is not None:
                                    try:
                                        # a is in thousandths of percent sometimes (0..100000)? Fallback to 1 if unknown
                                        aval = float(a)
                                        if aval > 1:
                                            aval = aval / 100000.0
                                        aval = max(0.0, min(1.0, aval))
                                        alpha_hex = f"{int(aval*255):02x}"
                                    except Exception:
                                        alpha_hex = 'ff'
                                return f"#{r:02X}{g:02X}{b:02X}{alpha_hex}"
                        except Exception:
                            pass
                        return f"#{r:02X}{g:02X}{b:02X}ff"
                    except:
                        pass
                # Default based on usage
                return "#000000ff" if is_text_color else "#FFFFFFff"
        except Exception as e:
            logger.warning(f"Failed to extract color: {e}")
            
        # Return appropriate default based on usage
        return "#000000ff" if is_text_color else "#00000000"