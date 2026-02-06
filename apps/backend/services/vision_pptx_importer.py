"""
Vision-based PPTX Importer - Perfect slide recreation using image analysis.

Strategy:
1. Convert PPTX slides to images (LibreOffice → PDF → PNG)
2. Use Gemini Vision to analyze each slide image
3. Generate custom components that visually match the original
4. Extract text separately for editability

This approach guarantees visual fidelity since we're matching pixels, not parsing XML.
"""

import os
import uuid
import logging
import tempfile
import subprocess
import base64
import asyncio
import shutil
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from io import BytesIO

from agents.config import VISION_IMPORT_MODEL

logger = logging.getLogger(__name__)


class VisionPPTXImporter:
    """
    Import PPTX by converting to images and using vision AI to recreate slides.
    """

    def __init__(self):
        self.temp_dir: Optional[str] = None
        self.slide_images: List[str] = []  # Paths to slide images
        self.slide_texts: List[str] = []   # Extracted text per slide
        self.slide_extracted_images: List[List[str]] = []  # Uploaded image URLs per slide

    async def import_file(self, file_path: str) -> Dict[str, Any]:
        """Import PPTX file and return deck data."""
        logger.info(f"[VisionPPTX] Starting import: {file_path}")

        try:
            # Create temp directory for conversion
            self.temp_dir = tempfile.mkdtemp(prefix="pptx_vision_")

            # Step 1: Convert PPTX to images
            self.slide_images = await self._convert_to_images(file_path)
            logger.info(f"[VisionPPTX] Converted to {len(self.slide_images)} slide images")

            if not self.slide_images:
                raise Exception("Failed to convert PPTX to images")

            # Step 2: Extract text from PPTX for editability
            self.slide_texts = self._extract_text_from_pptx(file_path)
            logger.info(f"[VisionPPTX] Extracted text from {len(self.slide_texts)} slides")

            # Step 3: Extract and upload images from PPTX
            self.slide_extracted_images = await self._extract_and_upload_images(file_path)
            logger.info(f"[VisionPPTX] Extracted images from {len(self.slide_extracted_images)} slides")

            # Step 4: Use vision AI to recreate each slide
            slides = await self._recreate_slides_with_vision()
            logger.info(f"[VisionPPTX] Recreated {len(slides)} slides with vision AI")

            # Build result
            deck_name = Path(file_path).stem
            return {
                "uuid": str(uuid.uuid4()),
                "name": deck_name,
                "slides": slides,
                "size": {"width": 1920, "height": 1080},
                "metadata": {
                    "source": "vision_pptx_import",
                    "import_stats": {
                        "slides": len(slides),
                        "method": "vision_ai"
                    }
                }
            }

        finally:
            # Cleanup temp files
            if self.temp_dir and os.path.exists(self.temp_dir):
                try:
                    shutil.rmtree(self.temp_dir)
                except Exception as e:
                    logger.warning(f"[VisionPPTX] Failed to cleanup temp dir: {e}")

    async def import_bytes(self, file_bytes: bytes, filename: str = "presentation.pptx") -> Dict[str, Any]:
        """Import PPTX from bytes."""
        # Save to temp file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        try:
            temp_file.write(file_bytes)
            temp_file.close()
            return await self.import_file(temp_file.name)
        finally:
            try:
                os.unlink(temp_file.name)
            except:
                pass

    async def _convert_to_images(self, pptx_path: str) -> List[str]:
        """Convert PPTX to PNG images using LibreOffice and pdftoppm."""
        images = []

        try:
            # Step 1: Convert PPTX to PDF using LibreOffice
            pdf_path = os.path.join(self.temp_dir, "presentation.pdf")

            # Find LibreOffice
            libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
            if not libreoffice:
                raise Exception("LibreOffice not found")

            # Convert to PDF
            logger.info(f"[VisionPPTX] Converting PPTX to PDF...")
            result = subprocess.run([
                libreoffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", self.temp_dir,
                pptx_path
            ], capture_output=True, text=True, timeout=120)

            if result.returncode != 0:
                logger.error(f"[VisionPPTX] LibreOffice error: {result.stderr}")
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")

            # Find the generated PDF
            pdf_files = list(Path(self.temp_dir).glob("*.pdf"))
            if not pdf_files:
                raise Exception("No PDF file generated")
            pdf_path = str(pdf_files[0])

            # Step 2: Convert PDF pages to PNG images
            logger.info(f"[VisionPPTX] Converting PDF to images...")
            output_prefix = os.path.join(self.temp_dir, "slide")

            # Use pdftoppm for high-quality conversion
            pdftoppm = shutil.which("pdftoppm")
            if pdftoppm:
                result = subprocess.run([
                    pdftoppm,
                    "-png",
                    "-r", "150",  # 150 DPI for good quality without huge files
                    pdf_path,
                    output_prefix
                ], capture_output=True, text=True, timeout=120)

                if result.returncode != 0:
                    logger.warning(f"[VisionPPTX] pdftoppm warning: {result.stderr}")

            # Collect generated images
            image_files = sorted(Path(self.temp_dir).glob("slide-*.png"))
            if not image_files:
                # Try alternate naming pattern
                image_files = sorted(Path(self.temp_dir).glob("slide*.png"))

            for img_path in image_files:
                images.append(str(img_path))

            logger.info(f"[VisionPPTX] Generated {len(images)} slide images")

        except subprocess.TimeoutExpired:
            logger.error("[VisionPPTX] Conversion timed out")
            raise Exception("PPTX conversion timed out")
        except Exception as e:
            logger.error(f"[VisionPPTX] Conversion error: {e}")
            raise

        return images

    def _extract_text_from_pptx(self, pptx_path: str) -> List[str]:
        """Extract text from each slide for editability."""
        texts = []

        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)

            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                texts.append("\n".join(slide_text))

        except Exception as e:
            logger.warning(f"[VisionPPTX] Text extraction error: {e}")
            # Return empty texts for each image
            texts = [""] * len(self.slide_images)

        return texts

    async def _extract_and_upload_images(self, pptx_path: str) -> List[List[str]]:
        """Extract images from each slide and upload to storage. Returns list of image URLs per slide."""
        from services.image_storage_service import ImageStorageService

        slide_images = []

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(pptx_path)
            storage = ImageStorageService()

            for slide_idx, slide in enumerate(prs.slides):
                images_for_slide = []
                img_counter = 0

                for shape in slide.shapes:
                    try:
                        # Check if shape has an image
                        shape_type = getattr(shape, 'shape_type', None)
                        has_picture = shape_type == getattr(MSO_SHAPE_TYPE, 'PICTURE', None)

                        if has_picture and hasattr(shape, 'image'):
                            image = shape.image
                            image_bytes = image.blob

                            # Determine content type
                            content_type = image.content_type or 'image/png'
                            ext = '.png'
                            if 'jpeg' in content_type or 'jpg' in content_type:
                                ext = '.jpg'
                            elif 'gif' in content_type:
                                ext = '.gif'
                            elif 'webp' in content_type:
                                ext = '.webp'

                            # Convert bytes to base64
                            image_b64 = base64.b64encode(image_bytes).decode('utf-8')

                            # Upload using base64 method
                            img_counter += 1
                            filename = f"slide_{slide_idx + 1}_img_{img_counter}{ext}"
                            result = await storage.upload_image_from_base64(
                                image_b64,
                                filename=filename,
                                content_type=content_type,
                                folder="pptx-import"  # Use specific folder for PPTX imports
                            )

                            if result and result.get('url'):
                                # Remove query params from URL for cleaner display
                                url = result['url'].split('?')[0]
                                images_for_slide.append(url)
                                logger.info(f"[VisionPPTX] Uploaded image from slide {slide_idx + 1}: {url[:80]}...")
                    except Exception as e:
                        logger.warning(f"[VisionPPTX] Failed to extract image from shape: {e}")
                        continue

                slide_images.append(images_for_slide)
                if images_for_slide:
                    logger.info(f"[VisionPPTX] Slide {slide_idx + 1}: extracted {len(images_for_slide)} images")

        except Exception as e:
            logger.error(f"[VisionPPTX] Image extraction error: {e}")
            # Return empty lists for each slide
            slide_images = [[] for _ in range(len(self.slide_images))]

        return slide_images

    async def _recreate_slides_with_vision(self) -> List[Dict[str, Any]]:
        """Use vision AI to recreate each slide as custom components - PARALLEL."""

        async def process_slide(idx: int, image_path: str) -> Dict[str, Any]:
            try:
                slide_text = self.slide_texts[idx] if idx < len(self.slide_texts) else ""
                image_urls = self.slide_extracted_images[idx] if idx < len(self.slide_extracted_images) else []
                slide = await self._recreate_single_slide(image_path, idx, slide_text, image_urls)
                logger.info(f"[VisionPPTX] Recreated slide {idx + 1}/{len(self.slide_images)}")
                return slide
            except Exception as e:
                logger.warning(f"[VisionPPTX] Failed to recreate slide {idx + 1}: {e}")
                return self._create_image_fallback_slide(image_path, idx)

        # Process all slides in parallel (batch of 10 at a time to avoid rate limits)
        batch_size = 10
        all_slides = []

        for i in range(0, len(self.slide_images), batch_size):
            batch = self.slide_images[i:i + batch_size]
            tasks = [process_slide(i + j, path) for j, path in enumerate(batch)]
            batch_results = await asyncio.gather(*tasks)
            all_slides.extend(batch_results)
            logger.info(f"[VisionPPTX] Completed batch {i // batch_size + 1}, total: {len(all_slides)}/{len(self.slide_images)}")

        return all_slides

    async def _recreate_single_slide(self, image_path: str, slide_idx: int, slide_text: str, image_urls: List[str] = None) -> Dict[str, Any]:
        """Recreate a single slide using vision AI."""
        image_urls = image_urls or []

        # Read image as base64
        with open(image_path, "rb") as f:
            image_data = f.read()
        image_b64 = base64.b64encode(image_data).decode("utf-8")
        image_data_url = f"data:image/png;base64,{image_b64}"

        # Build image URLs section for prompt
        images_section = ""
        if image_urls:
            images_section = "\nIMAGE URLs (use these exact URLs):\n"
            for i, url in enumerate(image_urls):
                images_section += f"- {url}\n"

        # Simplified prompt for faster processing
        prompt = f"""Recreate this slide as HTML. Match colors, fonts, layout exactly.

TEXT: {slide_text if slide_text else "(extract from image)"}
{images_section}
RULES:
- 1920x1080px, overflow:hidden
- Use Tailwind: <script src="https://cdn.tailwindcss.com"></script>
- Load fonts from Google Fonts
- Use provided image URLs only (no placeholders!)
- Match exact colors (#RRGGBB), spacing, font sizes

Output complete HTML starting with <!DOCTYPE html>. No markdown, just HTML."""

        try:
            # Call Gemini with vision to recreate the slide
            component_code = await self._call_gemini_vision(image_data_url, prompt)

            if component_code:
                # Post-process: ensure original slide images are used
                if image_urls:
                    component_code = self._inject_original_images(component_code, image_urls)
                return {
                    "id": str(uuid.uuid4()),
                    "title": f"Slide {slide_idx + 1}",
                    "components": [
                        {
                            "id": str(uuid.uuid4()),
                            "type": "CustomComponent",
                            "props": {
                                "position": {"x": 0, "y": 0},
                                "width": 1920,
                                "height": 1080,
                                "x": 0,
                                "y": 0,
                                "zIndex": 1,
                                "opacity": 1,
                                "rotation": 0,
                                "render": component_code,  # Use 'render' not 'code'
                                "originalImageUrl": image_data_url  # Keep reference to original
                            }
                        }
                    ]
                }
            else:
                return self._create_image_fallback_slide(image_path, slide_idx)

        except Exception as e:
            logger.warning(f"[VisionPPTX] Vision AI error for slide {slide_idx + 1}: {e}")
            return self._create_image_fallback_slide(image_path, slide_idx)

    async def _call_gemini_vision(self, image_data_url: str, prompt: str) -> Optional[str]:
        """Call Gemini with vision capabilities using google.genai SDK."""
        try:
            from google import genai
            from google.genai import types

            # Get API key
            api_key = os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY")
            if not api_key:
                logger.error("[VisionPPTX] No GOOGLE_API_KEY or GEMINI_API_KEY found!")
                return None

            logger.info(f"[VisionPPTX] Calling Gemini vision API...")

            # Create client
            client = genai.Client(api_key=api_key)

            # Decode image from data URL
            header, b64_data = image_data_url.split(",", 1)
            image_bytes = base64.b64decode(b64_data)

            # Determine mime type from header
            mime_type = "image/png"
            if "image/jpeg" in header:
                mime_type = "image/jpeg"
            elif "image/webp" in header:
                mime_type = "image/webp"

            # Build contents with image and prompt
            contents = [
                types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
                prompt
            ]

            response = client.models.generate_content(
                model=VISION_IMPORT_MODEL,
                contents=contents,
                config=types.GenerateContentConfig(
                    temperature=0.1,  # Low temp for precise recreation
                    max_output_tokens=16384,  # Enough for detailed HTML
                )
            )

            if response and response.text:
                logger.info(f"[VisionPPTX] Got response: {len(response.text)} chars")
                # Extract code from response
                code = self._extract_code_from_response(response.text)
                if code:
                    logger.info(f"[VisionPPTX] Extracted code: {len(code)} chars")
                    return code
                else:
                    logger.warning("[VisionPPTX] Failed to extract code from response")
            else:
                logger.warning("[VisionPPTX] Empty response from Gemini")

        except Exception as e:
            logger.error(f"[VisionPPTX] Gemini vision error: {e}", exc_info=True)

        return None

    @staticmethod
    def _inject_original_images(html: str, image_urls: List[str]) -> str:
        """Inject original PPTX images into AI-generated HTML.

        AI sometimes generates placeholders instead of using the provided URLs.
        This replaces any placeholder <img> src with the original uploaded URLs.
        """
        import re
        from services.image.placeholder_detector import is_placeholder_src, is_bucket_url

        if not image_urls:
            return html

        img_pattern = re.compile(r'<img([^>]*)>', re.IGNORECASE)
        url_idx = 0

        def replace_img(match):
            nonlocal url_idx
            attrs = match.group(1)
            src_match = re.search(r'src=["\']([^"\']*)["\']', attrs, re.IGNORECASE)
            current_src = src_match.group(1) if src_match else ''

            if current_src and is_bucket_url(current_src):
                return match.group(0)

            if url_idx < len(image_urls) and (not current_src or is_placeholder_src(current_src)):
                new_url = image_urls[url_idx]
                url_idx += 1
                if src_match:
                    new_attrs = re.sub(r'src=["\'][^"\']*["\']', f'src="{new_url}"', attrs, count=1, flags=re.IGNORECASE)
                else:
                    new_attrs = f' src="{new_url}"' + attrs
                return f'<img{new_attrs}>'

            return match.group(0)

        result = img_pattern.sub(replace_img, html)
        if url_idx > 0:
            logger.info(f"[VisionPPTX] Injected {url_idx} original slide images into HTML")
        return result

    def _extract_code_from_response(self, response: str) -> str:
        """Extract HTML/CSS code from AI response."""
        import re

        response = response.strip()

        # Try to find HTML code block first
        html_block = re.search(r"```html\s*([\s\S]*?)```", response, re.IGNORECASE)
        if html_block:
            code = html_block.group(1).strip()
            # Ensure it's a complete HTML document
            if code.lower().startswith('<!doctype') or code.lower().startswith('<html'):
                return code

        # Try generic code block
        generic_block = re.search(r"```\s*([\s\S]*?)```", response, re.IGNORECASE)
        if generic_block:
            code = generic_block.group(1).strip()
            if code.lower().startswith('<!doctype') or code.lower().startswith('<html'):
                return code

        # Try to find complete HTML document directly in response
        html_doc = re.search(r"(<!DOCTYPE html[\s\S]*?</html>)", response, re.IGNORECASE)
        if html_doc:
            return html_doc.group(1).strip()

        # Try to find <html> tag if no DOCTYPE
        html_tag = re.search(r"(<html[\s\S]*?</html>)", response, re.IGNORECASE)
        if html_tag:
            return "<!DOCTYPE html>\n" + html_tag.group(1).strip()

        # If response starts with DOCTYPE or html, use as-is
        if response.lower().startswith('<!doctype') or response.lower().startswith('<html'):
            return response

        # Last resort: if it has HTML tags, wrap it
        if "<" in response and ">" in response:
            logger.warning("[VisionPPTX] Response doesn't look like complete HTML, using as-is")
            return response

        return response

    def _create_background_component(self) -> Dict[str, Any]:
        """Create default background component."""
        return {
            "id": str(uuid.uuid4()),
            "type": "Background",
            "props": {
                "position": {"x": 0, "y": 0},
                "width": 1920,
                "height": 1080,
                "x": 0,
                "y": 0,
                "zIndex": 0,
                "opacity": 1,
                "rotation": 0,
                "backgroundType": "solid",
                "backgroundColor": "#ffffffff"
            }
        }

    def _create_image_fallback_slide(self, image_path: str, slide_idx: int) -> Dict[str, Any]:
        """Create slide with original image as fallback."""
        # Read image as base64
        with open(image_path, "rb") as f:
            image_data = f.read()
        image_b64 = base64.b64encode(image_data).decode("utf-8")
        image_data_url = f"data:image/png;base64,{image_b64}"

        return {
            "id": str(uuid.uuid4()),
            "title": f"Slide {slide_idx + 1}",
            "components": [
                self._create_background_component(),
                {
                    "id": str(uuid.uuid4()),
                    "type": "Image",
                    "props": {
                        "position": {"x": 0, "y": 0},
                        "width": 1920,
                        "height": 1080,
                        "x": 0,
                        "y": 0,
                        "zIndex": 1,
                        "opacity": 1,
                        "rotation": 0,
                        "src": image_data_url,
                        "alt": f"Slide {slide_idx + 1}",
                        "objectFit": "contain"
                    }
                }
            ]
        }


# Convenience function
async def import_pptx_with_vision(file_path_or_bytes, filename: str = "presentation.pptx") -> Dict[str, Any]:
    """Import PPTX using vision-based approach."""
    importer = VisionPPTXImporter()

    if isinstance(file_path_or_bytes, str):
        return await importer.import_file(file_path_or_bytes)
    else:
        return await importer.import_bytes(file_path_or_bytes, filename)
