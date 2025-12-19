import asyncio
from typing import Dict, Any, List, Optional

from services.file_design_extractor import (
    FileDesignExtractor, FileIntent, SlideStyle, FileAnalysis,
    design_to_theme_context, content_to_outline_context
)
from setup_logging_optimized import get_logger

logger = get_logger(__name__)


async def analyze_files_for_presentation(files: List['FileAttachment']) -> Dict[str, Any]:
    """
    Analyze uploaded files using Anthropic Claude for presentation creation.
    Supports: images (vision), PDFs, Excel/CSV, PPTX, documents.
    Returns structured analysis that can be used in outline generation.
    """
    if not files:
        return {"success": False, "analyses": [], "combined_context": ""}

    try:
        import base64
        from io import BytesIO

        client, model = get_client(OUTLINE_AGENT_MODEL, wrap_with_instructor=False)
        analyses = []
        combined_context_parts = []

        for file in files:
            file_type = file.type.lower() if file.type else ""
            filename = file.name.lower()
            analysis = {
                "file_id": file.id,
                "filename": file.name,
                "file_type": "unknown",
                "summary": "",
                "key_insights": [],
                "suggested_slides": [],
                "extracted_data": None
            }

            try:
                # If content is missing but URL is available, fetch the content
                file_content = file.content
                if file_content:
                    logger.info(f"[OutlineAgent] File {file.name} has base64 content ({len(file_content)} chars)")
                elif file.url:
                    logger.info(f"[OutlineAgent] Fetching file content from URL for: {file.name}")
                    try:
                        import httpx
                        async with httpx.AsyncClient(timeout=30.0) as http_client:
                            response = await http_client.get(file.url)
                            if response.status_code == 200:
                                file_content = base64.b64encode(response.content).decode('utf-8')
                                logger.info(f"[OutlineAgent] Successfully fetched {len(file_content)} chars of content from URL")
                            else:
                                logger.warning(f"[OutlineAgent] Failed to fetch file from URL: {response.status_code}")
                    except Exception as fetch_err:
                        logger.warning(f"[OutlineAgent] Error fetching file from URL: {fetch_err}")
                else:
                    logger.warning(f"[OutlineAgent] ⚠️ File {file.name} has NO content and NO URL - cannot process!")

                # Determine file category
                is_image = file_type.startswith("image/")
                is_pdf = file_type == "application/pdf" or filename.endswith(".pdf")
                is_excel = "spreadsheet" in file_type or "excel" in file_type or filename.endswith((".xlsx", ".xls", ".csv"))
                is_pptx = "presentation" in file_type or filename.endswith((".pptx", ".ppt"))
                is_doc = "document" in file_type or "word" in file_type or filename.endswith((".doc", ".docx", ".txt", ".md"))

                content_blocks = []

                if is_image and file_content:
                    # Use Claude vision for images
                    analysis["file_type"] = "image"
                    content_blocks.append({
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": file_type,
                            "data": file_content
                        }
                    })
                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this image for creating a presentation. Provide:
1. **Description**: What does this image show?
2. **Key Elements**: Important subjects, text, data, or objects
3. **Suggested Use**: Which slide types would this fit? (intro, data visualization, team, product, etc.)
4. **Extracted Text**: Any visible text (OCR)
5. **Data/Charts**: If it contains charts/graphs, extract the key data points

Be concise but thorough. This will be used to create presentation slides."""
                    })

                elif is_pdf and file_content:
                    # Check PDF size - Claude has ~25MB limit for PDFs
                    pdf_size_mb = len(file_content) * 3 / 4 / 1024 / 1024  # base64 to bytes to MB
                    logger.info(f"[OutlineAgent] PDF size: {pdf_size_mb:.1f}MB")

                    if pdf_size_mb > 20:  # Use text extraction for large PDFs
                        logger.info(f"[OutlineAgent] PDF too large ({pdf_size_mb:.1f}MB), using text extraction")
                        analysis["file_type"] = "document"
                        try:
                            import pypdf
                            raw_bytes = base64.b64decode(file_content)
                            pdf_reader = pypdf.PdfReader(BytesIO(raw_bytes))
                            text_content = []
                            for i, page in enumerate(pdf_reader.pages[:50]):  # First 50 pages
                                page_text = page.extract_text() or ""
                                if page_text.strip():
                                    text_content.append(f"--- Page {i+1} ---\n{page_text}")
                            extracted_text = "\n\n".join(text_content)[:50000]  # Limit to 50k chars
                            content_blocks.append({
                                "type": "text",
                                "text": f"""Analyze this PDF content for creating a presentation:

{extracted_text}

Extract:
1. **Main Topic**: What is this document about?
2. **Key Points**: The most important facts, statistics, and insights
3. **Slide Suggestions**: What slides should be created from this content?

Focus on content that would make great presentation slides."""
                            })
                        except Exception as pdf_err:
                            logger.warning(f"[OutlineAgent] PDF text extraction failed: {pdf_err}")
                            analysis["summary"] = f"PDF too large ({pdf_size_mb:.1f}MB) and text extraction failed"
                            continue
                    else:
                        # Use Claude's native PDF understanding for smaller PDFs
                        analysis["file_type"] = "document"
                        content_blocks.append({
                            "type": "document",
                            "source": {
                                "type": "base64",
                                "media_type": "application/pdf",
                                "data": file_content
                            }
                        })
                        content_blocks.append({
                            "type": "text",
                            "text": f"""Analyze this PDF document for creating a presentation. Extract:
1. **Main Topic**: What is this document about?
2. **Key Points**: The most important facts, statistics, and insights (bullet points)
3. **Structure**: How is the content organized?
4. **Slide Suggestions**: What slides should be created from this content?
5. **Data/Charts**: Any numerical data that could be visualized

Focus on content that would make great presentation slides."""
                        })

                elif is_excel and file_content:
                    # Parse Excel/CSV data
                    analysis["file_type"] = "spreadsheet"
                    raw_bytes = base64.b64decode(file_content)
                    data_preview = ""
                    extracted_data = None

                    if filename.endswith(".csv"):
                        import csv
                        from io import StringIO
                        text = raw_bytes.decode('utf-8', errors='ignore')
                        reader = csv.reader(StringIO(text))
                        rows = list(reader)[:500]  # Up to 500 rows (for rolodex, contact lists, etc.)
                        if rows:
                            headers = rows[0]
                            data_rows = rows[1:] if len(rows) > 1 else []
                            total_rows = len(data_rows)
                            # Show more rows in preview for better context
                            data_preview = f"Headers: {', '.join(headers[:15])}\n"
                            data_preview += f"Total rows: {total_rows}\n\n"
                            # Show first 20 rows in preview
                            for i, row in enumerate(data_rows[:20]):
                                data_preview += f"Row {i+1}: {', '.join(row[:10])}\n"
                            if total_rows > 20:
                                data_preview += f"... and {total_rows - 20} more rows\n"
                            # Store ALL rows (up to 500) for slide generation
                            extracted_data = {"headers": headers, "all_rows": data_rows, "sample_rows": data_rows[:20], "total_rows": total_rows}
                    else:
                        try:
                            import openpyxl
                            wb = openpyxl.load_workbook(BytesIO(raw_bytes), read_only=True)
                            sheet = wb.active
                            rows = []
                            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                                if i >= 30: break
                                rows.append([str(c) if c is not None else "" for c in row])
                            wb.close()
                            if rows:
                                headers = rows[0]
                                data_rows = rows[1:]
                                data_preview = f"Headers: {', '.join(headers[:10])}\n"
                                for i, row in enumerate(data_rows[:10]):
                                    data_preview += f"Row {i+1}: {', '.join(row[:8])}\n"
                                extracted_data = {"headers": headers, "sample_rows": data_rows[:20], "total_rows": len(rows)-1}
                        except Exception as e:
                            data_preview = f"Could not parse Excel: {str(e)}"

                    analysis["extracted_data"] = extracted_data
                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this spreadsheet data for creating presentation slides:

{data_preview}

Provide:
1. **Data Type**: Is this a list of items (contacts, products, team members) OR analytical data (metrics, trends)?
2. **Slide Strategy**:
   - For LISTS (contacts, team, products): Recommend creating individual slides or grouped cards per item
   - For ANALYTICS: Recommend charts and summary slides
3. **Key Fields**: Which columns are most important to display?
4. **Suggested Layout**: Cards, table, grid, or chart format?
5. **Total Items**: How many items total, and how should they be organized into slides?

IMPORTANT: If this looks like a contact list, team roster, product catalog, or similar LIST data:
- Each item should get its own space on slides
- Suggest creating slides that show individual entries (like contact cards, team member profiles)
- Do NOT summarize the list into charts - DISPLAY the actual items"""
                    })

                elif is_pptx and file_content:
                    # Parse PPTX content using UniversalPPTXImporter for better design extraction
                    analysis["file_type"] = "presentation"
                    raw_bytes = base64.b64decode(file_content)
                    slides_content = ""
                    design_info = ""

                    try:
                        # Use the new universal importer for robust extraction
                        from services.universal_pptx_importer import UniversalPPTXImporter

                        importer = UniversalPPTXImporter()
                        # Await the async import directly (we're already in async context)
                        deck = await importer.import_bytes(raw_bytes, file.name)

                        # Extract design summary
                        design_summary = importer.get_design_summary()

                        # Build content summary
                        slides = deck.get("slides", [])
                        slides_content = f"Total slides: {len(slides)}\n\n"

                        for i, slide in enumerate(slides[:15]):  # First 15 slides
                            slides_content += f"Slide {i+1}: {slide.get('title', 'Untitled')}\n"
                            # Extract text from components
                            for comp in slide.get("components", []):
                                if comp.get("type") == "TiptapTextBlock":
                                    texts = comp.get("props", {}).get("texts", [])
                                    for t in texts[:3]:  # First 3 text segments
                                        text = t.get("text", "")[:200]
                                        if text.strip():
                                            slides_content += f"  - {text}\n"
                            slides_content += "\n"

                        # Build design info
                        if design_summary:
                            design_info = "\n\n**DESIGN INFORMATION:**\n"
                            if design_summary.get("theme_name"):
                                design_info += f"- Theme: {design_summary['theme_name']}\n"
                            if design_summary.get("primary_color"):
                                design_info += f"- Primary Color: {design_summary['primary_color']}\n"
                            if design_summary.get("secondary_color"):
                                design_info += f"- Secondary Color: {design_summary['secondary_color']}\n"
                            if design_summary.get("background_color"):
                                design_info += f"- Background Color: {design_summary['background_color']}\n"
                            if design_summary.get("text_color"):
                                design_info += f"- Text Color: {design_summary['text_color']}\n"
                            if design_summary.get("heading_font"):
                                design_info += f"- Heading Font: {design_summary['heading_font']}\n"
                            if design_summary.get("body_font"):
                                design_info += f"- Body Font: {design_summary['body_font']}\n"

                            # Store design data for later use
                            analysis["design_data"] = design_summary

                        # Store import stats
                        import_stats = deck.get("metadata", {}).get("import_stats", {})
                        analysis["import_stats"] = import_stats

                    except Exception as e:
                        logger.warning(f"[OutlineAgent] Universal importer failed, falling back to basic: {e}")
                        # Fallback to basic extraction
                        try:
                            from pptx import Presentation
                            prs = Presentation(BytesIO(raw_bytes))
                            slides_content = f"Total slides: {len(prs.slides)}\n\n"
                            for i, slide in enumerate(prs.slides[:15]):
                                slides_content += f"Slide {i+1}:\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        slides_content += f"  - {shape.text[:300]}\n"
                                slides_content += "\n"
                        except Exception as e2:
                            slides_content = f"Could not parse PPTX: {str(e2)}"

                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this PowerPoint presentation:

{slides_content}{design_info}

Provide:
1. **Overview**: Main topic and purpose
2. **Structure**: How is it organized?
3. **Key Content**: Most important points from each slide
4. **Design Notes**: Describe the visual design style, colors, fonts, and layout patterns
5. **Design Replication**: If user wants to replicate this design, what are the key design elements?"""
                    })

                elif is_doc and file_content:
                    # Parse document text
                    analysis["file_type"] = "document"
                    raw_bytes = base64.b64decode(file_content)
                    text_content = raw_bytes.decode('utf-8', errors='ignore')[:15000]  # First 15k chars

                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this document for creating a presentation:

{text_content}

Provide:
1. **Main Topic**: What is this about?
2. **Key Points**: Most important facts and insights
3. **Slide Suggestions**: What slides should be created from this?
4. **Data/Statistics**: Any numbers that should be highlighted"""
                    })

                else:
                    # Unknown file type
                    analysis["summary"] = f"File type '{file_type}' not fully supported for deep analysis"
                    analyses.append(analysis)
                    continue

                # Call AI for analysis - support both Anthropic and Gemini
                if content_blocks:
                    is_gemini = 'gemini' in model.lower()

                    if is_gemini:
                        # Use Gemini API pattern
                        from google import genai

                        # Convert content blocks to Gemini format
                        gemini_parts = []
                        text_parts = []
                        system_prompt = "You are an expert at analyzing content for presentation creation. Be concise and focus on actionable insights.\n\n"

                        for block in content_blocks:
                            if block.get("type") == "text":
                                text_parts.append(block["text"])
                            elif block.get("type") == "image":
                                # Gemini supports inline images via Part
                                try:
                                    image_part = genai.types.Part.from_bytes(
                                        data=base64.b64decode(block["source"]["data"]),
                                        mime_type=block["source"]["media_type"]
                                    )
                                    gemini_parts.append(image_part)
                                except Exception as img_err:
                                    logger.warning(f"[OutlineAgent] Image conversion for Gemini failed: {img_err}")
                            elif block.get("type") == "document":
                                # For PDFs, extract text using pypdf since Gemini doesn't support document blocks directly
                                try:
                                    import pypdf
                                    raw_bytes = base64.b64decode(block["source"]["data"])
                                    pdf_reader = pypdf.PdfReader(BytesIO(raw_bytes))
                                    pdf_text_content = []
                                    for i, page in enumerate(pdf_reader.pages[:30]):
                                        page_text = page.extract_text() or ""
                                        if page_text.strip():
                                            pdf_text_content.append(f"--- Page {i+1} ---\n{page_text}")
                                    extracted_text = "\n\n".join(pdf_text_content)[:40000]
                                    text_parts.append(f"PDF Content:\n{extracted_text}")
                                except Exception as pdf_err:
                                    logger.warning(f"[OutlineAgent] PDF extraction for Gemini failed: {pdf_err}")
                                    text_parts.append("(PDF content could not be extracted)")

                        # Build final content: text first, then images
                        final_content = [system_prompt + "\n".join(text_parts)]
                        final_content.extend(gemini_parts)

                        gemini_client = genai.Client()
                        result = gemini_client.models.generate_content(
                            model="gemini-2.0-flash",
                            contents=final_content,
                            config=genai.types.GenerateContentConfig(
                                temperature=0.3,
                                max_output_tokens=2000,
                            )
                        )
                        analysis_text = result.text if result.text else ""
                        logger.info(f"[OutlineAgent] Gemini file analysis complete: {len(analysis_text)} chars")
                    else:
                        # Use Anthropic API pattern
                        messages = [{"role": "user", "content": content_blocks}]
                        result = client.messages.create(
                            model=model,
                            messages=messages,
                            max_tokens=2000,
                            system="You are an expert at analyzing content for presentation creation. Be concise and focus on actionable insights."
                        )
                        analysis_text = result.content[0].text if result.content else ""

                    analysis["summary"] = analysis_text

                    # Extract key insights as bullet points
                    lines = analysis_text.split('\n')
                    key_insights = []
                    for line in lines:
                        line = line.strip()
                        if line.startswith(('- ', '• ', '* ', '1.', '2.', '3.', '4.', '5.')):
                            key_insights.append(line.lstrip('-•* 0123456789.').strip())
                    analysis["key_insights"] = key_insights[:10]  # Top 10 insights

                    # Add to combined context
                    combined_context_parts.append(f"=== {file.name} ({analysis['file_type']}) ===\n{analysis_text}\n")

                    # For spreadsheets with list data, include ALL the actual data rows
                    # so slide generation can create slides for each item (contacts, team members, etc.)
                    if analysis.get("extracted_data") and analysis["file_type"] == "spreadsheet":
                        ed = analysis["extracted_data"]
                        all_rows = ed.get("all_rows", ed.get("sample_rows", []))
                        headers = ed.get("headers", [])
                        if all_rows and len(all_rows) > 0:
                            # Include full data for slide generation
                            data_section = f"\n=== FULL DATA FROM {file.name} (for slide generation) ===\n"
                            data_section += f"Headers: {', '.join(headers)}\n"
                            data_section += f"Total items: {len(all_rows)}\n\n"
                            # Include all rows (up to 200 for context limits)
                            for i, row in enumerate(all_rows[:200]):
                                data_section += f"Item {i+1}: {', '.join(str(cell) for cell in row)}\n"
                            if len(all_rows) > 200:
                                data_section += f"... ({len(all_rows) - 200} more items)\n"
                            combined_context_parts.append(data_section)

            except Exception as file_err:
                logger.error(f"[OutlineAgent] Error analyzing file {file.name}: {file_err}")
                analysis["summary"] = f"Error analyzing file: {str(file_err)}"

            analyses.append(analysis)

        return {
            "success": True,
            "analyses": analyses,
            "combined_context": "\n".join(combined_context_parts),
            "file_count": len(analyses)
        }

    except Exception as e:
        logger.error(f"[OutlineAgent] File analysis failed: {e}")
        return {
            "success": False,
            "analyses": [],
            "combined_context": "",
            "error": str(e)
        }

async def enhanced_file_analysis(
    files: List['FileAttachment'],
    user_message: str = ""
) -> Dict[str, Any]:
    """
    Enhanced file analysis that detects:
    - User intent (recreate, use design only, use content only, use both)
    - Slide style preference (traditional vs interactive)
    - Extracts design elements for theme matching
    - Extracts content for outline generation
    """
    if not files:
        return {
            "success": False,
            "intent": FileIntent.USE_BOTH.value,
            "slide_style": SlideStyle.AUTO.value,
            "design_context": None,
            "content_context": "",
            "file_analyses": []
        }

    try:
        extractor = FileDesignExtractor()
        file_analyses = []
        design_context = None
        content_parts = []
        all_extracted_images = []  # Collect images from all files
        all_slide_screenshots = []  # Collect slide screenshots for visual design reference

        # Detect intent and style from user message (same for all files)
        intent = extractor._detect_intent(user_message)
        slide_style = extractor._detect_slide_style(user_message)

        logger.info(f"[EnhancedAnalysis] Detected intent: {intent.value}, style: {slide_style.value}")

        for file in files:
            # If content is missing but URL is available, fetch the content
            file_content = file.content
            if not file_content and file.url:
                logger.info(f"[EnhancedAnalysis] Fetching file content from URL for: {file.name}")
                try:
                    import httpx
                    import base64
                    async with httpx.AsyncClient(timeout=30.0) as http_client:
                        response = await http_client.get(file.url)
                        if response.status_code == 200:
                            file_content = base64.b64encode(response.content).decode('utf-8')
                            logger.info(f"[EnhancedAnalysis] Successfully fetched content from URL")
                except Exception as fetch_err:
                    logger.warning(f"[EnhancedAnalysis] Error fetching file from URL: {fetch_err}")

            if not file_content:
                logger.warning(f"[EnhancedAnalysis] ⚠️ Skipping {file.name} - no content available")
                continue

            try:
                analysis = await extractor.analyze_file(
                    file_content=file_content,
                    filename=file.name,
                    file_type=file.type or "",
                    user_message=user_message
                )

                file_analyses.append({
                    "filename": analysis.filename,
                    "file_type": analysis.file_type,
                    "intent": analysis.intent.value,
                    "slide_style": analysis.slide_style.value,
                    "confidence": analysis.confidence,
                    "notes": analysis.analysis_notes
                })

                # Use the first file's design for theme context (usually the main reference)
                if design_context is None and intent in (FileIntent.USE_DESIGN_ONLY, FileIntent.USE_BOTH, FileIntent.RECREATE_EXACT):
                    design_context = design_to_theme_context(analysis.design)
                    logger.info(f"[EnhancedAnalysis] Extracted design: {design_context}")

                # Build content context
                if intent in (FileIntent.USE_CONTENT_ONLY, FileIntent.USE_BOTH, FileIntent.RECREATE_EXACT):
                    content_str = content_to_outline_context(analysis.content)
                    if content_str:
                        content_parts.append(f"=== {file.name} ===\n{content_str}")

                # Collect extracted images from uploaded files (logos, photos, etc.)
                if analysis.extracted_images:
                    all_extracted_images.extend(analysis.extracted_images)
                    logger.info(f"[EnhancedAnalysis] Collected {len(analysis.extracted_images)} images from {file.name}")

                # Collect slide screenshots for visual design reference
                if analysis.slide_screenshots:
                    all_slide_screenshots.extend(analysis.slide_screenshots)
                    logger.info(f"[EnhancedAnalysis] Collected {len(analysis.slide_screenshots)} slide screenshots from {file.name}")

            except Exception as e:
                logger.error(f"[EnhancedAnalysis] Error analyzing {file.name}: {e}")
                continue

        if all_extracted_images:
            logger.info(f"[EnhancedAnalysis] Total extracted images: {len(all_extracted_images)}")
        if all_slide_screenshots:
            logger.info(f"[EnhancedAnalysis] Total slide screenshots for visual reference: {len(all_slide_screenshots)}")

        return {
            "success": True,
            "intent": intent.value,
            "slide_style": slide_style.value,
            "design_context": design_context,
            "content_context": "\n\n".join(content_parts) if content_parts else "",
            "file_analyses": file_analyses,
            "extracted_images": all_extracted_images,  # Images from uploaded PPTX/PDF
            "slide_screenshots": all_slide_screenshots  # Base64 screenshots for visual design replication
        }

    except Exception as e:
        logger.error(f"[EnhancedAnalysis] Failed: {e}")
        return {
            "success": False,
            "intent": FileIntent.USE_BOTH.value,
            "slide_style": SlideStyle.AUTO.value,
            "design_context": None,
            "content_context": "",
            "file_analyses": [],
            "extracted_images": [],
            "slide_screenshots": [],
            "error": str(e)
        }
