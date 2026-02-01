"""
File Analysis API — provider-agnostic (Gemini, Claude, OpenAI)

Supports: images, PDFs, Excel, CSV, documents, presentations.
PDFs are limited to 3 pages for unauthenticated (free-preview) requests.
"""

import os
import base64
import logging
import mimetypes
from typing import List, Dict, Any, Optional
from pydantic import BaseModel, Field
from fastapi import APIRouter, HTTPException
from io import BytesIO

from agents.config import FILE_ANALYSIS_MODEL_FAST

logger = logging.getLogger(__name__)

router = APIRouter()

# ── Max free-preview pages for PDFs ──────────────────────────────────────────
FREE_PDF_PAGE_LIMIT = 3


# ═══════════════════════════════════════════════════════════════════════════════
# MODELS
# ═══════════════════════════════════════════════════════════════════════════════

class FileInput(BaseModel):
    """Input file for analysis"""
    id: str = Field(description="Unique file ID")
    name: str = Field(description="Original filename")
    type: str = Field(description="MIME type")
    content: Optional[str] = Field(default=None, description="Base64 encoded content")
    url: Optional[str] = Field(default=None, description="URL to the file")
    size: Optional[int] = Field(default=None, description="File size in bytes")


class FileAnalysisRequest(BaseModel):
    """Request for file analysis"""
    files: List[FileInput] = Field(description="Files to analyze")
    context: Optional[str] = Field(default=None, description="Additional context for analysis")
    chat_history: Optional[List[Dict[str, str]]] = Field(default=None, description="Chat history for context")


class FileAnalysisResult(BaseModel):
    """Result of file analysis"""
    file_id: str
    filename: str
    file_type: str
    analysis: str
    summary: str
    suggestions: List[str] = Field(default_factory=list)
    extracted_data: Optional[Dict[str, Any]] = None
    preview_url: Optional[str] = None


class FileAnalysisResponse(BaseModel):
    """Response from file analysis"""
    success: bool
    results: List[FileAnalysisResult]
    combined_analysis: str
    message: str
    # PDF page-limit metadata (present when a PDF was truncated)
    total_pages: Optional[int] = None
    pages_analyzed: Optional[int] = None


# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _call_llm(model_name: str, messages: list, max_tokens: int = 1500) -> str:
    """Call any LLM in a provider-agnostic way via the central invoke()."""
    from agents.ai.clients import get_client, invoke

    client, actual_model = get_client(model_name, wrap_with_instructor=False)
    result = invoke(client, actual_model, messages, max_tokens=max_tokens)
    return result or "No analysis available"


def _get_mime_type(filename: str, provided_type: str) -> str:
    if provided_type and provided_type != 'application/octet-stream':
        return provided_type
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or 'application/octet-stream'


def _is_image(mime_type: str) -> bool:
    return mime_type.startswith('image/')


def _is_document(mime_type: str) -> bool:
    doc_types = [
        'application/pdf',
        'application/msword',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'text/plain', 'text/markdown', 'text/html',
    ]
    return mime_type in doc_types


def _is_spreadsheet(mime_type: str, filename: str) -> bool:
    spreadsheet_types = [
        'application/vnd.ms-excel',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'text/csv',
    ]
    return mime_type in spreadsheet_types or filename.lower().endswith(('.csv', '.xlsx', '.xls'))


def _is_presentation(mime_type: str, filename: str) -> bool:
    ppt_types = [
        'application/vnd.ms-powerpoint',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    ]
    return mime_type in ppt_types or filename.lower().endswith(('.ppt', '.pptx'))


def _limit_pdf_pages(base64_content: str, max_pages: int = FREE_PDF_PAGE_LIMIT):
    """Limit a PDF to *max_pages* pages.

    Returns (limited_base64, total_pages, pages_kept).
    If the PDF has <= max_pages, returns the original content unchanged.
    """
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        try:
            from PyPDF2 import PdfReader, PdfWriter
        except ImportError:
            logger.warning("pypdf / PyPDF2 not installed — skipping page limit")
            return base64_content, None, None

    try:
        raw_bytes = base64.b64decode(base64_content)
        reader = PdfReader(BytesIO(raw_bytes))
        total_pages = len(reader.pages)

        if total_pages <= max_pages:
            return base64_content, total_pages, total_pages

        writer = PdfWriter()
        for i in range(max_pages):
            writer.add_page(reader.pages[i])

        output = BytesIO()
        writer.write(output)
        limited_b64 = base64.b64encode(output.getvalue()).decode('utf-8')
        return limited_b64, total_pages, max_pages
    except Exception as e:
        logger.warning(f"Could not limit PDF pages: {e}")
        return base64_content, None, None


# ═══════════════════════════════════════════════════════════════════════════════
# ANALYSERS (provider-agnostic)
# ═══════════════════════════════════════════════════════════════════════════════

async def _analyze_image(file_input: FileInput, context: str = "") -> FileAnalysisResult:
    """Analyze an image using the configured LLM's vision capabilities."""
    from agents.config import OUTLINE_AGENT_MODEL

    try:
        mime_type = _get_mime_type(file_input.name, file_input.type)
        content_blocks: list = []

        if file_input.content:
            content_blocks.append({
                "type": "image",
                "source": {"type": "base64", "media_type": mime_type, "data": file_input.content},
            })
        elif file_input.url:
            import httpx
            async with httpx.AsyncClient() as http_client:
                resp = await http_client.get(file_input.url)
                if resp.status_code != 200:
                    raise ValueError(f"Failed to fetch image: HTTP {resp.status_code}")
                image_data = base64.b64encode(resp.content).decode('utf-8')
                content_blocks.append({
                    "type": "image",
                    "source": {"type": "base64", "media_type": mime_type, "data": image_data},
                })

        prompt = (
            "Analyze this image for use in a presentation. Provide:\n"
            "1. **Description**: What the image shows\n"
            "2. **Key Elements**: Important elements, text, data\n"
            "3. **Quality Assessment**: Resolution suitability for presentations\n"
            "4. **Suggested Use**: How to use this in a presentation\n"
            "5. **Extracted Text**: Any visible text (OCR)\n"
            "6. **Data/Charts**: If charts/graphs, extract key figures\n"
        )
        if context:
            prompt += f"\nAdditional context: {context}"

        content_blocks.append({"type": "text", "text": prompt})
        messages = [{"role": "user", "content": content_blocks}]

        analysis_text = _call_llm(OUTLINE_AGENT_MODEL, messages, max_tokens=1500)
        summary = analysis_text.split('\n\n')[0][:200]

        return FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="image",
            analysis=analysis_text, summary=summary,
            suggestions=["Use as slide background", "Add as supporting visual", "Use in image gallery"],
            preview_url=file_input.url,
        )
    except Exception as e:
        logger.error(f"Error analyzing image {file_input.name}: {e}")
        return FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="image",
            analysis=f"Error analyzing image: {e}", summary="Analysis failed", suggestions=[],
        )


async def _analyze_document(file_input: FileInput, context: str = "") -> tuple[FileAnalysisResult, Optional[int], Optional[int]]:
    """Analyze a document (PDF, text, etc.).

    Returns (result, total_pages, pages_analyzed).  Pages info is only set for PDFs.
    """
    try:
        text_content = ""
        mime_type = _get_mime_type(file_input.name, file_input.type)
        total_pages = None
        pages_analyzed = None

        if file_input.content:
            try:
                raw_bytes = base64.b64decode(file_input.content)

                if mime_type == 'application/pdf':
                    # Limit to FREE_PDF_PAGE_LIMIT pages
                    limited_b64, total_pages, pages_analyzed = _limit_pdf_pages(file_input.content)

                    truncation_note = ""
                    if total_pages and pages_analyzed and total_pages > pages_analyzed:
                        truncation_note = (
                            f"\n\nNote: This PDF has {total_pages} pages. "
                            f"Only the first {pages_analyzed} pages are included in this preview."
                        )

                    messages = [{
                        "role": "user",
                        "content": [
                            {
                                "type": "document",
                                "source": {
                                    "type": "base64",
                                    "media_type": "application/pdf",
                                    "data": limited_b64,
                                },
                            },
                            {
                                "type": "text",
                                "text": (
                                    "Analyze this document for creating a presentation. Provide:\n"
                                    "1. **Summary**: Main topic and key points\n"
                                    "2. **Key Facts**: Important statistics, dates, or figures\n"
                                    "3. **Structure**: How the document is organized\n"
                                    "4. **Presentation Suggestions**: Key slides that could be created\n"
                                    + (f"\nAdditional context: {context}" if context else "")
                                    + "\n\nBe concise and focus on presentation-worthy content."
                                ),
                            },
                        ],
                    }]

                    analysis_text = _call_llm(FILE_ANALYSIS_MODEL_FAST, messages, max_tokens=2000)
                    if truncation_note:
                        analysis_text += truncation_note

                else:
                    # Plain text documents
                    text_content = raw_bytes.decode('utf-8', errors='ignore')[:10000]

                    messages = [{
                        "role": "user",
                        "content": (
                            "Analyze this document content for creating a presentation:\n\n"
                            f"---\n{text_content}\n---\n\n"
                            "Provide:\n"
                            "1. **Summary**: Main topic and key points\n"
                            "2. **Key Facts**: Important statistics, dates, or figures\n"
                            "3. **Presentation Suggestions**: Key slides that could be created\n"
                            + (f"\nAdditional context: {context}" if context else "")
                            + "\n\nBe concise."
                        ),
                    }]

                    analysis_text = _call_llm(FILE_ANALYSIS_MODEL_FAST, messages, max_tokens=1500)

            except Exception as decode_err:
                logger.error(f"Error decoding document content: {decode_err}")
                analysis_text = f"Error decoding document: {decode_err}"
        else:
            analysis_text = "No document content provided"

        summary = analysis_text.split('\n\n')[0][:200]

        suggestions = ["Extract key points for slides", "Create outline from document", "Use as reference material"]
        if total_pages and pages_analyzed and total_pages > pages_analyzed:
            suggestions.insert(0, f"Sign up free to convert all {total_pages} pages")

        result = FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="document",
            analysis=analysis_text, summary=summary, suggestions=suggestions,
            extracted_data={
                "total_pages": total_pages,
                "pages_analyzed": pages_analyzed,
            } if total_pages else None,
        )
        return result, total_pages, pages_analyzed

    except Exception as e:
        logger.error(f"Error analyzing document {file_input.name}: {e}")
        result = FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="document",
            analysis=f"Error analyzing document: {e}", summary="Analysis failed", suggestions=[],
        )
        return result, None, None


async def _analyze_spreadsheet(file_input: FileInput, context: str = "") -> FileAnalysisResult:
    """Analyze a spreadsheet (CSV, Excel)."""
    try:
        data_preview = ""
        extracted_data = None

        if file_input.content:
            raw_bytes = base64.b64decode(file_input.content)
            filename_lower = file_input.name.lower()

            if filename_lower.endswith('.csv'):
                import csv
                from io import StringIO
                text_content = raw_bytes.decode('utf-8', errors='ignore')
                reader = csv.reader(StringIO(text_content))
                rows = list(reader)[:20]

                if rows:
                    headers = rows[0]
                    data_rows = rows[1:]
                    data_preview = f"Headers: {', '.join(headers)}\nSample data ({len(data_rows)} rows shown):\n"
                    for row in data_rows[:5]:
                        data_preview += f"  {', '.join(row[:5])}\n"
                    extracted_data = {"headers": headers, "sample_rows": data_rows[:10], "total_rows": len(rows) - 1}

            elif filename_lower.endswith(('.xlsx', '.xls')):
                try:
                    import openpyxl

                    wb = openpyxl.load_workbook(BytesIO(raw_bytes), read_only=True)
                    sheet = wb.active
                    rows = []
                    for i, row in enumerate(sheet.iter_rows(values_only=True)):
                        if i >= 20:
                            break
                        rows.append([str(cell) if cell is not None else "" for cell in row])

                    if rows:
                        headers = rows[0]
                        data_rows = rows[1:]
                        data_preview = f"Headers: {', '.join(headers[:10])}\nSample data:\n"
                        for row in data_rows[:5]:
                            data_preview += f"  {', '.join(row[:5])}\n"
                        extracted_data = {"headers": headers[:20], "sample_rows": data_rows[:10], "total_rows": sheet.max_row - 1 if sheet.max_row else 0}
                    wb.close()
                except ImportError:
                    data_preview = "Excel file detected but openpyxl not available"
                except Exception as excel_err:
                    data_preview = f"Error reading Excel: {excel_err}"

        messages = [{
            "role": "user",
            "content": (
                "Analyze this spreadsheet data for creating presentation charts/visualizations:\n\n"
                f"{data_preview}\n\n"
                "Provide:\n"
                "1. **Data Summary**: What does this data represent?\n"
                "2. **Key Insights**: Notable trends, patterns, or statistics\n"
                "3. **Chart Recommendations**: Best chart types (bar, line, pie, etc.)\n"
                "4. **Presentation Ideas**: How to present this data effectively\n"
                + (f"\nAdditional context: {context}" if context else "")
                + "\n\nBe specific about visualization recommendations."
            ),
        }]

        analysis_text = _call_llm(FILE_ANALYSIS_MODEL_FAST, messages, max_tokens=1500)
        summary = analysis_text.split('\n\n')[0][:200]

        return FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="spreadsheet",
            analysis=analysis_text, summary=summary,
            suggestions=["Create bar chart", "Create line chart", "Create pie chart", "Add as data table"],
            extracted_data=extracted_data,
        )
    except Exception as e:
        logger.error(f"Error analyzing spreadsheet {file_input.name}: {e}")
        return FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="spreadsheet",
            analysis=f"Error analyzing spreadsheet: {e}", summary="Analysis failed", suggestions=[],
        )


async def _analyze_presentation(file_input: FileInput, context: str = "") -> FileAnalysisResult:
    """Analyze a PowerPoint file."""
    try:
        slides_info = ""

        if file_input.content:
            raw_bytes = base64.b64decode(file_input.content)
            try:
                from pptx import Presentation

                prs = Presentation(BytesIO(raw_bytes))
                slides_info = f"Total slides: {len(prs.slides)}\n\n"
                for i, slide in enumerate(prs.slides[:10]):
                    slides_info += f"Slide {i+1}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides_info += f"  - {shape.text[:200]}\n"
                    slides_info += "\n"
            except ImportError:
                slides_info = "PowerPoint file detected but python-pptx not available"
            except Exception as pptx_err:
                slides_info = f"Error reading PowerPoint: {pptx_err}"

        messages = [{
            "role": "user",
            "content": (
                f"Analyze this PowerPoint presentation:\n\n{slides_info}\n\n"
                "Provide:\n"
                "1. **Overview**: Main topic and purpose\n"
                "2. **Structure**: How it is organized\n"
                "3. **Key Points**: Main takeaways\n"
                "4. **Suggestions**: How to improve or adapt this content\n"
                + (f"\nAdditional context: {context}" if context else "")
            ),
        }]

        analysis_text = _call_llm(FILE_ANALYSIS_MODEL_FAST, messages, max_tokens=1500)
        summary = analysis_text.split('\n\n')[0][:200]

        return FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="presentation",
            analysis=analysis_text, summary=summary,
            suggestions=["Import slides", "Extract content for new deck", "Use as reference"],
        )
    except Exception as e:
        logger.error(f"Error analyzing presentation {file_input.name}: {e}")
        return FileAnalysisResult(
            file_id=file_input.id, filename=file_input.name, file_type="presentation",
            analysis=f"Error analyzing presentation: {e}", summary="Analysis failed", suggestions=[],
        )


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN ANALYSIS ENTRYPOINT
# ═══════════════════════════════════════════════════════════════════════════════

async def analyze_files(request: FileAnalysisRequest) -> FileAnalysisResponse:
    """Analyze multiple files. Routes each file to the appropriate analyser."""
    results = []
    response_total_pages = None
    response_pages_analyzed = None

    for file_input in request.files:
        mime_type = _get_mime_type(file_input.name, file_input.type)
        context = request.context or ""
        logger.info(f"Analyzing file: {file_input.name} (type: {mime_type})")

        if _is_image(mime_type):
            result = await _analyze_image(file_input, context)
        elif _is_spreadsheet(mime_type, file_input.name):
            result = await _analyze_spreadsheet(file_input, context)
        elif _is_presentation(mime_type, file_input.name):
            result = await _analyze_presentation(file_input, context)
        elif _is_document(mime_type):
            result, total_pages, pages_analyzed = await _analyze_document(file_input, context)
            if total_pages is not None:
                response_total_pages = total_pages
                response_pages_analyzed = pages_analyzed
        else:
            result = FileAnalysisResult(
                file_id=file_input.id, filename=file_input.name, file_type="unknown",
                analysis=f"File type '{mime_type}' — basic analysis only. This file type may not be fully supported.",
                summary=f"Uploaded: {file_input.name}",
                suggestions=["Try converting to a supported format"],
            )

        results.append(result)

    # Combined analysis
    combined = ""
    if len(results) > 1:
        combined = f"Analyzed {len(results)} files:\n\n"
        for r in results:
            combined += f"**{r.filename}** ({r.file_type}): {r.summary}\n\n"
    elif results:
        combined = results[0].analysis

    return FileAnalysisResponse(
        success=True,
        results=results,
        combined_analysis=combined,
        message=f"Successfully analyzed {len(results)} file(s)",
        total_pages=response_total_pages,
        pages_analyzed=response_pages_analyzed,
    )


# ═══════════════════════════════════════════════════════════════════════════════
# ENDPOINTS
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/analyze", response_model=FileAnalysisResponse)
async def analyze_files_endpoint(request: FileAnalysisRequest):
    """Analyze uploaded files (images, PDFs, docs, spreadsheets, presentations)."""
    try:
        return await analyze_files(request)
    except Exception as e:
        logger.error(f"Error in file analysis endpoint: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/chat-with-files")
async def chat_with_files_endpoint(request: Dict[str, Any]):
    """Chat endpoint that can understand uploaded files."""
    try:
        files = request.get("files", [])
        message = request.get("message", "")
        chat_history = request.get("chat_history", [])

        # Analyse non-image files first
        file_analyses = []
        content_blocks: list = []

        for file_data in files:
            file_input = FileInput(**file_data)
            mime_type = _get_mime_type(file_input.name, file_input.type)

            if _is_image(mime_type) and file_input.content:
                content_blocks.append({
                    "type": "image",
                    "source": {"type": "base64", "media_type": mime_type, "data": file_input.content},
                })
                file_analyses.append(f"[Image: {file_input.name}]")
            else:
                analysis_request = FileAnalysisRequest(files=[file_input], context=message)
                analysis_result = await analyze_files(analysis_request)
                if analysis_result.results:
                    file_analyses.append(f"[{file_input.name}]: {analysis_result.results[0].analysis}")

        # Build user message
        user_message = message
        if file_analyses and not any(_is_image(_get_mime_type(f.get("name", ""), f.get("type", ""))) for f in files):
            user_message = f"{message}\n\nAttached files analysis:\n" + "\n".join(file_analyses)

        content_blocks.append({"type": "text", "text": user_message or "Please analyze these files."})

        # Build messages with history + system instruction
        from agents.config import OUTLINE_AGENT_MODEL
        messages: list = [
            {"role": "system", "content": "You are a helpful presentation assistant. When analyzing files, provide clear, actionable insights for creating presentations. Be concise but thorough."},
        ]
        for msg in chat_history[-10:]:
            messages.append({"role": msg.get("role", "user"), "content": msg.get("content", "")})
        messages.append({"role": "user", "content": content_blocks})

        response_text = _call_llm(OUTLINE_AGENT_MODEL, messages, max_tokens=2000)

        return {"success": True, "response": response_text, "file_analyses": file_analyses}

    except Exception as e:
        logger.error(f"Error in chat-with-files: {e}")
        return {"success": False, "error": str(e), "response": f"Sorry, I encountered an error analyzing your files: {e}"}
