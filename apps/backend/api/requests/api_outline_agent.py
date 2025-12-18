"""
Outline Generation Agent - Conversational AI for creating presentation outlines.

Supports both Anthropic (Claude) and Google (Gemini) models with tool calling for web search.
Model is configured via OUTLINE_AGENT_MODEL in agents/config.py.
"""
import os
import logging
import json
import re
import asyncio
from typing import Dict, Any, List, Optional, AsyncGenerator
from datetime import datetime

from fastapi import APIRouter, HTTPException, Depends
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

from agents.ai.clients import get_client, invoke
from agents.config import OUTLINE_AGENT_MODEL
from services.supabase_auth_service import get_auth_service
from api.requests.api_auth import get_auth_header
from setup_logging_optimized import get_logger

# Import Gemini types for function calling
try:
    from google import genai
    from google.genai import types as genai_types
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    genai = None
    genai_types = None
from services.file_design_extractor import (
    FileDesignExtractor, FileIntent, SlideStyle, FileAnalysis,
    design_to_theme_context, content_to_outline_context
)

logger = get_logger(__name__)


# Tool definition for web search - model decides when to use it
SEARCH_TOOL = {
    "name": "web_search",
    "description": """Search the web for current information, facts, and data.

✅ USE SEARCH when:
1. User needs FACTUAL DATA you don't have: company financials, market stats, current events, scientific facts
2. User mentions a topic that requires REAL INFORMATION: courses, lectures, research, analysis
3. User implies they want you to FIND or RESEARCH something (e.g., "pick a lecture online", "find examples")
4. Creating content about SPECIFIC TOPICS where accuracy matters (science, business, technology, academics)
5. User asks about a KNOWN ENTITY (company, university, course) and needs accurate facts about it

🚫 DO NOT search when:
- User uploaded files - USE THEIR CONTENT INSTEAD
- User already provided the specific content they want
- File analysis already contains the needed information

🚫🚫 NEVER SEARCH FOR BRAND COLORS, LOGOS, OR FONTS - the theme system handles branding separately.

💡 DEFAULT TO SEARCHING: If you're about to generate content on a topic and you're not 100% certain of the facts,
SEARCH FIRST. It's better to search and get accurate info than to make up content.
Academic topics, scientific content, business data - always search to ensure accuracy.""",
    "input_schema": {
        "type": "object",
        "properties": {
            "query": {
                "type": "string",
                "description": "The search query"
            }
        },
        "required": ["query"]
    }
}

# Gemini function declaration for web search
def get_gemini_search_tool():
    """Get Gemini-compatible function declaration for web search."""
    if not GEMINI_AVAILABLE:
        return None
    return genai_types.Tool(
        function_declarations=[
            genai_types.FunctionDeclaration(
                name="web_search",
                description=SEARCH_TOOL["description"],
                parameters=genai_types.Schema(
                    type=genai_types.Type.OBJECT,
                    properties={
                        "query": genai_types.Schema(
                            type=genai_types.Type.STRING,
                            description="The search query"
                        )
                    },
                    required=["query"]
                )
            )
        ]
    )


def is_gemini_model(model: str) -> bool:
    """Check if the model is a Gemini model."""
    return model.startswith("gemini") if model else False


async def research_with_perplexity(query: str) -> Dict[str, Any]:
    """
    Use Perplexity Sonar to research a topic. Simple and direct - let the AI do the thinking.
    """
    try:
        # Check if API key is available
        if not (os.getenv("PPLX_API_KEY") or os.getenv("PERPLEXITY_API_KEY")):
            logger.warning("[OutlineAgent] Perplexity API key not set - skipping research")
            return {
                "success": False,
                "content": None,
                "citations": [],
                "error": "PPLX_API_KEY not configured"
            }

        client, model = get_client("perplexity-sonar", wrap_with_instructor=False)

        # Simple prompt - Perplexity is smart, let it figure out what's needed
        system_prompt = """You are a research assistant helping create a presentation.
Provide accurate, current information with specific facts, numbers, and statistics.
Always cite your sources. Focus on what would be useful for presentation slides."""

        # Run Perplexity call with timeout to prevent hanging
        response = await asyncio.wait_for(
            asyncio.to_thread(
                client.chat.completions.create,
                model="sonar",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": query}
                ],
                max_tokens=3000,
                extra_body={
                    "return_citations": True,
                    "search_recency_filter": "month"
                }
            ),
            timeout=60.0  # 60 second timeout per search
        )

        content = response.choices[0].message.content

        # Extract citations
        citations = []
        if hasattr(response, 'citations'):
            for cit in response.citations:
                if isinstance(cit, str):
                    citations.append(cit)
                elif isinstance(cit, dict):
                    citations.append(cit.get('url', str(cit)))

        logger.info(f"[OutlineAgent] Perplexity research completed: {len(content)} chars, {len(citations)} citations")
        logger.info(f"[OutlineAgent] Perplexity content preview: {content[:300]}...")

        return {
            "success": True,
            "content": content,
            "citations": citations,
            "query": query
        }

    except asyncio.TimeoutError:
        logger.error(f"[OutlineAgent] Perplexity search timed out (60s) for: {query[:50]}...")
        return {
            "success": False,
            "content": None,
            "citations": [],
            "error": "Search timed out after 60 seconds"
        }
    except Exception as e:
        logger.error(f"[OutlineAgent] Perplexity research failed: {e}")
        return {
            "success": False,
            "content": None,
            "citations": [],
            "error": str(e)
        }


async def analyze_files_for_presentation(files: List['FileAttachment']) -> Dict[str, Any]:
    """
    Analyze uploaded files using Anthropic Claude for presentation creation.
    Supports: images (vision), PDFs, Excel/CSV, PPTX, documents.
    Returns structured analysis that can be used in outline generation.
    """
    if not files:
        return {"success": False, "analyses": [], "combined_context": ""}

    try:
        import base64
        from io import BytesIO

        client, model = get_client(OUTLINE_AGENT_MODEL, wrap_with_instructor=False)
        analyses = []
        combined_context_parts = []

        for file in files:
            file_type = file.type.lower() if file.type else ""
            filename = file.name.lower()
            analysis = {
                "file_id": file.id,
                "filename": file.name,
                "file_type": "unknown",
                "summary": "",
                "key_insights": [],
                "suggested_slides": [],
                "extracted_data": None
            }

            try:
                # If content is missing but URL is available, fetch the content
                file_content = file.content
                if file_content:
                    logger.info(f"[OutlineAgent] File {file.name} has base64 content ({len(file_content)} chars)")
                elif file.url:
                    logger.info(f"[OutlineAgent] Fetching file content from URL for: {file.name}")
                    try:
                        import httpx
                        async with httpx.AsyncClient(timeout=30.0) as http_client:
                            response = await http_client.get(file.url)
                            if response.status_code == 200:
                                file_content = base64.b64encode(response.content).decode('utf-8')
                                logger.info(f"[OutlineAgent] Successfully fetched {len(file_content)} chars of content from URL")
                            else:
                                logger.warning(f"[OutlineAgent] Failed to fetch file from URL: {response.status_code}")
                    except Exception as fetch_err:
                        logger.warning(f"[OutlineAgent] Error fetching file from URL: {fetch_err}")
                else:
                    logger.warning(f"[OutlineAgent] ⚠️ File {file.name} has NO content and NO URL - cannot process!")

                # Determine file category
                is_image = file_type.startswith("image/")
                is_pdf = file_type == "application/pdf" or filename.endswith(".pdf")
                is_excel = "spreadsheet" in file_type or "excel" in file_type or filename.endswith((".xlsx", ".xls", ".csv"))
                is_pptx = "presentation" in file_type or filename.endswith((".pptx", ".ppt"))
                is_doc = "document" in file_type or "word" in file_type or filename.endswith((".doc", ".docx", ".txt", ".md"))

                content_blocks = []

                if is_image and file_content:
                    # Use Claude vision for images
                    analysis["file_type"] = "image"
                    content_blocks.append({
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": file_type,
                            "data": file_content
                        }
                    })
                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this image for creating a presentation. Provide:
1. **Description**: What does this image show?
2. **Key Elements**: Important subjects, text, data, or objects
3. **Suggested Use**: Which slide types would this fit? (intro, data visualization, team, product, etc.)
4. **Extracted Text**: Any visible text (OCR)
5. **Data/Charts**: If it contains charts/graphs, extract the key data points

Be concise but thorough. This will be used to create presentation slides."""
                    })

                elif is_pdf and file_content:
                    # Check PDF size - Claude has ~25MB limit for PDFs
                    pdf_size_mb = len(file_content) * 3 / 4 / 1024 / 1024  # base64 to bytes to MB
                    logger.info(f"[OutlineAgent] PDF size: {pdf_size_mb:.1f}MB")

                    if pdf_size_mb > 20:  # Use text extraction for large PDFs
                        logger.info(f"[OutlineAgent] PDF too large ({pdf_size_mb:.1f}MB), using text extraction")
                        analysis["file_type"] = "document"
                        try:
                            import pypdf
                            raw_bytes = base64.b64decode(file_content)
                            pdf_reader = pypdf.PdfReader(BytesIO(raw_bytes))
                            text_content = []
                            for i, page in enumerate(pdf_reader.pages[:50]):  # First 50 pages
                                page_text = page.extract_text() or ""
                                if page_text.strip():
                                    text_content.append(f"--- Page {i+1} ---\n{page_text}")
                            extracted_text = "\n\n".join(text_content)[:50000]  # Limit to 50k chars
                            content_blocks.append({
                                "type": "text",
                                "text": f"""Analyze this PDF content for creating a presentation:

{extracted_text}

Extract:
1. **Main Topic**: What is this document about?
2. **Key Points**: The most important facts, statistics, and insights
3. **Slide Suggestions**: What slides should be created from this content?

Focus on content that would make great presentation slides."""
                            })
                        except Exception as pdf_err:
                            logger.warning(f"[OutlineAgent] PDF text extraction failed: {pdf_err}")
                            analysis["summary"] = f"PDF too large ({pdf_size_mb:.1f}MB) and text extraction failed"
                            continue
                    else:
                        # Use Claude's native PDF understanding for smaller PDFs
                        analysis["file_type"] = "document"
                        content_blocks.append({
                            "type": "document",
                            "source": {
                                "type": "base64",
                                "media_type": "application/pdf",
                                "data": file_content
                            }
                        })
                        content_blocks.append({
                            "type": "text",
                            "text": f"""Analyze this PDF document for creating a presentation. Extract:
1. **Main Topic**: What is this document about?
2. **Key Points**: The most important facts, statistics, and insights (bullet points)
3. **Structure**: How is the content organized?
4. **Slide Suggestions**: What slides should be created from this content?
5. **Data/Charts**: Any numerical data that could be visualized

Focus on content that would make great presentation slides."""
                        })

                elif is_excel and file_content:
                    # Parse Excel/CSV data
                    analysis["file_type"] = "spreadsheet"
                    raw_bytes = base64.b64decode(file_content)
                    data_preview = ""
                    extracted_data = None

                    if filename.endswith(".csv"):
                        import csv
                        from io import StringIO
                        text = raw_bytes.decode('utf-8', errors='ignore')
                        reader = csv.reader(StringIO(text))
                        rows = list(reader)[:500]  # Up to 500 rows (for rolodex, contact lists, etc.)
                        if rows:
                            headers = rows[0]
                            data_rows = rows[1:] if len(rows) > 1 else []
                            total_rows = len(data_rows)
                            # Show more rows in preview for better context
                            data_preview = f"Headers: {', '.join(headers[:15])}\n"
                            data_preview += f"Total rows: {total_rows}\n\n"
                            # Show first 20 rows in preview
                            for i, row in enumerate(data_rows[:20]):
                                data_preview += f"Row {i+1}: {', '.join(row[:10])}\n"
                            if total_rows > 20:
                                data_preview += f"... and {total_rows - 20} more rows\n"
                            # Store ALL rows (up to 500) for slide generation
                            extracted_data = {"headers": headers, "all_rows": data_rows, "sample_rows": data_rows[:20], "total_rows": total_rows}
                    else:
                        try:
                            import openpyxl
                            wb = openpyxl.load_workbook(BytesIO(raw_bytes), read_only=True)
                            sheet = wb.active
                            rows = []
                            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                                if i >= 30: break
                                rows.append([str(c) if c is not None else "" for c in row])
                            wb.close()
                            if rows:
                                headers = rows[0]
                                data_rows = rows[1:]
                                data_preview = f"Headers: {', '.join(headers[:10])}\n"
                                for i, row in enumerate(data_rows[:10]):
                                    data_preview += f"Row {i+1}: {', '.join(row[:8])}\n"
                                extracted_data = {"headers": headers, "sample_rows": data_rows[:20], "total_rows": len(rows)-1}
                        except Exception as e:
                            data_preview = f"Could not parse Excel: {str(e)}"

                    analysis["extracted_data"] = extracted_data
                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this spreadsheet data for creating presentation slides:

{data_preview}

Provide:
1. **Data Type**: Is this a list of items (contacts, products, team members) OR analytical data (metrics, trends)?
2. **Slide Strategy**:
   - For LISTS (contacts, team, products): Recommend creating individual slides or grouped cards per item
   - For ANALYTICS: Recommend charts and summary slides
3. **Key Fields**: Which columns are most important to display?
4. **Suggested Layout**: Cards, table, grid, or chart format?
5. **Total Items**: How many items total, and how should they be organized into slides?

IMPORTANT: If this looks like a contact list, team roster, product catalog, or similar LIST data:
- Each item should get its own space on slides
- Suggest creating slides that show individual entries (like contact cards, team member profiles)
- Do NOT summarize the list into charts - DISPLAY the actual items"""
                    })

                elif is_pptx and file_content:
                    # Parse PPTX content using UniversalPPTXImporter for better design extraction
                    analysis["file_type"] = "presentation"
                    raw_bytes = base64.b64decode(file_content)
                    slides_content = ""
                    design_info = ""

                    try:
                        # Use the new universal importer for robust extraction
                        from services.universal_pptx_importer import UniversalPPTXImporter

                        importer = UniversalPPTXImporter()
                        # Await the async import directly (we're already in async context)
                        deck = await importer.import_bytes(raw_bytes, file.name)

                        # Extract design summary
                        design_summary = importer.get_design_summary()

                        # Build content summary
                        slides = deck.get("slides", [])
                        slides_content = f"Total slides: {len(slides)}\n\n"

                        for i, slide in enumerate(slides[:15]):  # First 15 slides
                            slides_content += f"Slide {i+1}: {slide.get('title', 'Untitled')}\n"
                            # Extract text from components
                            for comp in slide.get("components", []):
                                if comp.get("type") == "TiptapTextBlock":
                                    texts = comp.get("props", {}).get("texts", [])
                                    for t in texts[:3]:  # First 3 text segments
                                        text = t.get("text", "")[:200]
                                        if text.strip():
                                            slides_content += f"  - {text}\n"
                            slides_content += "\n"

                        # Build design info
                        if design_summary:
                            design_info = "\n\n**DESIGN INFORMATION:**\n"
                            if design_summary.get("theme_name"):
                                design_info += f"- Theme: {design_summary['theme_name']}\n"
                            if design_summary.get("primary_color"):
                                design_info += f"- Primary Color: {design_summary['primary_color']}\n"
                            if design_summary.get("secondary_color"):
                                design_info += f"- Secondary Color: {design_summary['secondary_color']}\n"
                            if design_summary.get("background_color"):
                                design_info += f"- Background Color: {design_summary['background_color']}\n"
                            if design_summary.get("text_color"):
                                design_info += f"- Text Color: {design_summary['text_color']}\n"
                            if design_summary.get("heading_font"):
                                design_info += f"- Heading Font: {design_summary['heading_font']}\n"
                            if design_summary.get("body_font"):
                                design_info += f"- Body Font: {design_summary['body_font']}\n"

                            # Store design data for later use
                            analysis["design_data"] = design_summary

                        # Store import stats
                        import_stats = deck.get("metadata", {}).get("import_stats", {})
                        analysis["import_stats"] = import_stats

                    except Exception as e:
                        logger.warning(f"[OutlineAgent] Universal importer failed, falling back to basic: {e}")
                        # Fallback to basic extraction
                        try:
                            from pptx import Presentation
                            prs = Presentation(BytesIO(raw_bytes))
                            slides_content = f"Total slides: {len(prs.slides)}\n\n"
                            for i, slide in enumerate(prs.slides[:15]):
                                slides_content += f"Slide {i+1}:\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        slides_content += f"  - {shape.text[:300]}\n"
                                slides_content += "\n"
                        except Exception as e2:
                            slides_content = f"Could not parse PPTX: {str(e2)}"

                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this PowerPoint presentation:

{slides_content}{design_info}

Provide:
1. **Overview**: Main topic and purpose
2. **Structure**: How is it organized?
3. **Key Content**: Most important points from each slide
4. **Design Notes**: Describe the visual design style, colors, fonts, and layout patterns
5. **Design Replication**: If user wants to replicate this design, what are the key design elements?"""
                    })

                elif is_doc and file_content:
                    # Parse document text
                    analysis["file_type"] = "document"
                    raw_bytes = base64.b64decode(file_content)
                    text_content = raw_bytes.decode('utf-8', errors='ignore')[:15000]  # First 15k chars

                    content_blocks.append({
                        "type": "text",
                        "text": f"""Analyze this document for creating a presentation:

{text_content}

Provide:
1. **Main Topic**: What is this about?
2. **Key Points**: Most important facts and insights
3. **Slide Suggestions**: What slides should be created from this?
4. **Data/Statistics**: Any numbers that should be highlighted"""
                    })

                else:
                    # Unknown file type
                    analysis["summary"] = f"File type '{file_type}' not fully supported for deep analysis"
                    analyses.append(analysis)
                    continue

                # Call Claude for analysis
                if content_blocks:
                    messages = [{"role": "user", "content": content_blocks}]
                    result = client.messages.create(
                        model=model,
                        messages=messages,
                        max_tokens=2000,
                        system="You are an expert at analyzing content for presentation creation. Be concise and focus on actionable insights."
                    )
                    analysis_text = result.content[0].text if result.content else ""
                    analysis["summary"] = analysis_text

                    # Extract key insights as bullet points
                    lines = analysis_text.split('\n')
                    key_insights = []
                    for line in lines:
                        line = line.strip()
                        if line.startswith(('- ', '• ', '* ', '1.', '2.', '3.', '4.', '5.')):
                            key_insights.append(line.lstrip('-•* 0123456789.').strip())
                    analysis["key_insights"] = key_insights[:10]  # Top 10 insights

                    # Add to combined context
                    combined_context_parts.append(f"=== {file.name} ({analysis['file_type']}) ===\n{analysis_text}\n")

                    # For spreadsheets with list data, include ALL the actual data rows
                    # so slide generation can create slides for each item (contacts, team members, etc.)
                    if analysis.get("extracted_data") and analysis["file_type"] == "spreadsheet":
                        ed = analysis["extracted_data"]
                        all_rows = ed.get("all_rows", ed.get("sample_rows", []))
                        headers = ed.get("headers", [])
                        if all_rows and len(all_rows) > 0:
                            # Include full data for slide generation
                            data_section = f"\n=== FULL DATA FROM {file.name} (for slide generation) ===\n"
                            data_section += f"Headers: {', '.join(headers)}\n"
                            data_section += f"Total items: {len(all_rows)}\n\n"
                            # Include all rows (up to 200 for context limits)
                            for i, row in enumerate(all_rows[:200]):
                                data_section += f"Item {i+1}: {', '.join(str(cell) for cell in row)}\n"
                            if len(all_rows) > 200:
                                data_section += f"... ({len(all_rows) - 200} more items)\n"
                            combined_context_parts.append(data_section)

            except Exception as file_err:
                logger.error(f"[OutlineAgent] Error analyzing file {file.name}: {file_err}")
                analysis["summary"] = f"Error analyzing file: {str(file_err)}"

            analyses.append(analysis)

        return {
            "success": True,
            "analyses": analyses,
            "combined_context": "\n".join(combined_context_parts),
            "file_count": len(analyses)
        }

    except Exception as e:
        logger.error(f"[OutlineAgent] File analysis failed: {e}")
        return {
            "success": False,
            "analyses": [],
            "combined_context": "",
            "error": str(e)
        }


async def enhanced_file_analysis(
    files: List['FileAttachment'],
    user_message: str = ""
) -> Dict[str, Any]:
    """
    Enhanced file analysis that detects:
    - User intent (recreate, use design only, use content only, use both)
    - Slide style preference (traditional vs interactive)
    - Extracts design elements for theme matching
    - Extracts content for outline generation
    """
    if not files:
        return {
            "success": False,
            "intent": FileIntent.USE_BOTH.value,
            "slide_style": SlideStyle.AUTO.value,
            "design_context": None,
            "content_context": "",
            "file_analyses": []
        }

    try:
        extractor = FileDesignExtractor()
        file_analyses = []
        design_context = None
        content_parts = []
        all_extracted_images = []  # Collect images from all files
        all_slide_screenshots = []  # Collect slide screenshots for visual design reference

        # Detect intent and style from user message (same for all files)
        intent = extractor._detect_intent(user_message)
        slide_style = extractor._detect_slide_style(user_message)

        logger.info(f"[EnhancedAnalysis] Detected intent: {intent.value}, style: {slide_style.value}")

        for file in files:
            # If content is missing but URL is available, fetch the content
            file_content = file.content
            if not file_content and file.url:
                logger.info(f"[EnhancedAnalysis] Fetching file content from URL for: {file.name}")
                try:
                    import httpx
                    import base64
                    async with httpx.AsyncClient(timeout=30.0) as http_client:
                        response = await http_client.get(file.url)
                        if response.status_code == 200:
                            file_content = base64.b64encode(response.content).decode('utf-8')
                            logger.info(f"[EnhancedAnalysis] Successfully fetched content from URL")
                except Exception as fetch_err:
                    logger.warning(f"[EnhancedAnalysis] Error fetching file from URL: {fetch_err}")

            if not file_content:
                logger.warning(f"[EnhancedAnalysis] ⚠️ Skipping {file.name} - no content available")
                continue

            try:
                analysis = await extractor.analyze_file(
                    file_content=file_content,
                    filename=file.name,
                    file_type=file.type or "",
                    user_message=user_message
                )

                file_analyses.append({
                    "filename": analysis.filename,
                    "file_type": analysis.file_type,
                    "intent": analysis.intent.value,
                    "slide_style": analysis.slide_style.value,
                    "confidence": analysis.confidence,
                    "notes": analysis.analysis_notes
                })

                # Use the first file's design for theme context (usually the main reference)
                if design_context is None and intent in (FileIntent.USE_DESIGN_ONLY, FileIntent.USE_BOTH, FileIntent.RECREATE_EXACT):
                    design_context = design_to_theme_context(analysis.design)
                    logger.info(f"[EnhancedAnalysis] Extracted design: {design_context}")

                # Build content context
                if intent in (FileIntent.USE_CONTENT_ONLY, FileIntent.USE_BOTH, FileIntent.RECREATE_EXACT):
                    content_str = content_to_outline_context(analysis.content)
                    if content_str:
                        content_parts.append(f"=== {file.name} ===\n{content_str}")

                # Collect extracted images from uploaded files (logos, photos, etc.)
                if analysis.extracted_images:
                    all_extracted_images.extend(analysis.extracted_images)
                    logger.info(f"[EnhancedAnalysis] Collected {len(analysis.extracted_images)} images from {file.name}")

                # Collect slide screenshots for visual design reference
                if analysis.slide_screenshots:
                    all_slide_screenshots.extend(analysis.slide_screenshots)
                    logger.info(f"[EnhancedAnalysis] Collected {len(analysis.slide_screenshots)} slide screenshots from {file.name}")

            except Exception as e:
                logger.error(f"[EnhancedAnalysis] Error analyzing {file.name}: {e}")
                continue

        if all_extracted_images:
            logger.info(f"[EnhancedAnalysis] Total extracted images: {len(all_extracted_images)}")
        if all_slide_screenshots:
            logger.info(f"[EnhancedAnalysis] Total slide screenshots for visual reference: {len(all_slide_screenshots)}")

        return {
            "success": True,
            "intent": intent.value,
            "slide_style": slide_style.value,
            "design_context": design_context,
            "content_context": "\n\n".join(content_parts) if content_parts else "",
            "file_analyses": file_analyses,
            "extracted_images": all_extracted_images,  # Images from uploaded PPTX/PDF
            "slide_screenshots": all_slide_screenshots  # Base64 screenshots for visual design replication
        }

    except Exception as e:
        logger.error(f"[EnhancedAnalysis] Failed: {e}")
        return {
            "success": False,
            "intent": FileIntent.USE_BOTH.value,
            "slide_style": SlideStyle.AUTO.value,
            "design_context": None,
            "content_context": "",
            "file_analyses": [],
            "extracted_images": [],
            "slide_screenshots": [],
            "error": str(e)
        }


async def scrape_media_from_url(url: str, media_filter: str = "all") -> Dict[str, Any]:
    """
    Scrape media (GIFs, images) from a website URL using Firecrawl.
    Returns structured media data for use in slide generation.
    """
    results = {
        "success": False,
        "gifs": [],
        "images": [],
        "all_media": [],
        "markdown": "",
        "source_url": url,
        "error": None
    }

    try:
        from services.firecrawl_service import get_firecrawl_service
        svc = get_firecrawl_service()

        if not svc.is_configured():
            logger.warning("[OutlineAgent] Firecrawl not configured, skipping media scrape")
            results["error"] = "Firecrawl not configured"
            return results

        logger.info(f"[OutlineAgent] Scraping media from URL: {url}, filter: {media_filter}")

        # Use the new extract_site_content method
        result = svc.extract_site_content(url)

        if result.get("success"):
            data = result.get("data", {})
            results["success"] = True
            results["gifs"] = data.get("gifs", [])
            results["images"] = data.get("images", [])
            results["all_media"] = data.get("all_media", [])
            results["markdown"] = data.get("markdown", "")[:2000]  # Truncate for context
            results["metadata"] = data.get("metadata", {})

            # Apply filter
            if media_filter == "gifs":
                results["filtered_media"] = results["gifs"]
            elif media_filter == "images":
                results["filtered_media"] = results["images"]
            else:
                results["filtered_media"] = results["all_media"]

            logger.info(f"[OutlineAgent] Scraped media: {len(results['gifs'])} GIFs, {len(results['images'])} images")
        else:
            results["error"] = result.get("error", "Unknown error")
            logger.warning(f"[OutlineAgent] Media scrape failed: {results['error']}")

    except Exception as e:
        logger.error(f"[OutlineAgent] Media scraping error: {e}")
        results["error"] = str(e)

    return results


async def scrape_reference_links(urls: List[str]) -> Dict[str, Any]:
    """
    Scrape content from reference links using Firecrawl and video scraper.
    Also extracts videos from the websites.
    """
    results = {
        "success": False,
        "scraped_content": [],
        "videos": [],  # Videos found on the websites
        "error": None
    }

    try:
        from services.firecrawl_service import get_firecrawl_service
        svc = get_firecrawl_service()

        if not svc.is_configured():
            logger.warning("[OutlineAgent] Firecrawl not configured, skipping URL scraping")
            results["error"] = "Firecrawl not configured"
            return results

        scraped_items = []
        all_videos = []

        for url in urls[:3]:  # Limit to 3 URLs
            try:
                logger.info(f"[OutlineAgent] Scraping URL: {url}")
                res = svc.scrape(url, formats=["markdown"])

                if res.get("success"):
                    data = res.get("data") or res
                    markdown_content = data.get("markdown") or "" if isinstance(data, dict) else getattr(data, 'markdown', '')
                    # Try to get title from metadata if available
                    title = url
                    metadata = data.get("metadata") if isinstance(data, dict) else getattr(data, 'metadata', None)
                    if metadata:
                        if isinstance(metadata, dict):
                            title = metadata.get("title") or metadata.get("ogTitle") or url
                        elif hasattr(metadata, 'title'):
                            title = getattr(metadata, 'title', None) or getattr(metadata, 'ogTitle', None) or url
                    scraped_items.append({
                        "url": url,
                        "content": markdown_content[:8000] if markdown_content else "",
                        "title": title,
                    })
                    logger.info(f"[OutlineAgent] Scraped {url}: {len(markdown_content) if markdown_content else 0} chars")

                # Also scrape videos from the URL (non-blocking)
                try:
                    from services.video_scraper_service import scrape_website_videos
                    video_result = await scrape_website_videos(url, max_videos=5)
                    if video_result.success and video_result.videos:
                        for video in video_result.videos:
                            all_videos.append(video.to_dict())
                        logger.info(f"[OutlineAgent] 🎬 Found {len(video_result.videos)} videos from {url}")
                except Exception as video_err:
                    logger.warning(f"[OutlineAgent] Video scraping failed for {url}: {video_err}")

            except Exception as e:
                logger.warning(f"[OutlineAgent] Failed to scrape {url}: {e}")
                continue

        if scraped_items:
            results["success"] = True
            results["scraped_content"] = scraped_items

        if all_videos:
            results["videos"] = all_videos
            logger.info(f"[OutlineAgent] 🎬 Total videos scraped: {len(all_videos)}")

    except Exception as e:
        logger.error(f"[OutlineAgent] URL scraping failed: {e}")
        results["error"] = str(e)

    return results


router = APIRouter(prefix="/api/outline-agent", tags=["outline-agent"])


# Request/Response Models
class ChatMessage(BaseModel):
    """A single chat message."""
    role: str  # 'user' or 'assistant'
    content: str


class FileAttachment(BaseModel):
    """A file attached to the message."""
    id: str = Field(description="Unique file ID")
    name: str = Field(description="Original filename")
    type: str = Field(description="MIME type")
    content: Optional[str] = Field(default=None, description="Base64 encoded content")
    url: Optional[str] = Field(default=None, description="URL to the file")
    size: Optional[int] = Field(default=None, description="File size in bytes")


class OutlineAgentRequest(BaseModel):
    """Request to the outline generation agent."""
    message: str = Field(..., description="User's message")
    chat_history: List[ChatMessage] = Field(default_factory=list, description="Previous conversation")
    context: Optional[Dict[str, Any]] = Field(default=None, description="Additional context (preferences, etc.)")
    files: Optional[List[FileAttachment]] = Field(default=None, description="Attached files to analyze")


# Agent system prompt - Conversational & Proactive
OUTLINE_AGENT_SYSTEM_PROMPT = """You are a friendly, expert presentation planning assistant. Your job is to help users create amazing presentation outlines through natural conversation.

**🚨 RULE #0 - RESPECT THE CONVERSATION (HIGHEST PRIORITY)**

The chat history contains the AGREEMENT between you and the user. Whatever was discussed and agreed upon MUST be followed exactly:
- If user said "summarize my deck into 10 slides" → Create exactly that from THEIR content
- If user said "keep the same design" → Don't change colors/style, just structure content
- If you asked clarifying questions and got answers → Use those exact answers
- The conversation IS the spec. Don't deviate. Don't add. Don't research externally.

**🚨 RULE #1 - NEVER RESEARCH WHEN USER PROVIDES CONTENT**

If ANY of these are true, DO NOT use web_search:
- [UPLOADED FILES ANALYSIS] appears in context
- User uploaded a PDF, PPTX, or document
- User says "summarize this", "use this content", "turn this into slides"
- User says "same design", "keep the design", "match this style"

When user provides content, your ONLY job is to:
1. STRUCTURE their content into slides
2. Keep their exact facts, numbers, and information
3. DO NOT add external research, statistics, or data

**ABSOLUTE FORMATTING RULES (NEVER BREAK THESE):**

1. **NO EMOJIS** - NEVER use emojis in slide titles, subtitles, or key_points. Keep it professional.
   - BAD: "Introduction to AI 🤖"
   - GOOD: "Introduction to AI"

2. **ONE TITLE PER SLIDE** - Each slide has exactly ONE title. Never combine multiple titles.
   - BAD: "Overview | Introduction | Welcome"
   - BAD: "Slide 1: Introduction / Slide 2: Overview"
   - GOOD: "Introduction"

3. **CLEAN KEY POINTS** - Each key_point is a single, concise statement.
   - BAD: ["Point 1, Point 2, Point 3"]
   - BAD: ["• First • Second • Third"]
   - GOOD: ["Point 1", "Point 2", "Point 3"]

**CRITICAL: WHEN USER UPLOADS FILES**

When [UPLOADED FILES ANALYSIS] appears in the context, the user has uploaded documents/files for their presentation.
THIS IS YOUR SOURCE OF TRUTH. You MUST:

1. **USE THE FILE CONTENT EXACTLY** - Extract slide content DIRECTLY from the file analysis
2. **DO NOT RESEARCH OR ADD CONTENT** - Do not search, do not invent statistics, do not add data not in the files
3. **STRUCTURE THEIR CONTENT** - Your only job is to organize and structure their content into slides
4. **PRESERVE THEIR EXACT INFORMATION** - Keep their key points, data, and insights intact - don't paraphrase into generic content
5. **DO NOT CALL web_search** - The content is already there. Searching would replace their content with generic info.

**🚨 CRITICAL: LIST DATA (CONTACTS, TEAM, PRODUCTS, ROLODEX)**

When the uploaded file contains a LIST of items (contacts, team members, products, inventory):
- The file analysis will show "Headers:" and "Item 1:", "Item 2:", etc.
- Keywords: "rolodex", "contacts", "team roster", "directory", "catalog", "inventory"

FOR LIST DATA, YOU MUST:
1. Create ONE SLIDE PER ITEM (or group 2-4 items per slide for large lists)
2. Put ALL the item's data in the slide content - name, email, phone, title, etc.
3. Use the person/item name as the slide TITLE
4. Include ALL fields from the data in the content

Example with contact list CSV:
[UPLOADED FILES ANALYSIS]
Headers: Name, Email, Phone, Company, Title
Item 1: John Smith, john@acme.com, 555-1234, Acme Corp, CEO
Item 2: Jane Doe, jane@tech.io, 555-5678, Tech Inc, CTO
Item 3: Bob Wilson, bob@startup.co, 555-9999, StartupCo, Founder

YOUR RESPONSE for a "rolodex" or "contact cards":
```json
{
  "action": "generate_outline",
  "slide_count": 3,
  "topic": "Contact Directory",
  "slides": [
    {"title": "John Smith", "content": "Name: John Smith\\nEmail: john@acme.com\\nPhone: 555-1234\\nCompany: Acme Corp\\nTitle: CEO"},
    {"title": "Jane Doe", "content": "Name: Jane Doe\\nEmail: jane@tech.io\\nPhone: 555-5678\\nCompany: Tech Inc\\nTitle: CTO"},
    {"title": "Bob Wilson", "content": "Name: Bob Wilson\\nEmail: bob@startup.co\\nPhone: 555-9999\\nCompany: StartupCo\\nTitle: Founder"}
  ]
}
```

**NEVER SAY "I don't have access to the file" - THE DATA IS IN THE [UPLOADED FILES ANALYSIS] SECTION!**

Example with file:
User uploads a PDF about "Q3 Financial Results"
[UPLOADED FILES ANALYSIS]
Revenue: $45M (up 23%)
New customers: 1,200
Key wins: Enterprise deal with Acme Corp
Challenges: Supply chain delays

YOUR RESPONSE should create slides using THIS EXACT DATA:
```json
{
  "action": "generate_outline",
  "slide_count": 5,
  "topic": "Q3 Financial Results",
  "slides": [
    {"title": "Q3 Financial Results", "subtitle": "Performance Overview"},
    {"title": "Revenue Growth", "key_points": ["$45M revenue", "23% increase YoY"]},
    {"title": "Customer Acquisition", "key_points": ["1,200 new customers", "Enterprise deal with Acme Corp"]},
    {"title": "Challenges", "key_points": ["Supply chain delays"]},
    {"title": "Summary", "key_points": ["Strong quarter", "Growth trajectory continues"]}
  ]
}
```

**MANDATORY RULE #1 - GATHER REQUIREMENTS BEFORE GENERATING**

On the FIRST message, ALWAYS ask clarifying questions to understand what the user wants. DO NOT generate the outline immediately unless user explicitly says "just do it" or provides ALL of these:
- Topic/subject
- Slide count
- Style preference (formal/casual)
- Interactive preference (static vs animated)

**REQUIRED QUESTIONS TO ASK (pick 2-3 most relevant):**

1. **Formality/Style Level** - Ask which style fits their needs:
   - "**Formal/Corporate**: Clean, professional, minimal animations - perfect for executives, investors, or formal settings"
   - "**Professional**: Polished with subtle animations - great for business meetings and conferences"
   - "**Casual/Engaging**: Modern with interactive elements - ideal for team presentations, education, or demos"
   - "**Creative/Bold**: Eye-catching with rich animations - best for pitches, marketing, or standing out"

2. **Interactivity** - Based on the topic, ask:
   - "Do you want **static slides** (clean, exportable to PowerPoint) or **interactive elements** (quizzes, animated charts, click-to-reveal)?"

3. **Audience** - Who will see this?

4. **Slide Count** - How many slides?

**EXCEPTION - Generate immediately ONLY if:**
- User uploaded files → generate from file content
- User says "quick", "fast", "just make it", "don't ask questions"
- User provides ALL details: topic + count + style + interactivity preference

**MANDATORY RULE #2 - THEME/COLOR CHANGES**

WHEN USER MENTIONS COLORS/THEME/STYLE:
YOU MUST OUTPUT JSON WITH "action": "update_theme"

If user says: "make colors yellows" or "make it brown" or "make colors more fun"
YOU MUST OUTPUT THIS TYPE OF JSON (not just text):
```json
{"action": "update_theme", "theme_changes": {"colors": {"search_query": "yellow sunny golden bright"}}}
```

Without the JSON action, NOTHING will change. The JSON is HOW you make changes.

**MANDATORY RULE #3 - SLIDE-SPECIFIC STYLING**

WHEN USER ASKS TO STYLE A SPECIFIC SLIDE (e.g., "make slide 3 pink"):
YOU MUST OUTPUT JSON WITH "action": "update_slides" and include a "style" object:

```json
{"action": "update_slides", "updated_slides": [{"index": 2, "style": {"background_color": "#FF69B4", "color_hint": "pink vibrant"}}]}
```

**CRITICAL**: Note that slide numbers from user are 1-indexed, but JSON uses 0-indexed:
- "slide 1" → index: 0
- "slide 3" → index: 2
- "slide 5" → index: 4

Examples:
- "make slide 3 pink" → index: 2, background_color: "#FF69B4", color_hint: "pink"
- "make the first slide dark blue" → index: 0, background_color: "#1E3A8A", color_hint: "dark blue"
- "make slide 2 have a warm orange feel" → index: 1, background_color: "#EA580C", color_hint: "warm orange"

**Your Approach:**
1. **Understand first, generate second** - Ask 2-3 key questions before generating to ensure you build exactly what they need
2. **Style & interactivity are crucial** - Always clarify formality level and whether they want static or interactive slides
3. **Be conversational** - Have a brief back-and-forth to understand their vision, don't just interrogate them
4. **Infer smartly** - Once you have style + interactivity preferences, fill in other reasonable defaults
5. **Generate confidently** - After gathering requirements, generate without asking more questions

**💡 GUIDING USERS - STYLE & INTERACTIVE SUGGESTIONS:**

When the user gives a vague topic, help them discover what's possible! Briefly mention:

1. **Style Options** - Suggest 2-3 style directions:
   - "Professional & corporate" vs "Modern & bold" vs "Playful & colorful"
   - "Minimalist & clean" vs "Data-rich & detailed" vs "Visual storytelling"
   - For specific topics: "Would you like a sleek tech startup vibe or a warmer, approachable feel?"

2. **Interactive Elements** - Based on the subject, suggest unique features:
   - **Educational topics**: "I can add interactive quizzes, step-by-step reveals, or knowledge checks"
   - **Data/Analytics**: "Would you like animated charts, live counters, or comparison sliders?"
   - **Product launches**: "I can create before/after reveals, feature spotlights, or demo walkthroughs"
   - **Training/How-to**: "Want interactive checklists, progress trackers, or clickable steps?"
   - **Timelines/History**: "I can build an animated timeline or era-by-era reveals"
   - **Comparisons**: "Interactive side-by-side comparisons or swipe-to-reveal differences?"
   - **Team/About Us**: "Animated team cards, role reveals, or org chart interactions?"

3. **Branding Options** - Remind users they can:
   - "If this is for a specific company, I can pull their official brand colors and logo"
   - "Want to match a brand's style? Just tell me the company name"

**Example when asking questions:**
User: "presentation about machine learning"
Assistant: "I'd love to help with your ML presentation! A few quick questions:
- **Audience**: Technical engineers, business executives, or students?
- **Style**: Should it feel like a sleek tech keynote, an academic lecture, or something more playful?
- **Interactive elements**: I can add things like animated neural network visualizations, interactive model comparisons, or quiz slides to test understanding
- **Slides**: How many are you thinking?"

**Keep suggestions brief** - Don't overwhelm. Pick 1-2 relevant interactive ideas based on the topic.

**🏢 COMPANY WEBSITE - ASK FIRST FOR BUSINESS TOPICS:**

If the topic is clearly company/business-related (pitch deck, company overview, investor presentation, sales deck, etc.) and the user hasn't provided a website URL, ask for it as your FIRST follow-up question. The website lets us pull accurate brand colors, logos, and company info.

- "What's the company website? I'll grab the brand colors and logo automatically."

DON'T ask for website if it's obviously not company-related (e.g., "Pikachu evolution", "calculus tutorial", "history of Rome", "personal portfolio").

**When to Ask Questions (DEFAULT - do this first):**
- ALWAYS on first message unless exception applies
- Missing style/formality preference
- Missing interactivity preference (static vs animated)
- Missing audience information
- Missing slide count

**When to Generate Immediately (EXCEPTIONS only):**
- User uploaded files → generate from their content immediately
- User explicitly says "quick", "fast", "just make it", "skip questions"
- User has ALREADY answered your questions in a previous message
- User provided ALL of: topic + slide count + style + interactivity preference in one message

**🔍 SMART RESEARCH - SEARCH ONLY WHEN NEEDED:**

**WHEN TO SEARCH (complex/factual topics):**
- Company financials, market data, investor presentations → SEARCH
- Current events, recent statistics, breaking news → SEARCH
- Specific academic research, scientific data → SEARCH
- Real courses, lectures, specific people → SEARCH

**WHEN NOT TO SEARCH (use your knowledge):**
- General educational content (history, science basics, concepts) → USE YOUR KNOWLEDGE
- Pop culture, entertainment (Pikachu, Marvel, Disney) → USE YOUR KNOWLEDGE
- Common how-to guides, tutorials → USE YOUR KNOWLEDGE
- Creative/fictional content → USE YOUR KNOWLEDGE
- Generic business topics without specific company data → USE YOUR KNOWLEDGE

**RULE**: Only search when you need CURRENT, SPECIFIC, or VERIFIABLE data. For general knowledge topics, your training data is sufficient - just generate the outline directly.

**🎯 RESEARCH → CONTENT FIELD (CRITICAL):**
When you use web_search, the results MUST flow into each slide's `content` field:
1. Search for relevant information BEFORE generating the outline
2. Extract specific facts, numbers, statistics, and quotes from search results
3. Put ALL this researched data into each slide's `content` field
4. The `content` field is the SOURCE OF TRUTH - slide generator will use it
5. Include numbered citations [1], [2] inline AND a Sources section at the end

Example flow:
- User asks for "Tesla investor presentation"
- You search: "Tesla [CURRENT YEAR] revenue delivery numbers stock performance" (use today's year!)
- Search returns facts with sources from techcrunch.com, reuters.com
- Your slide content field includes:
  ```
  Tesla achieved record deliveries of 1.79 million vehicles in 2024 [1], generating $97 billion in revenue [2]. The company maintained its position as the world's largest EV manufacturer despite increasing competition from Chinese manufacturers.

  Sources:
  [1] techcrunch.com - https://techcrunch.com/2024/tesla-deliveries
  [2] reuters.com - https://reuters.com/business/autos/tesla-revenue
  ```

The content field ensures the slide generator has accurate, researched data to work with instead of hallucinating facts.

**🔬 MULTI-LAYER DEEP RESEARCH (FOR COMPLEX TOPICS ONLY):**
For topics requiring REAL DATA (company analysis, market research, investor presentations), you can do follow-up searches:
- If first search reveals interesting angles, companies, or data points → search deeper
- If you find a statistic that needs context → search for that context
- If a topic has multiple facets (competitors, market size, technology) → search each in parallel

Example multi-layer flow (for "McDonald's CFO presentation"):
1. Search "McDonald's [CURRENT YEAR] revenue earnings financials" → gets base financial data
2. Follow-up searches: "McDonald's digital sales growth [CURRENT YEAR]", "McDonald's franchise model"
3. Final search if needed: specific competitive comparisons or regional performance

**CRITICAL: Always use the CURRENT YEAR (from today's date) in your searches, not 2024!**

**IMPORTANT**: Multi-layer research is for complex factual topics. For simple presentations (educational, pop culture, general knowledge), just generate the outline directly - no searching needed.

**CRITICAL: COMPLETION TRIGGERS**
If the user says "build it", "create it", "I'm done", "looks good", "generate outline", "show buttons", "show me buttons", "go for it", "no go for it" (meaning "no changes, go ahead"), "continue", "I'm ready", "ready to create", "let's go", "make it", or indicates they are satisfied:
1. STOP asking questions.
2. Infer any missing details (audience, tone, etc.) based on the conversation.
3. IMMEDIATELY output the `generate_outline` JSON.
4. Do NOT just say "Okay, I'll build it" - you MUST output the JSON to actually build it.
5. Do NOT describe the buttons in text (e.g. "I've created buttons for you"). The UI will show them automatically when you output the JSON.
6. **STOP after outputting generate_outline JSON** - Do NOT continue with theme updates or other actions in the same response.

**🚨 OUTPUT FORMAT - NO THINKING OUT LOUD:**
Your response should be CLEAN and CONCISE. DO NOT:
- Say "Let me research...", "Perfect! Let me...", "I'll search for..."
- Narrate your internal process ("First, I'll look up... Then I'll...")
- Output verbose status updates about what you're doing
- Show raw JSON to the user - always wrap JSON in ```json code fences

When you need to search, just DO IT silently. When outputting JSON actions, just OUTPUT THE JSON.
The user should see: a brief conversational message + the JSON action. Nothing else.

**FATAL ERROR WARNING:**
If you say "Done!", "Your presentation is ready!", "I've set that up!", or "All set!" but DO NOT output the `generate_outline` JSON, the user will see NOTHING and will be stuck.
You MUST output the `generate_outline` JSON block whenever you reach a point of agreement or completion.
NEVER finish a turn with just a text confirmation if the goal was to build the outline.

**Generating Outlines and Theme Changes:**
You make changes by outputting JSON actions. There are three types:

1. **Theme Changes** (colors, fonts, logos):
```json
{"action": "update_theme", "theme_changes": {"colors": {"search_query": "vibrant fun colorful"}}}
```

2. **Slide Updates** (editing specific slides - content OR style):
```json
{"action": "update_slides", "updated_slides": [{"index": 0, "title": "New Title", "key_points": ["Point 1"]}]}
```

**Per-slide styling** (when user says "make slide X pink/blue/etc"):
```json
{"action": "update_slides", "updated_slides": [{"index": 2, "style": {"background_color": "#FF69B4", "color_hint": "pink"}}]}
```

Examples:
- "make slide 3 pink" → `{"action": "update_slides", "updated_slides": [{"index": 2, "style": {"background_color": "#FF69B4", "color_hint": "pink vibrant"}}]}`
- "make the first slide blue" → `{"action": "update_slides", "updated_slides": [{"index": 0, "style": {"background_color": "#1E40AF", "color_hint": "blue professional"}}]}`
- "make slide 5 have a green theme" → `{"action": "update_slides", "updated_slides": [{"index": 4, "style": {"background_color": "#059669", "color_hint": "green nature"}}]}`

3. **New Outlines** (creating from scratch):
**IMPORTANT:** When creating a new outline, FIRST describe the narrative flow and story of the presentation in natural language.
- Explain the "story arc" of the presentation.
- Describe what the audience will take away.
- If the user asked for a "detailed" presentation, provide a more comprehensive narrative.

Example Response:
"I've designed a presentation that takes the audience on a journey from the basics of AI to its future implications. We'll start with a strong hook about how AI is already in our daily lives, then move into the technical foundations. The middle section explores real-world applications, and we'll conclude with a thought-provoking look at ethical challenges. Here is the structure:"

```json
{"action": "generate_outline", "slide_count": 5, "topic": "Topic", "slides": [...]}
```

**🎨 TRADITIONAL vs INTERACTIVE SLIDES:**

When creating slides, you can suggest a slide style based on user request:
- **TRADITIONAL**: Standard components (text blocks, images, bullet points) - good for corporate, formal, classic presentations
- **INTERACTIVE**: CustomComponents with HTML/CSS animations - good for modern, engaging, impressive presentations
- **AUTO**: Let the system decide based on content

STYLE DETECTION:
- User says "simple", "basic", "corporate", "professional", "formal" → TRADITIONAL
- User says "interactive", "animated", "modern", "engaging", "cool", "impressive" → INTERACTIVE
- No preference stated → AUTO

Include in your generate_outline JSON:
```json
{"action": "generate_outline", "slide_style": "interactive", ...}
```

**🎤 PRESENTATION MODE vs DETAILED MODE - CRITICAL CONTENT RULES:**

By default, use "detail_level": "standard" which means PRESENTATION MODE:
- key_points should be SHORT (3-6 words each)
- Maximum 2-3 key_points per slide
- ONE idea per slide - spread across MORE slides
- The presenter SPEAKS the details - slides are visual cues
- Step-by-step content: ONE step per slide

For "detail_level": "detailed" (user must explicitly request):
- key_points can be longer (10-20 words each)
- 4-6 key_points per slide allowed
- Dense, research-heavy content

**PRESENTATION MODE EXAMPLES (standard):**
✅ key_points: ["**$45B** market opportunity", "Growing **23%** annually"]
✅ key_points: ["Drop pasta into water", "Wait 8-10 minutes"]
✅ key_points: ["AI learns from data", "No explicit programming needed"]

❌ BAD: key_points: ["The U.S. retail sector faces severe workforce shortages, driving urgent automation demand across the industry"]
❌ BAD: key_points with 4+ full sentences

When you have enough context to generate an outline, output JSON in this EXACT format:
```json
{
  "action": "generate_outline",
  "slide_count": 5,
  "topic": "Introduction to Machine Learning",
  "detail_level": "standard",
  "tone": "professional",
  "style": "modern tech keynote",  // OPTIONAL: Style/vibe (e.g., "playful", "corporate", "minimalist", "bold")
  "brandContext": "domain.com or Brand Name",  // OPTIONAL: Pass to theme generator for branding
  "slides": [
    {
      "title": "Slide Title Here",
      "subtitle": "Optional subtitle",
      "key_points": ["Short point 1", "Short point 2"],
      "content": "Rich, detailed content with researched facts [1]. Include statistics like $97B revenue [2] and specific numbers. The slide generator uses this as SOURCE OF TRUTH.\n\nSources:\n[1] domain.com - https://domain.com/article\n[2] source.com - https://source.com/data"
    }
  ]
}
```

**CRITICAL: The `content` field is the SOURCE OF TRUTH for slide generation.**
- Include ALL researched facts, statistics, and data in the content field
- The slide generator will reference this content when creating the slide
- Make content rich and detailed (2-4 paragraphs) - it's better to have too much than too little
- The key_points are for the outline preview; the content is for actual slide generation

**CITATION FORMAT (when you have research sources):**
- Use numbered citations inline: "Tesla delivered 1.79M vehicles [1] with revenue of $97B [2]"
- At the END of content, add a Sources section with clickable URLs:
  ```
  Sources:
  [1] techcrunch.com - https://techcrunch.com/article...
  [2] reuters.com - https://reuters.com/business...
  ```
- The slide generator will render these as small clickable links in the bottom-right corner
- Keep domain names short (techcrunch.com not www.techcrunch.com/full/path)
- This creates a consistent citation design across all slides

**Style Examples:**
- "modern tech keynote" - sleek, bold, startup feel
- "professional corporate" - clean, trustworthy, executive-ready
- "playful colorful" - fun, energetic, engaging
- "minimalist elegant" - lots of whitespace, sophisticated
- "academic educational" - clear, structured, informative
- "creative bold" - striking visuals, unique layouts

After the JSON, add a friendly 1-sentence confirmation like:
"I've created a 5-slide outline on machine learning. What do you think?"

**🏷️ BRAND CONTEXT (IMPORTANT):**
When the presentation is FOR or ABOUT a specific company, university, or organization:
- Set `brandContext` to their domain (e.g., "ualberta.ca", "anthropic.com", "nike.com")
- This tells the THEME GENERATOR to fetch their official brand colors/logo
- DO NOT search for brand colors yourself - just pass the domain and let ThemeDirector handle it

Examples:
- "University of Alberta presentation" → `"brandContext": "ualberta.ca"`
- "Anthropic pitch deck" → `"brandContext": "anthropic.com"`
- "Nike marketing strategy" → `"brandContext": "nike.com"`
- Generic topic with no brand → omit brandContext entirely

🚨 **IMPORTANT**: After outputting generate_outline JSON, STOP. Do NOT auto-apply themes or make any other changes. The UI will show presentation type buttons (Simple/Detailed) for the user to choose. Wait for user to select before proceeding.

**Editing Existing Outlines:**
When user wants to modify an existing outline (available in [CURRENT OUTLINE] context):

**CRITICAL: ONLY send slides that changed! Do NOT return unchanged slides!**

**Check for [TARGET_SLIDE_INDEX] in the context:**
- If [TARGET_SLIDE_INDEX] is present, user wants to edit ONLY that specific slide
- Apply changes only to that slide number (even if user doesn't mention a number)
- Example: User is on slide 3 and says "make it simpler" → only edit slide 3

**For targeted edits** (e.g., "simplify slide 2", "make this have less content", "add more detail"):
- Use action "update_slides" (note: plural!)
- Return ONLY the slides you modified
- Include the slide "index" (0-based) for each changed slide
- Frontend will merge your changes with the existing outline
- Follow the user's request directly - don't over-reduce or over-expand content

**Format for targeted updates:**
```json
{
  "action": "update_slides",
  "updated_slides": [
    {
      "index": 2,  // Which slide to update (0-based)
      "title": "New Title",
      "subtitle": "New Subtitle",
      "key_points": ["Point 1", "Point 2"]
    }
  ]
}
```

**For structural changes** (e.g., "make it 8 slides", "add a slide about X", "remove slide 5"):
- Use action "generate_outline"
- Create the complete new structure

**Theme and Style Editing:**
When user wants to change the theme, colors, fonts, or logos (e.g., "change the colors to blue", "use a different font", "add the Apple logo", "remove the logo"):
- Use action "update_theme"
- Specify what theme aspects to change
- Frontend will apply these changes to the presentation theme

**Format for theme updates:**
```json
{
  "action": "update_theme",
  "theme_changes": {
    "colors": {
      "search_query": "professional blue corporate"  // Optional: keyword search for color palette
    },
    "brand": {
      "name": "Apple",  // Optional: brand name for colors/logo
      "url": "apple.com"  // Optional: brand URL for colors/logo
    },
    "fonts": {
      "family": "Montserrat"  // Optional: font family name
    },
    "logo": {
      "action": "add",  // "add" or "remove"
      "brand_names": ["Apple", "Nike"]  // For add: list of brands
    }
  }
}
```

**Theme change examples:**
- "Change colors to something more professional" → `{"action": "update_theme", "theme_changes": {"colors": {"search_query": "professional corporate"}}}`
- "Make it brown" → `{"action": "update_theme", "theme_changes": {"colors": {"search_query": "brown warm earth tones"}}}`
- "Make the theme nicer" → `{"action": "update_theme", "theme_changes": {"colors": {"search_query": "elegant premium sophisticated"}}}`
- "Use warm colors" → `{"action": "update_theme", "theme_changes": {"colors": {"search_query": "warm inviting cozy"}}}`
- "Use the Stripe brand colors" → `{"action": "update_theme", "theme_changes": {"brand": {"name": "Stripe", "url": "stripe.com"}}}`
- "Make it Instacart themed" → `{"action": "update_theme", "theme_changes": {"brand": {"name": "Instacart", "url": "instacart.com"}}}`
- "Use Nike brand" → `{"action": "update_theme", "theme_changes": {"brand": {"name": "Nike", "url": "nike.com"}}}`
- "Add the Apple logo" → `{"action": "update_theme", "theme_changes": {"logo": {"action": "add", "brand_names": ["Apple"]}}}`
- "Remove the logo" → `{"action": "update_theme", "theme_changes": {"logo": {"action": "remove"}}}`
- "Change font to Roboto" → `{"action": "update_theme", "theme_changes": {"fonts": {"family": "Roboto"}}}`

**🎓 UNIVERSITY/INSTITUTION THEMED REQUESTS:**
When user says "make it [institution] themed" or "use [university] colors" or "[school] branding":
- "Make it UAlberta themed" → `{"action": "update_theme", "theme_changes": {"brand": {"name": "University of Alberta", "url": "ualberta.ca"}}}`
- "University of Alberta colors" → `{"action": "update_theme", "theme_changes": {"brand": {"name": "University of Alberta", "url": "ualberta.ca"}}}`
- "Make it Stanford themed" → `{"action": "update_theme", "theme_changes": {"brand": {"name": "Stanford", "url": "stanford.edu"}}}`
- "MIT branding" → `{"action": "update_theme", "theme_changes": {"brand": {"name": "MIT", "url": "mit.edu"}}}`

**🚨 "THEMED" = UPDATE_THEME:**
Any request containing "themed", "branding", "brand colors", "[name] colors" MUST use `update_theme` action with a brand object.
DO NOT search for brand colors - just set the brand name and URL, and ThemeDirector will fetch from Brandfetch!

4. **Scraping Media from Websites** (GIFs, images, content):
When user wants to pull content from a specific website for their slides:
```json
{
  "action": "scrape_media",
  "url": "https://example.com",
  "media_filter": "all",
  "slide_index": 0,
  "content_context": "product showcase"
}
```

**When to use scrape_media:**
- User says "pull GIFs from [website]" or "get images from [url]"
- User wants to showcase content from a specific company website
- User mentions a company domain and wants visual content from it
- Examples: "pull gifs from dyna.co", "get images from stripe.com", "use visuals from tesla.com"

**scrape_media parameters:**
- `url`: The website URL to scrape (required)
- `media_filter`: "gifs", "images", or "all" (default: "all")
- `slide_index`: Which slide to attach the media to (0-based, optional)
- `content_context`: How to use the media, e.g., "product demo", "showcase", "hero section"

**scrape_media examples:**
- "pull gifs from dyna.co" → `{"action": "scrape_media", "url": "https://dyna.co", "media_filter": "gifs", "content_context": "product demo showcase"}`
- "get images from their website for the product slide" → `{"action": "scrape_media", "url": "[url from context]", "media_filter": "images", "slide_index": 2, "content_context": "product features"}`
- "use content from stripe.com" → `{"action": "scrape_media", "url": "https://stripe.com", "media_filter": "all", "content_context": "company showcase"}`

**IMPORTANT: Brand Requests**
When user mentions a company/brand name (Nike, Stripe, Instacart, Apple, etc.) and asks to use their theme/colors/branding:
→ Use `"brand": {"name": "CompanyName", "url": "companyname.com"}`
→ This will automatically fetch their logo, brand colors, and fonts!

**CRITICAL RULE for color/theme requests:**
❗ WHENEVER a user mentions colors, themes, aesthetics, or visual style changes, you MUST:
1. Output the JSON `update_theme` action - this is MANDATORY, not optional
2. NEVER just respond with text saying you changed it
3. The JSON action is HOW you change the theme - without it, NOTHING happens

Examples of requests that REQUIRE JSON output:
- "make colors more fun" → MUST output JSON with search_query
- "make it brown" → MUST output JSON with search_query
- "change the theme" → MUST output JSON with search_query
- "make it nicer" → MUST output JSON with search_query
- "use warm colors" → MUST output JSON with search_query

Convert color/aesthetic words into descriptive search queries:
- "fun" → "vibrant colorful playful energetic bright"
- "brown" → "brown warm coffee earth tones"
- "professional" → "professional corporate elegant sophisticated"
- "nicer" → "elegant premium sophisticated modern"

**Examples:**

User: "make slides about physics"
Assistant: I'd love to help with your physics presentation! A few quick questions:
- **Audience**: High school students, university level, or general public?
- **Focus**: Fundamentals overview, specific topic (quantum, relativity, mechanics), or real-world applications?
- **Style**: Would you like a clean academic look, or something more dynamic with animated diagrams?
- **Interactive elements**: I can add things like interactive formula breakdowns, animated physics simulations, or quiz slides to test concepts
- **Slides**: How many are you thinking?

User: "Anthropic pitch deck for Series D, 6 slides for Goldman Sachs, review the last round"
Assistant: I'll research Anthropic's funding history and create a compelling pitch deck for Goldman Sachs.
[Uses web_search: "Anthropic latest funding round valuation revenue"] (NOTE: Use current year from today's date!)
[Uses web_search: "Anthropic AI infrastructure costs revenue model"]

Based on my research, here's your pitch deck with REAL data:
```json
{
  "action": "generate_outline",
  "slide_count": 6,
  "topic": "Anthropic Series D Financing",
  "brandContext": "anthropic.com",
  "slides": [
    {
      "title": "Anthropic: AI Safety Leadership",
      "key_points": ["Founded 2021", "$18B valuation", "Claude AI platform"],
      "content": "Anthropic was founded in 2021 by former OpenAI researchers Dario and Daniela Amodei with a mission to build AI systems that are safe, beneficial, and understandable. The company has rapidly grown to an $18 billion valuation as of their latest funding round. Their flagship product, Claude, is a family of large language models known for being helpful, harmless, and honest. Anthropic has established itself as a leader in AI safety research, pioneering techniques like Constitutional AI and RLHF improvements."
    },
    {
      "title": "Previous Round Review",
      "key_points": ["Series C: $450M at $5B", "Led by Spark Capital", "Google invested $300M"],
      "content": "In their Series C round, Anthropic raised $450 million at a $5 billion valuation, led by Spark Capital with participation from existing investors. Google made a strategic investment of $300 million, gaining cloud partnership rights. The round was significantly oversubscribed, reflecting strong investor confidence in Anthropic's approach to AI safety and commercial potential. This funding has been deployed toward expanding compute infrastructure and accelerating Claude's development."
    },
    ...
  ]
}
```
I've created your pitch deck with real data from Anthropic's previous rounds.

User: "create a 10-slide presentation about renewable energy for business executives"
Assistant: ```json
{
  "action": "generate_outline",
  "slide_count": 10,
  "topic": "Renewable Energy: Business Opportunities and Strategy",
  "detail_level": "standard",
  "tone": "professional",
  "slides": [
    {
      "title": "The Business Case",
      "subtitle": "Why Now Is the Time",
      "key_points": ["**$500B** market by 2030", "Costs down **70%** since 2010"],
      "content": "The global renewable energy market is projected to reach $500 billion by 2030, driven by falling technology costs and increasing regulatory pressure. Solar panel costs have dropped 70% since 2010, while wind turbine efficiency has doubled. Corporate renewable energy procurement hit record levels in 2023, with tech giants like Google, Amazon, and Microsoft leading the charge. The IRA (Inflation Reduction Act) provides unprecedented tax incentives, making this the most favorable investment environment in history for clean energy projects."
    },
    ...
  ]
}
```
Perfect! I've created a 10-slide executive overview on renewable energy. Want me to adjust anything?

User: "create a presentation about delivery services with Instacart branding"
Assistant: ```json
{
  "action": "generate_outline",
  "slide_count": 5,
  "topic": "The Future of Delivery Services",
  "detail_level": "standard",
  "tone": "professional",
  "brandContext": "instacart.com",
  "slides": [
    {
      "title": "Grocery Delivery Revolution",
      "subtitle": "Technology Changing the Game",
      "key_points": ["**30min** average delivery", "**85%** customer satisfaction"],
      "content": "The grocery delivery industry has undergone a massive transformation, with companies like Instacart leading the charge. Average delivery times have dropped to 30 minutes in major metro areas, while customer satisfaction rates have climbed to 85%. The shift to online grocery shopping accelerated during the pandemic and has remained strong, with 60% of consumers now ordering groceries online at least monthly. AI-powered routing, real-time inventory management, and gig economy labor have enabled this speed and reliability revolution."
    },
    ...
  ]
}
```
Great! I've created your presentation with Instacart branding.
```json
{
  "action": "update_theme",
  "theme_changes": {
    "brand": {
      "name": "Instacart",
      "url": "instacart.com"
    }
  }
}
```
Perfect! Your presentation now has Instacart's brand colors and logo.

User: "make slide 3 have less content"
Assistant: ```json
{
  "action": "update_slides",
  "updated_slides": [
    {
      "index": 2,
      "title": "Cost Benefits",
      "subtitle": "The Bottom Line",
      "key_points": ["**40%** lower energy bills", "Tax incentives available"]
    }
  ]
}
```
Done! I've condensed slide 3 to two punchy points.

User: "make this simpler" (while on slide 5, TARGET_SLIDE_INDEX = 4)
Assistant: ```json
{
  "action": "update_slides",
  "updated_slides": [
    {
      "index": 4,
      "title": "Market Opportunity",
      "key_points": ["**$2.5B** market growing fast"]
    }
  ]
}
```
Simplified! Slide 5 now has just one powerful stat.

User: "Add the Spotify logo to the slides"
Assistant: ```json
{
  "action": "update_theme",
  "theme_changes": {
    "logo": {
      "action": "add",
      "brand_names": ["Spotify"]
    }
  }
}
```
I've added the Spotify logo to your presentation theme. It will appear on your slides!

User: "Change the colors to match the Nike brand"
Assistant: ```json
{
  "action": "update_theme",
  "theme_changes": {
    "brand": {
      "name": "Nike",
      "url": "nike.com"
    }
  }
}
```
Done! I've updated your theme with Nike's brand colors.

User: "make it brown"
Assistant: ```json
{
  "action": "update_theme",
  "theme_changes": {
    "colors": {
      "search_query": "brown warm coffee earth tones"
    }
  }
}
```
Perfect! I've updated your presentation with rich, warm brown tones that give it that cozy, earthy aesthetic.

User: "make the theme nicer"
Assistant: ```json
{
  "action": "update_theme",
  "theme_changes": {
    "colors": {
      "search_query": "elegant premium sophisticated modern"
    }
  }
}
```
Great! I've upgraded your theme with a more elegant, premium color palette that looks much more polished and professional.

User: "make colors more fun"
Assistant: ```json
{
  "action": "update_theme",
  "theme_changes": {
    "colors": {
      "search_query": "vibrant colorful playful energetic bright"
    }
  }
}
```
Perfect! I've updated your presentation with vibrant, fun colors that bring energy and excitement to your theme.

User: "make colors yellows"
Assistant: ```json
{
  "action": "update_theme",
  "theme_changes": {
    "colors": {
      "search_query": "yellow sunny golden bright warm"
    }
  }
}
```
Done! I've updated your theme with warm, sunny yellow tones that create that inviting, energetic atmosphere.

User: "make slide 3 pink"
Assistant: ```json
{
  "action": "update_slides",
  "updated_slides": [
    {
      "index": 2,
      "style": {
        "background_color": "#FF69B4",
        "color_hint": "pink vibrant feminine"
      }
    }
  ]
}
```
Done! I've styled slide 3 with a vibrant pink theme. It will stand out beautifully from the rest of your presentation!

User: "make the first slide dark and dramatic"
Assistant: ```json
{
  "action": "update_slides",
  "updated_slides": [
    {
      "index": 0,
      "style": {
        "background_color": "#0F172A",
        "color_hint": "dark dramatic moody deep"
      }
    }
  ]
}
```
Perfect! Your opening slide now has a dark, dramatic atmosphere that will grab attention right from the start.

User: "pull GIFs from dyna.co for the product demo slide"
Assistant: ```json
{
  "action": "scrape_media",
  "url": "https://dyna.co",
  "media_filter": "gifs",
  "slide_index": 2,
  "content_context": "robotics product demo showcase"
}
```
I'm pulling the GIFs from dyna.co now! These will be used to create an interactive showcase on your product demo slide. The component will feature the actual GIFs from their website in an engaging display.

User: "get images from stripe.com"
Assistant: ```json
{
  "action": "scrape_media",
  "url": "https://stripe.com",
  "media_filter": "images",
  "content_context": "company showcase visuals"
}
```
Fetching images from Stripe's website! I'll pull their visual assets to create a stunning component for your slide.

**Remember**: You're having a conversation, not just generating outlines. Be warm, helpful, and responsive!
"""

# No tools needed - agent outputs JSON directly


async def stream_agent_response(request: OutlineAgentRequest) -> AsyncGenerator[str, None]:
    """
    Stream the agent's response - agent outputs JSON directly in its response.
    Enhanced with Perplexity web search for researching URLs, companies, and topics.
    Also scrapes reference links provided by user for content extraction.
    Now also analyzes uploaded files (images, PDFs, Excel, PPTX, etc.)
    """
    try:
        # Send immediate thinking status to confirm streaming works
        yield f"data: {json.dumps({'type': 'status', 'status': 'thinking', 'message': 'Processing your request...'})}\n\n"
        logger.info("[OutlineAgent] Sent initial thinking status")

        # Get the outline agent client from config
        from agents.config import OUTLINE_AGENT_MODEL
        client, model = get_client(OUTLINE_AGENT_MODEL, wrap_with_instructor=False)

        scraped_context = ""
        scrape_result = None  # Will hold videos and scraped content if URL scraping happens
        file_context = ""

        # Variables to track file analysis results
        detected_intent = None
        detected_slide_style = None
        extracted_design_context = None
        extracted_file_images = []  # Images extracted from uploaded PPTX/PDF files
        extracted_slide_screenshots = []  # Visual reference screenshots from uploaded PPTX/PDF

        # Analyze uploaded files if present
        if request.files and len(request.files) > 0:
            file_names = [f.name for f in request.files]
            logger.info(f"[OutlineAgent] Analyzing {len(request.files)} files: {file_names}")

            # Send status event for each file
            for i, file in enumerate(request.files):
                yield f"data: {json.dumps({'type': 'status', 'status': 'analyzing_file', 'message': f'Analyzing {file.name}...', 'file_index': i, 'file_name': file.name, 'total_files': len(request.files)})}\n\n"

            # Run both analyses in parallel for efficiency
            basic_analysis_task = analyze_files_for_presentation(request.files)
            enhanced_analysis_task = enhanced_file_analysis(request.files, request.message)

            file_analysis, enhanced_result = await asyncio.gather(
                basic_analysis_task,
                enhanced_analysis_task
            )

            # Extract enhanced analysis results
            if enhanced_result.get("success"):
                detected_intent = enhanced_result.get("intent")
                detected_slide_style = enhanced_result.get("slide_style")
                extracted_design_context = enhanced_result.get("design_context")
                extracted_file_images = enhanced_result.get("extracted_images", [])
                extracted_slide_screenshots = enhanced_result.get("slide_screenshots", [])
                logger.info(f"[OutlineAgent] Enhanced analysis: intent={detected_intent}, style={detected_slide_style}, images={len(extracted_file_images)}, screenshots={len(extracted_slide_screenshots)}")

            if file_analysis["success"] and file_analysis["combined_context"]:
                # Build file context with intent and style information
                intent_info = ""
                if detected_intent:
                    intent_info = f"\n[FILE INTENT]: {detected_intent}"
                    if detected_intent == "use_design_only":
                        intent_info += " (User wants to USE THE DESIGN/STYLE from these files, NOT the content)"
                    elif detected_intent == "use_content_only":
                        intent_info += " (User wants to USE THE CONTENT from these files, NOT the design)"
                    elif detected_intent == "recreate_exact":
                        intent_info += " (User wants to RECREATE these files exactly)"
                    elif detected_intent == "use_both":
                        intent_info += " (User wants to use BOTH design AND content from these files)"

                style_info = ""
                if detected_slide_style:
                    style_info = f"\n[PREFERRED SLIDE STYLE]: {detected_slide_style}"
                    if detected_slide_style == "interactive":
                        style_info += " (User wants INTERACTIVE slides with animations)"
                    elif detected_slide_style == "traditional":
                        style_info += " (User wants TRADITIONAL simple slides)"

                design_info = ""
                if extracted_design_context:
                    design_info = f"\n[EXTRACTED DESIGN]:\n{json.dumps(extracted_design_context, indent=2)}"
                    design_info += "\n(Use these colors/fonts when generating theme)"

                file_context = f"\n\n[UPLOADED FILES ANALYSIS]{intent_info}{style_info}{design_info}\n{file_analysis['combined_context']}\n[END FILES ANALYSIS]\n"
                file_count = file_analysis['file_count']
                logger.info(f"[OutlineAgent] File analysis complete: {file_count} files, {len(file_context)} chars context")

                # Send file analysis complete event with results
                # Include extracted design for frontend to preserve across chat turns
                event_data = {
                    'type': 'status',
                    'status': 'files_analyzed',
                    'message': f'Analyzed {file_count} file(s)',
                    'analyses': file_analysis['analyses'],
                    'detected_intent': detected_intent,
                    'detected_slide_style': detected_slide_style,
                    'has_design': extracted_design_context is not None,
                    # CRITICAL: Send extracted design so frontend can preserve it across turns
                    'extracted_design': extracted_design_context,
                    'slide_screenshots_count': len(extracted_slide_screenshots) if extracted_slide_screenshots else 0
                }
                yield f"data: {json.dumps(event_data)}\n\n"
            else:
                logger.warning(f"[OutlineAgent] File analysis failed or empty: {file_analysis.get('error', 'No context')}")
                yield f"data: {json.dumps({'type': 'status', 'status': 'file_analysis_error', 'message': file_analysis.get('error', 'Could not analyze files')})}\n\n"

        # If no files sent but previousFileAnalysis exists in context, use it
        # This allows continued conversation without re-analyzing files
        if not file_context and request.context and request.context.get('previousFileAnalysis'):
            previous_analysis = request.context['previousFileAnalysis']
            file_context = f"\n\n[PREVIOUSLY ANALYZED FILES]\n{previous_analysis}\n(Files were analyzed earlier in this conversation - use the chat history for full context)\n"
            logger.info(f"[OutlineAgent] Using previous file analysis context: {len(file_context)} chars")

            # CRITICAL: Also restore extracted design context, intent, and screenshots from previous analysis
            # This ensures follow-up messages preserve the user's extracted design colors
            if request.context.get('previousExtractedDesign'):
                extracted_design_context = request.context['previousExtractedDesign']
                logger.info(f"[OutlineAgent] ✅ Restored extracted design from previous turn: {list(extracted_design_context.get('color_palette', {}).keys())}")
            if request.context.get('previousFileIntent'):
                detected_intent = request.context['previousFileIntent']
                logger.info(f"[OutlineAgent] ✅ Restored file_intent from previous turn: {detected_intent}")
            if request.context.get('previousSlideScreenshots'):
                extracted_slide_screenshots = request.context['previousSlideScreenshots']
                logger.info(f"[OutlineAgent] ✅ Restored {len(extracted_slide_screenshots)} slide screenshots from previous turn")

        # Detect URLs in the message and auto-scrape them
        url_pattern = re.compile(r'https?://[^\s<>"{}|\\^`\[\]]+|(?:www\.)?([a-zA-Z0-9][-a-zA-Z0-9]*\.(?:life|com|co|io|org|net|ai|app|xyz|dev)(?:/[^\s]*)?)', re.IGNORECASE)
        detected_urls = url_pattern.findall(request.message)

        # Also check for domain mentions like "numi.life" without http
        domain_pattern = re.compile(r'\b([a-zA-Z0-9][-a-zA-Z0-9]*\.(?:life|com|co|io|org|net|ai|app|xyz|dev))\b', re.IGNORECASE)
        detected_domains = domain_pattern.findall(request.message)

        # Combine and dedupe
        urls_to_scrape = []
        for url in detected_urls:
            if url and not url.startswith('http'):
                url = f'https://{url}'
            if url:
                urls_to_scrape.append(url)
        for domain in detected_domains:
            url = f'https://{domain}'
            if url not in urls_to_scrape:
                urls_to_scrape.append(url)

        # Handle reference links if user explicitly provided them (from link button)
        reference_links = []
        if request.context:
            reference_links = request.context.get("reference_links") or request.context.get("referenceLinks") or []

        # Add detected URLs to reference links
        urls_to_scrape.extend(reference_links)
        urls_to_scrape = list(set(urls_to_scrape))[:3]  # Dedupe and limit

        if urls_to_scrape:
            logger.info(f"[OutlineAgent] Auto-detected URLs to scrape: {urls_to_scrape}")
            yield f"data: {json.dumps({'type': 'status', 'status': 'scraping', 'message': f'Reading content from {urls_to_scrape[0]}...'})}\n\n"

            scrape_result = await scrape_reference_links(urls_to_scrape)
            if scrape_result["success"] and scrape_result["scraped_content"]:
                scraped_parts = []
                for item in scrape_result["scraped_content"]:
                    title = item.get("title") or item.get("url", "Reference")
                    scraped_parts.append(f"--- {title} ---\n{item['content']}\n---")
                scraped_context = "\n\n[REFERENCE CONTENT]\n" + "\n".join(scraped_parts) + "\n[END REFERENCE CONTENT]\n\n"

                # Add video information if found
                videos = scrape_result.get("videos", [])
                if videos:
                    video_parts = ["\n\n[AVAILABLE VIDEOS FROM WEBSITE]"]
                    video_parts.append(f"Found {len(videos)} video(s) that can be embedded in slides:\n")
                    for i, video in enumerate(videos[:5]):  # Limit to top 5
                        video_url = video.get('url', video.get('embed_url', ''))
                        video_type = video.get('source_type', 'direct')
                        video_title = video.get('title', 'Untitled')
                        thumbnail = video.get('thumbnail', '')
                        video_parts.append(f"{i+1}. [{video_type}] {video_title}")
                        video_parts.append(f"   URL: {video_url}")
                        if thumbnail:
                            video_parts.append(f"   Thumbnail: {thumbnail}")
                    video_parts.append("\nYou can reference these videos in slide content suggestions.")
                    video_parts.append("[END AVAILABLE VIDEOS]\n")
                    scraped_context += "\n".join(video_parts)
                    logger.info(f"[OutlineAgent] 🎬 Added {len(videos)} videos to context")
                    yield f"data: {json.dumps({'type': 'status', 'status': 'videos_found', 'message': f'Found {len(videos)} video(s) from website'})}\n\n"

                logger.info(f"[OutlineAgent] Added scraped context: {len(scraped_context)} chars")
                count = len(scrape_result["scraped_content"])
                yield f"data: {json.dumps({'type': 'status', 'status': 'scraped', 'message': f'Extracted content from {count} link(s)'})}\n\n"

        # Build message history - filter out empty messages
        messages = []
        for msg in request.chat_history:
            if msg.content and msg.content.strip():
                messages.append({
                    "role": msg.role,
                    "content": msg.content
                })

        # Build user message with context if available
        user_content = request.message + scraped_context + file_context

        # If context has current_outline, append it to the message
        if request.context and "current_outline" in request.context:
            outline = request.context["current_outline"]
            outline_json = json.dumps({
                "title": outline.get("title", ""),
                "slides": [{
                    "index": slide["index"],
                    "title": slide["title"],
                    "subtitle": slide.get("subtitle", ""),
                    "content": slide.get("content", ""),
                    "key_points": slide.get("key_points", [])
                } for slide in outline.get("slides", [])]
            }, indent=2)
            user_content = f"{user_content}\n\n[CURRENT OUTLINE]\n```json\n{outline_json}\n```"
            logger.info(f"[OutlineAgent] Added outline context with {len(outline.get('slides', []))} slides")

            # If target_slide_index is set, add it to the context
            if "target_slide_index" in request.context:
                target_idx = request.context["target_slide_index"]
                user_content = f"{user_content}\n\n[TARGET_SLIDE_INDEX]\n{target_idx}\n(User wants to edit slide {target_idx + 1} specifically)"
                logger.info(f"[OutlineAgent] Targeting specific slide: {target_idx}")

        # Add current user message
        messages.append({
            "role": "user",
            "content": user_content
        })

        logger.info(f"[OutlineAgent] Processing message with {len(messages)} messages in history")

        # Call Anthropic API with tool support - model decides when to search
        async def call_model_with_tools(msgs, depth=0):
            """Call Claude with optional tool use, handle search tool if called.

            Supports multi-layer research:
            - depth=0: Initial response, may trigger first round of searches
            - depth=1: After first search results, may do follow-up deep dives
            - depth=2: After follow-up searches, may do final refinements
            - depth=3: Force synthesis (no more searches)
            """
            if depth >= 3:
                logger.warning(f"[OutlineAgent] Max tool depth reached (depth={depth}), forcing synthesis")
                # Send status event so frontend shows "compiling research" instead of last search query
                yield f"data: {json.dumps({'type': 'status', 'status': 'compiling', 'message': 'Compiling research into outline...'})}\n\n"
                yield ("text", "I've gathered extensive research. Let me synthesize everything into your outline.")

                # Force a final call without tools to get the synthesis
                logger.info("[OutlineAgent] Making final synthesis call (no tools)...")

                # Build synthesis messages - msgs already ends with user (tool_results)
                # We need to append our synthesis request to that user message, not add a new one
                synthesis_msgs = []
                for i, msg in enumerate(msgs):
                    if i == len(msgs) - 1 and msg["role"] == "user":
                        # Append synthesis instruction to the last user message (tool results)
                        existing_content = msg["content"]
                        if isinstance(existing_content, list):
                            # Tool results are a list - add our text as another item
                            new_content = existing_content + [{
                                "type": "text",
                                "text": """\n\n---\nYou have completed your research. Now create the presentation outline using ALL the research data above.

Output ONLY a JSON code block with the generate_outline action in this EXACT format:

```json
{"action": "generate_outline", "slide_count": 10, "topic": "Presentation Topic", "brandContext": "company.com", "slides": [{"title": "Slide Title", "key_points": ["**$86B** Q4 revenue", "**35%** profit growth"], "content": "Detailed narrative with researched facts [1]. Include specific statistics and numbers [2].\\n\\nSources:\\n[1] source.com\\n[2] data.com"}]}
```

IMPORTANT FORMAT RULES:
- "slides" array is at the TOP LEVEL (not nested under "outline")
- "topic" field contains the presentation title
- "slide_count" matches the number of slides
- Each slide MUST have BOTH:
  * "key_points": Array of 2-4 SHORT bullet points (5-10 words each) with **bold** numbers
  * "content": Detailed narrative paragraph with [1] [2] citations and Sources section
- Include 8-12 slides covering all the key topics researched
- Put specific facts, numbers, and statistics from the research in BOTH key_points AND content"""
                            }]
                            synthesis_msgs.append({"role": "user", "content": new_content})
                        else:
                            # String content - append directly
                            synthesis_msgs.append({"role": "user", "content": existing_content + "\n\n---\nNow create the presentation outline using ALL the research above. Output ONLY a JSON code block with generate_outline action. Use format: {\"action\": \"generate_outline\", \"slide_count\": N, \"topic\": \"Title\", \"slides\": [...]}"})
                    else:
                        synthesis_msgs.append(msg)

                logger.info(f"[OutlineAgent] Synthesis messages: {len(synthesis_msgs)} messages, last role={synthesis_msgs[-1]['role'] if synthesis_msgs else 'EMPTY'}")

                current_date = datetime.now().strftime("%B %d, %Y")
                current_year = datetime.now().strftime("%Y")
                synthesis_system = f"""Today's date is {current_date}. Current year: {current_year}. You are a presentation outline generator.

The user has gathered research and now wants you to synthesize it into a presentation outline.
Output ONLY a JSON code block with the generate_outline action. No other text.

Required JSON format (slides at TOP LEVEL, not nested):
{{"action": "generate_outline", "slide_count": N, "topic": "Title", "brandContext": "domain.com", "slides": [{{"title": "...", "key_points": ["**$X** metric", "**Y%** growth"], "content": "Detailed narrative with facts [1][2]...\\n\\nSources:\\n[1] url"}}]}}

CRITICAL: Each slide MUST have BOTH key_points (short array for preview) AND content (detailed narrative for generation)."""

                try:
                    # Detect provider for synthesis call
                    use_gemini_synthesis = is_gemini_model(model)

                    if use_gemini_synthesis:
                        # Gemini synthesis call - use simple dict format
                        gemini_synthesis_contents = []
                        for msg in synthesis_msgs:
                            role = "user" if msg["role"] == "user" else "model"
                            content = msg["content"]
                            if isinstance(content, str):
                                gemini_synthesis_contents.append({"role": role, "parts": [{"text": content}]})
                            elif isinstance(content, list):
                                text_parts = [{"text": item.get("text", "") if isinstance(item, dict) else str(item)} for item in content]
                                if text_parts:
                                    gemini_synthesis_contents.append({"role": role, "parts": text_parts})

                        final_response = await asyncio.wait_for(
                            asyncio.to_thread(
                                client.models.generate_content,
                                model=model,
                                contents=gemini_synthesis_contents,
                                config=genai_types.GenerateContentConfig(
                                    system_instruction=synthesis_system,
                                    max_output_tokens=8192,
                                    temperature=0.7
                                )
                            ),
                            timeout=90.0
                        )

                        # Extract text from Gemini response
                        synthesis_text = ""
                        if final_response.candidates and len(final_response.candidates) > 0:
                            candidate = final_response.candidates[0]
                            if candidate.content and candidate.content.parts:
                                for part in candidate.content.parts:
                                    if hasattr(part, 'text') and part.text:
                                        synthesis_text += part.text

                        logger.info(f"[OutlineAgent] Gemini synthesis response: {len(synthesis_text)} chars")
                        if synthesis_text:
                            logger.info(f"[OutlineAgent] Synthesis preview: {synthesis_text[:500]}...")
                            yield ("text", synthesis_text)
                        else:
                            logger.warning("[OutlineAgent] Gemini synthesis returned empty text")
                    else:
                        # Anthropic synthesis call
                        final_response = await asyncio.wait_for(
                            asyncio.to_thread(
                                client.messages.create,
                                model=model,
                                max_tokens=8192,
                                system=synthesis_system,
                                messages=synthesis_msgs,
                                temperature=0.7
                            ),
                            timeout=90.0  # 90 second timeout for synthesis
                        )
                        logger.info(f"[OutlineAgent] Synthesis response: stop_reason={final_response.stop_reason}, blocks={len(final_response.content)}")
                        for i, block in enumerate(final_response.content):
                            logger.info(f"[OutlineAgent] Block {i}: type={type(block).__name__}, has_text={hasattr(block, 'text')}")
                            if hasattr(block, 'text') and block.text:
                                logger.info(f"[OutlineAgent] Final synthesis response: {len(block.text)} chars")
                                logger.info(f"[OutlineAgent] Synthesis preview: {block.text[:500]}...")
                                yield ("text", block.text)
                            else:
                                logger.warning(f"[OutlineAgent] Block {i} has no text or empty text")
                except asyncio.TimeoutError:
                    logger.error(f"[OutlineAgent] Final synthesis timed out (90s)")
                    yield ("text", "I'm taking too long to synthesize. Please try with a simpler request.")
                except Exception as e:
                    logger.error(f"[OutlineAgent] Final synthesis failed: {e}")
                return

            # Add current date to system prompt so model knows the year
            current_date = datetime.now().strftime("%B %d, %Y")
            current_year = datetime.now().strftime("%Y")
            system_with_date = f"TODAY'S DATE: {current_date}. CURRENT YEAR: {current_year}. ALWAYS use {current_year} (not 2024) in your web searches for recent data, financials, or statistics.\n\n{OUTLINE_AGENT_SYSTEM_PROMPT}"

            # Log before API call for debugging
            logger.info(f"[OutlineAgent] 📡 Making API call at depth={depth} with {len(msgs)} messages...")

            # Detect provider based on model name
            use_gemini = is_gemini_model(model)

            # Run API call in thread pool with timeout to prevent hanging
            try:
                if use_gemini:
                    # Gemini API call with function calling
                    logger.info(f"[OutlineAgent] Using Gemini provider for model: {model}")

                    # Convert messages to Gemini format - use simple dict format
                    gemini_contents = []
                    for msg in msgs:
                        role = "user" if msg["role"] == "user" else "model"
                        content = msg["content"]
                        if isinstance(content, str):
                            gemini_contents.append({"role": role, "parts": [{"text": content}]})
                        elif isinstance(content, list):
                            # Handle structured content (tool results)
                            text_parts = []
                            for item in content:
                                if isinstance(item, dict):
                                    if item.get("type") == "text":
                                        text_parts.append({"text": item.get("text", "")})
                                    elif item.get("type") == "tool_result":
                                        text_parts.append({"text": f"[Search Result]: {item.get('content', '')}"})
                                elif isinstance(item, str):
                                    text_parts.append({"text": item})
                            if text_parts:
                                gemini_contents.append({"role": role, "parts": text_parts})

                    response = await asyncio.wait_for(
                        asyncio.to_thread(
                            client.models.generate_content,
                            model=model,
                            contents=gemini_contents,
                            config=genai_types.GenerateContentConfig(
                                system_instruction=system_with_date,
                                max_output_tokens=8192,
                                temperature=0.7,
                                tools=[get_gemini_search_tool()]
                            )
                        ),
                        timeout=120.0
                    )
                else:
                    # Anthropic API call (existing code)
                    response = await asyncio.wait_for(
                        asyncio.to_thread(
                            client.messages.create,
                            model=model,
                            max_tokens=8192,
                            system=system_with_date,
                            messages=msgs,
                            tools=[SEARCH_TOOL],
                            temperature=0.7
                        ),
                        timeout=120.0  # 2 minute timeout
                    )
            except asyncio.TimeoutError:
                logger.error(f"[OutlineAgent] API call timed out at depth={depth}")
                yield ("text", "I'm taking too long to process. Let me give you what I have so far.")
                return
            except Exception as e:
                logger.error(f"[OutlineAgent] API call failed at depth={depth}: {type(e).__name__}: {e}")
                yield ("text", f"I encountered an issue while researching. Let me work with what I have so far.")
                return

            # Parse response based on provider
            if use_gemini:
                # Gemini response handling
                stop_reason = "end_turn"
                text_content = ""
                function_calls = []

                if response.candidates and len(response.candidates) > 0:
                    candidate = response.candidates[0]
                    if candidate.content and candidate.content.parts:
                        for part in candidate.content.parts:
                            if hasattr(part, 'text') and part.text:
                                text_content += part.text
                            elif hasattr(part, 'function_call') and part.function_call:
                                stop_reason = "tool_use"
                                function_calls.append(part.function_call)

                logger.info(f"[OutlineAgent] ✅ Gemini response received at depth={depth}, stop_reason={stop_reason}, function_calls={len(function_calls)}")
            else:
                # Anthropic response handling (existing)
                stop_reason = response.stop_reason
                logger.info(f"[OutlineAgent] ✅ API response received at depth={depth}, stop_reason={stop_reason}")

            # Handle max_tokens case - response was truncated
            if stop_reason == "max_tokens" or (use_gemini and response.candidates and response.candidates[0].finish_reason == "MAX_TOKENS"):
                logger.warning("[OutlineAgent] Response hit max_tokens limit - output may be incomplete")
                if use_gemini:
                    if text_content:
                        yield ("text", text_content)
                else:
                    for block in response.content:
                        if hasattr(block, 'text') and block.text:
                            yield ("text", block.text)
                yield f"data: {json.dumps({'type': 'error', 'message': 'Response was truncated. Try asking for fewer slides or less detail.'})}\n\n"
                return

            # Check if model wants to use tools
            if stop_reason == "tool_use":
                # Collect ALL tool calls and their results
                tool_results = []

                # First, yield any text that came before the tool call
                if use_gemini:
                    if text_content:
                        yield ("text", text_content)
                else:
                    for block in response.content:
                        if hasattr(block, 'text') and block.text:
                            yield ("text", block.text)

                # Collect all search queries to run in parallel
                search_tasks = []
                brand_blocked = []
                brand_search_keywords = ['brand color', 'brand colours', 'official color', 'hex code', 'logo', 'brand guideline', 'brand identity', 'color palette', 'official font', 'typography guide']

                if use_gemini:
                    # Gemini function calls
                    for i, fc in enumerate(function_calls):
                        if fc.name == "web_search":
                            query = fc.args.get("query", "") if hasattr(fc, 'args') else ""
                            query_lower = query.lower()
                            is_brand_search = any(kw in query_lower for kw in brand_search_keywords)

                            if is_brand_search:
                                logger.info(f"[OutlineAgent] 🚫 BLOCKED brand search: {query}")
                                brand_blocked.append({
                                    "block_id": f"gemini_fc_{i}",
                                    "content": "DO NOT search for brand colors, logos, or fonts. Instead, set brandContext in your generate_outline or update_theme JSON."
                                })
                            else:
                                logger.info(f"[OutlineAgent] Model requested search: {query}")
                                search_tasks.append({
                                    "block_id": f"gemini_fc_{i}",
                                    "query": query
                                })
                else:
                    # Anthropic tool use blocks
                    for block in response.content:
                        if block.type == "tool_use" and block.name == "web_search":
                            query = block.input.get("query", "")
                            query_lower = query.lower()
                            is_brand_search = any(kw in query_lower for kw in brand_search_keywords)

                            if is_brand_search:
                                logger.info(f"[OutlineAgent] 🚫 BLOCKED brand search: {query}")
                                brand_blocked.append({
                                    "block_id": block.id,
                                    "content": "DO NOT search for brand colors, logos, or fonts. Instead, set brandContext in your generate_outline or update_theme JSON (e.g., 'brandContext': 'ualberta.ca') and the ThemeDirector will fetch official brand data from Brandfetch."
                                })
                            else:
                                logger.info(f"[OutlineAgent] Model requested search: {query}")
                                search_tasks.append({
                                    "block_id": block.id,
                                    "query": query
                                })

                # Add blocked brand searches to results
                for blocked in brand_blocked:
                    tool_results.append({
                        "type": "tool_result",
                        "tool_use_id": blocked["block_id"],
                        "content": blocked["content"]
                    })

                # Run all searches in PARALLEL
                if search_tasks:
                    # Emit individual status event for EACH query (line by line)
                    queries = [t["query"] for t in search_tasks]
                    logger.info(f"[OutlineAgent] 🚀 Running {len(search_tasks)} searches in PARALLEL: {queries}")

                    for task in search_tasks:
                        status_event = f"data: {json.dumps({'type': 'status', 'status': 'researching', 'query': task['query']})}\n\n"
                        yield status_event

                    # Execute all searches concurrently
                    async def search_with_id(task):
                        result = await research_with_perplexity(task["query"])
                        return {"block_id": task["block_id"], "query": task["query"], "result": result}

                    search_results = await asyncio.gather(*[search_with_id(t) for t in search_tasks])

                    # Process all results
                    for sr in search_results:
                        if sr["result"]["success"]:
                            tool_result = sr["result"]["content"]
                            # Truncate research to prevent context overflow
                            if len(tool_result) > 2500:
                                tool_result = tool_result[:2500] + "\n\n[Truncated for brevity - use key facts above]"
                            logger.info(f"[OutlineAgent] Search success for '{sr['query'][:50]}': {len(tool_result)} chars")
                            research_event = f"data: {json.dumps({'type': 'research', 'content': sr['result']['content'][:500] + '...', 'citations': sr['result']['citations'][:5], 'query': sr['query']})}\n\n"
                            yield research_event
                        else:
                            tool_result = f"Search failed: {sr['result'].get('error', 'Unknown error')}"
                            logger.warning(f"[OutlineAgent] Search failed for '{sr['query']}': {tool_result}")
                            yield f"data: {json.dumps({'type': 'status', 'status': 'research_failed', 'message': sr['result'].get('error'), 'query': sr['query']})}\n\n"

                        tool_results.append({
                            "type": "tool_result",
                            "tool_use_id": sr["block_id"],
                            "content": tool_result
                        })

                    logger.info(f"[OutlineAgent] ✅ All {len(search_tasks)} parallel searches completed")

                # Build proper message format for next call
                if tool_results:
                    if use_gemini:
                        # For Gemini, build simpler message format with tool results as text
                        assistant_text = text_content if text_content else "I'll search for that information."
                        results_text = "\n\n".join([f"[Search Result for '{t.get('query', 'query')}']:\n{r['content']}" for t, r in zip(search_tasks, tool_results)])

                        new_msgs = msgs + [
                            {"role": "assistant", "content": assistant_text},
                            {"role": "user", "content": f"Here are the search results:\n\n{results_text}"}
                        ]
                    else:
                        # Anthropic format - convert response.content to proper format
                        assistant_content = []
                        for block in response.content:
                            if block.type == "text":
                                assistant_content.append({"type": "text", "text": block.text})
                            elif block.type == "tool_use":
                                assistant_content.append({
                                    "type": "tool_use",
                                    "id": block.id,
                                    "name": block.name,
                                    "input": block.input
                                })

                        new_msgs = msgs + [
                            {"role": "assistant", "content": assistant_content},
                            {"role": "user", "content": tool_results}
                        ]

                    # Log context size to help debug slow responses
                    total_chars = sum(len(str(m)) for m in new_msgs)
                    logger.info(f"[OutlineAgent] Calling model again with tool results (depth={depth+1}, context_size={total_chars} chars, {len(tool_results)} results)")
                    logger.info(f"[OutlineAgent] Tool result preview: {tool_results[0]['content'][:200] if tool_results else 'EMPTY'}...")

                    # Call model again with results
                    async for event in call_model_with_tools(new_msgs, depth + 1):
                        yield event
                    return

            # No tool use or end_turn - return the text content
            if use_gemini:
                if text_content:
                    yield ("text", text_content)
            else:
                for block in response.content:
                    if hasattr(block, 'text') and block.text:
                        yield ("text", block.text)

        # Run the model and collect response
        full_response = ""
        in_json_block = False
        streamed_text = ""  # Track what we've already streamed

        # Patterns that indicate "thinking" text we shouldn't stream
        THINKING_PATTERNS = [
            "let me research", "let me search", "let me look", "let me find",
            "i'll research", "i'll search", "i'll look up", "i'll find",
            "researching", "searching for", "looking up", "gathering",
            "i've gathered", "i've researched", "i've found",
            "now let me", "first, let me", "perfect!", "great!",
        ]

        def is_thinking_text(text: str) -> bool:
            """Check if text is internal 'thinking' that shouldn't be shown."""
            lower = text.lower()
            return any(pattern in lower for pattern in THINKING_PATTERNS)

        def contains_json_start(text: str) -> bool:
            """Check if text contains the start of JSON (fenced or raw)."""
            return '```json' in text or '{"action"' in text or '"action":' in text

        async for result in call_model_with_tools(messages):
            if isinstance(result, str) and result.startswith("data:"):
                # This is a status event, yield it directly
                logger.info(f"[OutlineAgent] Yielding status event: {result[:100]}...")
                yield result
            elif isinstance(result, tuple) and result[0] == "text":
                text = result[1]
                full_response += text
                logger.info(f"[OutlineAgent] 📝 Received text chunk: {len(text)} chars (total: {len(full_response)} chars)")

                # Detect JSON blocks (fenced or raw)
                if contains_json_start(full_response) and not in_json_block:
                    in_json_block = True
                    # Don't stream text before JSON - it's usually "thinking" text
                    logger.info(f"[OutlineAgent] Detected JSON block, suppressing pre-JSON text")
                elif not in_json_block:
                    # Only stream text that isn't "thinking" and isn't about to be JSON
                    if not is_thinking_text(text) and not contains_json_start(text):
                        yield f"data: {json.dumps({'type': 'text', 'content': text})}\n\n"
                        streamed_text += text

        logger.info(f"[OutlineAgent] 🏁 Model loop complete. Total response: {len(full_response)} chars")

        # After streaming, extract JSON and any text after it
        # Find ALL JSON blocks to handle cases where agent outputs multiple actions
        found_json_blocks = []

        # 1. Search for code blocks
        for match in re.finditer(r'```json\s*(\{[\s\S]*?\})\s*```', full_response):
            try:
                data = json.loads(match.group(1))
                found_json_blocks.append({
                    'data': data,
                    'end_index': match.end()
                })
            except:
                pass

        # 2. If no code blocks, search for raw JSON with action
        if not found_json_blocks:
             # Basic regex for simple cases (nested braces are hard with regex)
             match = re.search(r'(\{[\s\S]*?"action"[\s\S]*?\})', full_response)
             if match:
                 try:
                     data = json.loads(match.group(1))
                     found_json_blocks.append({
                         'data': data,
                         'end_index': match.end()
                     })
                 except:
                     pass

        # Check if there's already an outline in context
        has_existing_outline = request.context and "current_outline" in request.context

        # Choose the appropriate action based on context
        chosen_block = None

        if has_existing_outline:
            # When there's an existing outline, prioritize UPDATE actions over generate
            # This prevents accidental regeneration when user just wants to modify theme/slides
            for block in found_json_blocks:
                action = block['data'].get('action')
                if action in ('update_theme', 'update_slides', 'update_outline', 'scrape_media'):
                    chosen_block = block
                    logger.info(f"[OutlineAgent] Existing outline present - chose update action: {action}")
                    break

            # Only fall back to generate_outline if it's the ONLY option
            if not chosen_block and found_json_blocks:
                chosen_block = found_json_blocks[0]
                if chosen_block['data'].get('action') == 'generate_outline':
                    logger.warning(f"[OutlineAgent] Existing outline present but only generate_outline found - may cause regeneration")
        else:
            # No existing outline - prioritize generate_outline for new presentations
            for block in found_json_blocks:
                if block['data'].get('action') == 'generate_outline':
                    chosen_block = block
                    break

            # If no generate_outline, take the first one
            if not chosen_block and found_json_blocks:
                chosen_block = found_json_blocks[0]
            
        if chosen_block:
            outline_data = chosen_block['data']
            text_after_json = full_response[chosen_block['end_index']:].strip()

            logger.info(f"[OutlineAgent] Extracted outline data: {outline_data.get('action')}")

            # Handle scrape_media action - scrape media from URL before sending response
            if outline_data.get('action') == 'scrape_media':
                url = outline_data.get('url')
                media_filter = outline_data.get('media_filter', 'all')
                slide_index = outline_data.get('slide_index')
                content_context = outline_data.get('content_context', '')

                if url:
                    # Send status event
                    yield f"data: {json.dumps({'type': 'status', 'status': 'scraping_media', 'message': f'Pulling media from {url}...'})}\n\n"

                    # Scrape the media
                    media_result = await scrape_media_from_url(url, media_filter)

                    if media_result.get('success'):
                        # Include scraped media in the action data
                        outline_data['scraped_media'] = {
                            'gifs': media_result.get('gifs', []),
                            'images': media_result.get('images', []),
                            'all_media': media_result.get('all_media', []),
                            'filtered_media': media_result.get('filtered_media', []),
                            'source_url': url,
                            'markdown': media_result.get('markdown', ''),
                            'content_context': content_context,
                        }
                        gif_count = len(media_result.get('gifs', []))
                        img_count = len(media_result.get('images', []))
                        yield f"data: {json.dumps({'type': 'status', 'status': 'media_scraped', 'message': f'Found {gif_count} GIFs, {img_count} images'})}\n\n"
                        logger.info(f"[OutlineAgent] Scraped media attached: {gif_count} GIFs, {img_count} images")
                    else:
                        error_msg = media_result.get('error', 'Unknown error')
                        yield f"data: {json.dumps({'type': 'status', 'status': 'media_scrape_failed', 'message': f'Could not fetch media: {error_msg}'})}\n\n"
                        logger.warning(f"[OutlineAgent] Media scrape failed: {error_msg}")

            # Attach scraped videos to generate_outline action
            if outline_data.get('action') == 'generate_outline' and scrape_result and scrape_result.get('videos'):
                outline_data['scraped_videos'] = scrape_result['videos']
                logger.info(f"[OutlineAgent] 🎬 Attached {len(scrape_result['videos'])} scraped videos to outline")

            # Attach extracted design and slide style to generate_outline action
            if outline_data.get('action') == 'generate_outline':
                if extracted_design_context:
                    outline_data['extracted_design'] = extracted_design_context
                    logger.info(f"[OutlineAgent] 🎨 Attached extracted design context to outline")

                    # CRITICAL: Convert extracted_design to stylePreferences.colors for frontend theme tab
                    # Frontend expects stylePreferences.colors with background, text, accent1, accent2
                    color_palette = extracted_design_context.get('color_palette', {})
                    if color_palette:
                        style_colors = {
                            'type': 'custom',
                            'background': color_palette.get('background') or color_palette.get('primary_background') or '#ffffff',
                            'text': color_palette.get('text') or color_palette.get('primary_text') or '#000000',
                            'accent1': color_palette.get('primary') or color_palette.get('accent') or color_palette.get('accent_1') or '#007bff',
                            'accent2': color_palette.get('secondary') or color_palette.get('accent_2') or None,
                            'accent3': color_palette.get('accent_3') or None
                        }
                        # Remove None values
                        style_colors = {k: v for k, v in style_colors.items() if v is not None}

                        # Set stylePreferences with colors so frontend theme tab works
                        if 'stylePreferences' not in outline_data:
                            outline_data['stylePreferences'] = {}
                        outline_data['stylePreferences']['colors'] = style_colors

                        # Also add font if available from typography
                        typography = extracted_design_context.get('typography', {})
                        if typography.get('hero_font'):
                            outline_data['stylePreferences']['font'] = typography['hero_font']
                        if typography.get('body_font'):
                            outline_data['stylePreferences']['bodyFont'] = typography['body_font']

                        logger.info(f"[OutlineAgent] 🎨 Set stylePreferences.colors from extracted design: {style_colors}")
                if extracted_file_images:
                    outline_data['extracted_images'] = extracted_file_images
                    logger.info(f"[OutlineAgent] 🖼️ Attached {len(extracted_file_images)} extracted images from uploaded files")
                # Attach slide screenshots for visual design replication
                # These are base64 PNGs that the CustomComponentGenerator will use to SEE and replicate the design
                if extracted_slide_screenshots:
                    outline_data['slide_screenshots'] = extracted_slide_screenshots
                    logger.info(f"[OutlineAgent] 📸 Attached {len(extracted_slide_screenshots)} slide screenshots for visual design reference")

                    # CRITICAL: Also set stylePreferences.referenceImages so slide generator can use them!
                    # The CustomComponent generator expects reference images in stylePreferences.referenceImages
                    if 'stylePreferences' not in outline_data:
                        outline_data['stylePreferences'] = {}
                    # Convert base64 screenshots to data URLs for the vision model
                    reference_data_urls = []
                    for screenshot in extracted_slide_screenshots[:3]:  # Limit to 3 for performance
                        if screenshot.startswith('data:'):
                            reference_data_urls.append(screenshot)
                        else:
                            # Assume PNG if no prefix
                            reference_data_urls.append(f"data:image/png;base64,{screenshot}")
                    outline_data['stylePreferences']['referenceImages'] = reference_data_urls
                    logger.info(f"[OutlineAgent] 📸 Set stylePreferences.referenceImages with {len(reference_data_urls)} screenshots for CustomComponent")
                if detected_slide_style and detected_slide_style != 'auto':
                    # Only override if AI didn't already set it
                    if not outline_data.get('slide_style'):
                        outline_data['slide_style'] = detected_slide_style
                        logger.info(f"[OutlineAgent] 🎨 Set slide_style to: {detected_slide_style}")
                if detected_intent:
                    outline_data['file_intent'] = detected_intent
                    logger.info(f"[OutlineAgent] 📋 Set file_intent to: {detected_intent}")

                # Pass brandContext and/or style to theme generator via stylePreferences.vibeContext
                # Priority: brandContext (specific brand) > style (general vibe)
                brand_context = outline_data.get('brandContext')
                style_context = outline_data.get('style')

                if brand_context or style_context:
                    if 'stylePreferences' not in outline_data:
                        outline_data['stylePreferences'] = {}
                    # Combine brand and style if both present, otherwise use whichever exists
                    if brand_context and style_context:
                        outline_data['stylePreferences']['vibeContext'] = brand_context
                        outline_data['stylePreferences']['style'] = style_context
                        logger.info(f"[OutlineAgent] 🏷️ Set vibeContext={brand_context}, style={style_context}")
                    elif brand_context:
                        outline_data['stylePreferences']['vibeContext'] = brand_context
                        logger.info(f"[OutlineAgent] 🏷️ Set vibeContext for theme generator: {brand_context}")
                    else:
                        outline_data['stylePreferences']['vibeContext'] = style_context
                        outline_data['stylePreferences']['style'] = style_context
                        logger.info(f"[OutlineAgent] 🎨 Set style for theme generator: {style_context}")

            # Attach uploaded files to generate_outline action so they're used in slide generation
            logger.info(f"[OutlineAgent] 📎 Checking file attachment: action={outline_data.get('action')}, has_files={bool(request.files)}, file_count={len(request.files) if request.files else 0}")
            if outline_data.get('action') == 'generate_outline' and request.files:
                uploaded_media = []
                for f in request.files:
                    logger.info(f"[OutlineAgent] 📄 Processing file: {f.name}, type={f.type}, is_image={f.type and f.type.startswith('image/')}")
                    # Only include images for now (they're the ones that should appear on slides)
                    if f.type and f.type.startswith('image/'):
                        uploaded_media.append({
                            'id': f.id,
                            'name': f.name,
                            'type': f.type,
                            'content': f.content,  # Base64 content for rendering
                            'size': f.size
                        })

                if uploaded_media:
                    outline_data['uploadedMedia'] = uploaded_media
                    logger.info(f"[OutlineAgent] 📎 Attached {len(uploaded_media)} uploaded images to outline")
                else:
                    logger.warning(f"[OutlineAgent] ⚠️ No image files to attach (files had non-image types)")

            yield f"data: {json.dumps({'type': 'outline', 'data': outline_data})}\n\n"

            if text_after_json:
                 yield f"data: {json.dumps({'type': 'text', 'content': text_after_json})}\n\n"

        # Send done event
        yield f"data: {json.dumps({'type': 'done'})}\n\n"

    except Exception as e:
        logger.error(f"[OutlineAgent] Error in stream: {str(e)}", exc_info=True)
        error_msg = f"I encountered an error: {str(e)}. Could you try rephrasing your request?"
        yield f"data: {json.dumps({'type': 'error', 'message': error_msg})}\n\n"


@router.post("/chat")
async def outline_agent_chat(
    request: OutlineAgentRequest,
    token: Optional[str] = Depends(get_auth_header)
):
    """
    Chat with the outline generation agent.
    Returns a streaming response with the agent's conversation.
    """
    try:
        logger.info(f"[OutlineAgent] Received chat request: {request.message[:100]}")
        logger.info(f"[OutlineAgent] 📎 FILES RECEIVED: {len(request.files) if request.files else 0}")
        if request.files:
            for f in request.files:
                logger.info(f"[OutlineAgent] 📄 File: {f.name}, type={f.type}, has_content={bool(f.content)}, content_len={len(f.content) if f.content else 0}")

        return StreamingResponse(
            stream_agent_response(request),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no"
            }
        )

    except Exception as e:
        logger.error(f"[OutlineAgent] Error: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


class GenerateSlideContentRequest(BaseModel):
    """Request to generate detailed content for a single slide."""
    slide_title: str = Field(..., description="The title of the slide")
    slide_index: int = Field(..., description="The index of the slide in the presentation")
    total_slides: int = Field(..., description="Total number of slides in the presentation")
    presentation_topic: str = Field(..., description="The main topic of the presentation")
    presentation_context: Optional[str] = Field(None, description="Additional context about the presentation")
    existing_key_points: Optional[List[str]] = Field(None, description="Any existing key points for this slide")
    file_content: Optional[str] = Field(None, description="Content from uploaded files relevant to this slide")


class GenerateSlideContentResponse(BaseModel):
    """Response with generated slide content."""
    content: str = Field(..., description="The detailed narrative content for the slide")
    key_points: List[str] = Field(default_factory=list, description="Key bullet points for the slide")


@router.post("/generate-slide-content", response_model=GenerateSlideContentResponse)
async def generate_slide_content(
    request: GenerateSlideContentRequest,
    token: Optional[str] = Depends(get_auth_header)
):
    """
    Generate detailed content for a single slide on demand.
    Uses AI to create compelling narrative content and key points.
    """
    try:
        logger.info(f"[OutlineAgent] Generating content for slide: {request.slide_title}")

        client, model = get_client(OUTLINE_AGENT_MODEL, wrap_with_instructor=False)

        # Build the prompt
        prompt = f"""Generate detailed content for a presentation slide.

PRESENTATION CONTEXT:
- Topic: {request.presentation_topic}
- Slide {request.slide_index + 1} of {request.total_slides}
- Slide Title: "{request.slide_title}"
{f'- Additional Context: {request.presentation_context}' if request.presentation_context else ''}
{f'- Existing Key Points: {", ".join(request.existing_key_points)}' if request.existing_key_points else ''}
{f'- Source Material: {request.file_content[:2000]}...' if request.file_content and len(request.file_content) > 2000 else f'- Source Material: {request.file_content}' if request.file_content else ''}

Generate:
1. A compelling narrative paragraph (2-4 sentences) that explains the main point of this slide
2. 3-5 key bullet points that capture the essential information

Format your response as JSON:
{{
    "content": "The narrative paragraph...",
    "key_points": ["Point 1", "Point 2", "Point 3"]
}}

Make the content:
- Informative and specific (use real facts/data when available)
- Engaging and professional
- Appropriate for the slide's position in the presentation flow
- Concise but substantive"""

        response = await asyncio.to_thread(
            client.messages.create,
            model=model,
            max_tokens=1000,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )

        # Parse the response
        response_text = response.content[0].text

        # Try to extract JSON from the response
        try:
            # Find JSON in the response
            json_match = re.search(r'\{[\s\S]*\}', response_text)
            if json_match:
                result = json.loads(json_match.group())
                return GenerateSlideContentResponse(
                    content=result.get("content", ""),
                    key_points=result.get("key_points", [])
                )
        except json.JSONDecodeError:
            pass

        # Fallback: use the response as content
        return GenerateSlideContentResponse(
            content=response_text,
            key_points=request.existing_key_points or []
        )

    except Exception as e:
        logger.error(f"[OutlineAgent] Error generating slide content: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))
