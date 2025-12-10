"""
Unified Attachment Analyzer for Editing Agent

Handles all file types uploaded by users:
- Images: Downloads, resizes, encodes as base64 for Claude vision
- Spreadsheets: Parses CSV/Excel, extracts data, suggests charts
- Documents: Extracts text content for context
- Presentations: Extracts slide content

This provides a clean, consistent interface for the editing agent to work with any file type.
"""

import base64
import csv
import logging
import requests
from io import BytesIO, StringIO
from typing import List, Dict, Any, Optional, Tuple, TYPE_CHECKING
from dataclasses import dataclass, field
from enum import Enum
from PIL import Image

logger = logging.getLogger(__name__)

# Maximum image dimensions to keep token usage reasonable
# Claude vision works well with images up to 1568x1568, but we keep smaller
# to leave room for the rest of the prompt (which can be 100k+ tokens)
MAX_IMAGE_DIMENSION = 800   # pixels - smaller to reduce token usage
MAX_IMAGE_BYTES = 200_000   # 200KB max for base64 encoded image


class FileType(Enum):
    IMAGE = "image"
    SPREADSHEET = "spreadsheet"
    DOCUMENT = "document"
    PRESENTATION = "presentation"
    UNKNOWN = "unknown"


@dataclass
class AnalyzedAttachment:
    """Result of analyzing an attachment"""
    name: str
    file_type: FileType
    mime_type: str
    url: str

    # Original URL from user upload (before any processing)
    original_url: str = ""

    # For images - base64 encoded for Claude vision
    base64_data: Optional[str] = None

    # For spreadsheets - parsed data
    headers: List[str] = field(default_factory=list)
    rows: List[List[str]] = field(default_factory=list)
    chart_suggestion: Optional[Dict[str, str]] = None

    # For documents - extracted text
    text_content: Optional[str] = None

    # For PDFs - extracted images (list of (base64_data, mime_type) tuples)
    extracted_images: List[Tuple[str, str]] = field(default_factory=list)

    # Human-readable summary for prompt context
    summary: str = ""

    # Whether this attachment should be sent as vision content
    is_vision_content: bool = False


def _detect_file_type(mime_type: str, filename: str) -> FileType:
    """Detect file type from MIME type and filename"""
    mime_lower = mime_type.lower()
    name_lower = filename.lower()

    # Images
    if mime_lower.startswith('image/'):
        return FileType.IMAGE

    # Spreadsheets
    spreadsheet_mimes = [
        'text/csv',
        'application/vnd.ms-excel',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    ]
    if mime_lower in spreadsheet_mimes or name_lower.endswith(('.csv', '.xlsx', '.xls')):
        return FileType.SPREADSHEET

    # Presentations
    ppt_mimes = [
        'application/vnd.ms-powerpoint',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    ]
    if mime_lower in ppt_mimes or name_lower.endswith(('.ppt', '.pptx')):
        return FileType.PRESENTATION

    # Documents
    doc_mimes = [
        'application/pdf',
        'text/plain',
        'text/markdown',
        'application/msword',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    ]
    if mime_lower in doc_mimes or name_lower.endswith(('.pdf', '.txt', '.md', '.doc', '.docx')):
        return FileType.DOCUMENT

    return FileType.UNKNOWN


def _resize_image_for_vision(content: bytes, max_dim: int = MAX_IMAGE_DIMENSION) -> Tuple[bytes, str]:
    """
    Resize image to fit within max dimensions while preserving aspect ratio.
    Returns (resized_bytes, mime_type).

    This is critical for keeping token usage reasonable - large screenshots
    can easily exceed Claude's 200k token limit when base64 encoded.
    """
    try:
        img = Image.open(BytesIO(content))
        original_size = img.size

        # Convert RGBA to RGB if needed (for JPEG output)
        if img.mode in ('RGBA', 'P'):
            # Create white background for transparency
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[3] if img.mode == 'RGBA' else None)
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')

        # Calculate new size maintaining aspect ratio
        width, height = img.size
        if width > max_dim or height > max_dim:
            if width > height:
                new_width = max_dim
                new_height = int(height * (max_dim / width))
            else:
                new_height = max_dim
                new_width = int(width * (max_dim / height))

            img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            logger.info(f"[AttachmentAnalyzer] Resized image from {original_size} to {img.size}")

        # Save as JPEG with reasonable quality
        output = BytesIO()
        img.save(output, format='JPEG', quality=85, optimize=True)
        resized_bytes = output.getvalue()

        # If still too large, reduce quality further
        if len(resized_bytes) > MAX_IMAGE_BYTES:
            output = BytesIO()
            img.save(output, format='JPEG', quality=60, optimize=True)
            resized_bytes = output.getvalue()
            logger.info(f"[AttachmentAnalyzer] Reduced quality, size now {len(resized_bytes)} bytes")

        return resized_bytes, 'image/jpeg'

    except Exception as e:
        logger.warning(f"[AttachmentAnalyzer] Could not resize image: {e}")
        # Return original if resize fails
        return content, 'image/png'


def _suggest_chart_type(headers: List[str], rows: List[List[str]], filename: str) -> Dict[str, str]:
    """Suggest appropriate chart type based on data structure"""
    if not headers or not rows:
        return {"type": "bar", "reason": "Default chart type"}

    # Check for time-series patterns
    time_keywords = ['year', 'month', 'date', 'time', 'quarter', 'q1', 'q2', 'q3', 'q4']
    has_time = any(any(kw in h.lower() for kw in time_keywords) for h in headers)

    # Check for percentage data
    has_percentages = any('%' in str(cell) for row in rows for cell in row)

    # Detect numeric columns
    numeric_cols = 0
    for col_idx in range(len(headers)):
        is_numeric = True
        for row in rows[:10]:
            if col_idx < len(row):
                try:
                    float(str(row[col_idx]).replace(',', '').replace('$', '').replace('%', ''))
                except:
                    is_numeric = False
                    break
        if is_numeric:
            numeric_cols += 1

    if has_time and numeric_cols >= 1:
        return {"type": "line", "reason": "Time series data detected"}
    elif has_percentages and len(rows) < 10:
        return {"type": "pie", "reason": "Percentage/distribution data"}
    elif len(rows) < 20 and numeric_cols == 1:
        return {"type": "bar", "reason": "Categorical comparison"}

    # Fallback based on filename
    name_lower = filename.lower()
    if any(w in name_lower for w in ['revenue', 'sales', 'growth', 'trend']):
        return {"type": "line", "reason": "Filename suggests trend data"}
    elif any(w in name_lower for w in ['distribution', 'breakdown']):
        return {"type": "pie", "reason": "Filename suggests distribution"}

    return {"type": "bar", "reason": "Default for data visualization"}


def _parse_csv(content: bytes) -> tuple[List[str], List[List[str]]]:
    """Parse CSV content and return headers + rows"""
    try:
        text = content.decode('utf-8', errors='ignore')
        reader = csv.reader(StringIO(text))
        rows = list(reader)

        if not rows:
            return [], []

        headers = rows[0] if rows else []
        data_rows = rows[1:20] if len(rows) > 1 else []  # Limit to 20 rows

        return headers, data_rows
    except Exception as e:
        logger.warning(f"Error parsing CSV: {e}")
        return [], []


def _parse_excel(content: bytes) -> tuple[List[str], List[List[str]]]:
    """Parse Excel content and return headers + rows"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(BytesIO(content), read_only=True)
        sheet = wb.active

        rows = []
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= 20:  # Limit to 20 rows
                break
            rows.append([str(cell) if cell is not None else "" for cell in row])

        wb.close()

        if not rows:
            return [], []

        headers = rows[0]
        data_rows = rows[1:]

        return headers, data_rows
    except ImportError:
        logger.warning("openpyxl not available for Excel parsing")
        return [], []
    except Exception as e:
        logger.warning(f"Error parsing Excel: {e}")
        return [], []


def _extract_pdf_text(content: bytes, max_pages: int = 50, max_chars: int = 100000) -> str:
    """Extract text from PDF

    Args:
        content: PDF file bytes
        max_pages: Maximum number of pages to extract (default 50)
        max_chars: Maximum characters to return (default 100k)
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        text = ""
        for page in doc[:max_pages]:
            text += page.get_text()
            if len(text) >= max_chars:
                break
        doc.close()
        logger.info(f"[AttachmentAnalyzer] PDF extraction: {len(doc)} pages, extracted {len(text)} chars")
        return text[:max_chars]
    except ImportError:
        logger.warning("PyMuPDF not available for PDF extraction")
        return ""
    except Exception as e:
        logger.warning(f"Error extracting PDF text: {e}")
        return ""


def _extract_pdf_images(
    content: bytes,
    max_pages: int = 10,
    max_images: int = 5,
    min_image_size: int = 10000  # Skip tiny images (icons, etc.)
) -> List[Tuple[bytes, str]]:
    """Extract images from PDF pages.

    Args:
        content: PDF file bytes
        max_pages: Maximum pages to scan for images
        max_images: Maximum number of images to extract
        min_image_size: Minimum image size in bytes (skip smaller)

    Returns:
        List of (image_bytes, mime_type) tuples
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        images = []

        for page_num, page in enumerate(doc[:max_pages]):
            if len(images) >= max_images:
                break

            # Get images from page
            image_list = page.get_images(full=True)

            for img_index, img_info in enumerate(image_list):
                if len(images) >= max_images:
                    break

                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)

                    if not base_image:
                        continue

                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]

                    # Skip small images (likely icons, logos, decorations)
                    if len(image_bytes) < min_image_size:
                        continue

                    # Map extension to mime type
                    mime_map = {
                        "png": "image/png",
                        "jpeg": "image/jpeg",
                        "jpg": "image/jpeg",
                        "jp2": "image/jp2",
                        "jxr": "image/jxr",
                    }
                    mime_type = mime_map.get(image_ext.lower(), f"image/{image_ext}")

                    images.append((image_bytes, mime_type))
                    logger.info(f"[AttachmentAnalyzer] Extracted image from PDF page {page_num + 1}: {len(image_bytes)} bytes, {image_ext}")

                except Exception as e:
                    logger.debug(f"[AttachmentAnalyzer] Could not extract image {img_index} from page {page_num}: {e}")
                    continue

        doc.close()
        logger.info(f"[AttachmentAnalyzer] Extracted {len(images)} images from PDF")
        return images

    except ImportError:
        logger.warning("PyMuPDF not available for PDF image extraction")
        return []
    except Exception as e:
        logger.warning(f"Error extracting PDF images: {e}")
        return []


async def _smart_extract_relevant_sections(
    full_text: str,
    user_request: str,
    max_output_chars: int = 15000
) -> str:
    """
    Use Haiku to identify and extract only the relevant sections from a document.

    This is much more token-efficient than sending the entire document.
    ~15k chars ≈ 4k tokens, which is very reasonable.
    """
    from agents.ai_clients import get_client
    from agents.config import CLAUDE_HAIKU

    # If text is already short, just return it
    if len(full_text) <= max_output_chars:
        return full_text

    try:
        client = get_client(CLAUDE_HAIKU)

        # Split into chunks for analysis (Haiku can handle ~100k tokens)
        # We'll send the full text and ask it to extract relevant parts
        prompt = f"""You are a document analyst. Extract ONLY the sections from this document that are relevant to the user's request.

USER REQUEST: {user_request}

DOCUMENT TEXT:
{full_text[:80000]}

INSTRUCTIONS:
1. Identify sections relevant to the user's request
2. Extract those sections verbatim (don't summarize)
3. Include key data, numbers, and facts
4. Skip irrelevant sections (boilerplate, disclaimers, etc.)
5. Keep total output under {max_output_chars} characters
6. If the document has financial data, prioritize: key metrics, highlights, YoY comparisons

OUTPUT FORMAT:
Return only the extracted relevant text, with section headers if present."""

        response = await client.messages.create(
            model=CLAUDE_HAIKU,
            max_tokens=4096,
            messages=[{"role": "user", "content": prompt}]
        )

        extracted = response.content[0].text if response.content else ""
        logger.info(f"[AttachmentAnalyzer] Smart extraction: {len(full_text)} chars -> {len(extracted)} chars")
        return extracted[:max_output_chars]

    except Exception as e:
        logger.warning(f"[AttachmentAnalyzer] Smart extraction failed, using truncation: {e}")
        # Fallback: return first portion of text
        return full_text[:max_output_chars]


def _extract_pptx_content(content: bytes) -> str:
    """Extract text from PowerPoint"""
    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(content))

        text = f"Total slides: {len(prs.slides)}\n\n"
        for i, slide in enumerate(prs.slides[:10]):
            text += f"Slide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += f"  - {shape.text[:200]}\n"
            text += "\n"

        return text
    except ImportError:
        logger.warning("python-pptx not available for PowerPoint extraction")
        return ""
    except Exception as e:
        logger.warning(f"Error extracting PPTX content: {e}")
        return ""


def analyze_attachment(attachment: Dict[str, Any], timeout: int = 15) -> AnalyzedAttachment:
    """
    Analyze a single attachment and return structured data.

    Args:
        attachment: Dict with 'name', 'url', 'mimeType' or 'type'
        timeout: HTTP timeout for downloading

    Returns:
        AnalyzedAttachment with parsed data
    """
    name = attachment.get('name') or 'unknown'
    url = attachment.get('url') or ''
    mime_type = attachment.get('mimeType') or attachment.get('type') or ''

    file_type = _detect_file_type(mime_type, name)

    result = AnalyzedAttachment(
        name=name,
        file_type=file_type,
        mime_type=mime_type,
        url=url,
        original_url=url  # Keep original URL for direct use when needed
    )

    if not url:
        result.summary = f"[{name}]: No URL provided"
        return result

    # Download the file
    try:
        logger.info(f"[AttachmentAnalyzer] Downloading {name} from {url[:80]}...")
        resp = requests.get(url, timeout=timeout)

        if resp.status_code != 200:
            result.summary = f"[{name}]: Failed to download (HTTP {resp.status_code})"
            logger.warning(f"[AttachmentAnalyzer] Failed to download {name}: HTTP {resp.status_code}")
            return result

        content = resp.content
        logger.info(f"[AttachmentAnalyzer] Downloaded {name} ({len(content)} bytes)")

    except Exception as e:
        result.summary = f"[{name}]: Download error - {str(e)}"
        logger.warning(f"[AttachmentAnalyzer] Error downloading {name}: {e}")
        return result

    # Process based on file type
    if file_type == FileType.IMAGE:
        # Resize image to keep token usage reasonable (large screenshots can exceed 200k tokens!)
        original_size = len(content)
        resized_content, resized_mime = _resize_image_for_vision(content)
        result.base64_data = base64.b64encode(resized_content).decode('utf-8')
        result.mime_type = resized_mime  # Update mime type after resize (usually jpeg)
        result.is_vision_content = True
        result.summary = f"[Image: {name}] - Ready for visual analysis"
        logger.info(f"[AttachmentAnalyzer] Encoded image {name}: {original_size} bytes -> {len(resized_content)} bytes -> {len(result.base64_data)} chars base64")

    elif file_type == FileType.SPREADSHEET:
        if name.lower().endswith('.csv') or mime_type == 'text/csv':
            headers, rows = _parse_csv(content)
        else:
            headers, rows = _parse_excel(content)

        result.headers = headers
        result.rows = rows
        result.chart_suggestion = _suggest_chart_type(headers, rows, name)

        # Build summary
        if headers and rows:
            result.summary = f"[Data: {name}] {len(rows)} rows, {len(headers)} columns. "
            result.summary += f"Headers: {', '.join(headers[:5])}{'...' if len(headers) > 5 else ''}. "
            result.summary += f"Suggested chart: {result.chart_suggestion['type']} ({result.chart_suggestion['reason']})"

            # Add sample data preview
            if rows:
                result.text_content = f"Data preview from {name}:\n"
                result.text_content += f"Headers: {', '.join(headers)}\n"
                for row in rows[:5]:
                    result.text_content += f"  {', '.join(str(c)[:30] for c in row[:6])}\n"
        else:
            result.summary = f"[Data: {name}] Could not parse data"

        logger.info(f"[AttachmentAnalyzer] Parsed spreadsheet {name}: {len(rows)} rows, {len(headers)} cols")

    elif file_type == FileType.DOCUMENT:
        if name.lower().endswith('.pdf') or mime_type == 'application/pdf':
            text = _extract_pdf_text(content)

            # Also extract images from PDF (charts, graphs, tables as images)
            pdf_images = _extract_pdf_images(content, max_pages=20, max_images=10)
            if pdf_images:
                for img_bytes, img_mime in pdf_images:
                    # Resize image for vision
                    resized_bytes, resized_mime = _resize_image_for_vision(img_bytes)
                    img_base64 = base64.b64encode(resized_bytes).decode('utf-8')
                    result.extracted_images.append((img_base64, resized_mime))
                logger.info(f"[AttachmentAnalyzer] Added {len(result.extracted_images)} images from PDF for vision analysis")
        else:
            # Plain text
            text = content.decode('utf-8', errors='ignore')[:10000]

        result.text_content = text

        # Generate summary (first 200 chars or first paragraph)
        if text:
            first_para = text.split('\n\n')[0][:300]
            img_note = f" + {len(result.extracted_images)} images" if result.extracted_images else ""
            result.summary = f"[Document: {name}] {len(text)} chars{img_note}. Preview: {first_para}..."
        else:
            result.summary = f"[Document: {name}] Could not extract text"

        logger.info(f"[AttachmentAnalyzer] Extracted {len(text) if text else 0} chars from document {name}")

    elif file_type == FileType.PRESENTATION:
        text = _extract_pptx_content(content)
        result.text_content = text

        if text:
            result.summary = f"[Presentation: {name}] {text.split(chr(10))[0]}"
        else:
            result.summary = f"[Presentation: {name}] Could not extract content"

        logger.info(f"[AttachmentAnalyzer] Extracted content from presentation {name}")

    else:
        result.summary = f"[{name}]: Unsupported file type ({mime_type})"
        logger.warning(f"[AttachmentAnalyzer] Unsupported file type for {name}: {mime_type}")

    return result


def analyze_attachments(attachments: List[Dict[str, Any]]) -> List[AnalyzedAttachment]:
    """
    Analyze multiple attachments.

    Args:
        attachments: List of attachment dicts with 'name', 'url', 'mimeType'

    Returns:
        List of AnalyzedAttachment objects
    """
    if not attachments:
        return []

    results = []
    for att in attachments:
        result = analyze_attachment(att)
        results.append(result)

    return results


async def analyze_attachments_smart(
    attachments: List[Dict[str, Any]],
    user_request: str,
    max_doc_chars: int = 15000
) -> List[AnalyzedAttachment]:
    """
    Analyze attachments with smart extraction for documents.

    For PDFs and long documents, uses Haiku to extract only relevant sections
    based on the user's request. This significantly reduces token usage.

    Args:
        attachments: List of attachment dicts
        user_request: The user's request (used to identify relevant sections)
        max_doc_chars: Max chars for document content after smart extraction

    Returns:
        List of AnalyzedAttachment objects with optimized content
    """
    if not attachments:
        return []

    results = []
    for att in attachments:
        result = analyze_attachment(att)

        # Apply smart extraction for large documents
        if (result.file_type == FileType.DOCUMENT and
            result.text_content and
            len(result.text_content) > max_doc_chars):

            logger.info(f"[AttachmentAnalyzer] Applying smart extraction for {result.name} ({len(result.text_content)} chars)")
            extracted = await _smart_extract_relevant_sections(
                result.text_content,
                user_request,
                max_doc_chars
            )
            result.text_content = extracted
            result.summary = f"[Document: {result.name}] Smart-extracted {len(extracted)} chars relevant to request"

        results.append(result)

    return results


def build_multimodal_content(
    analyzed: List[AnalyzedAttachment],
    base_prompt: str,
    max_images: int = 10,
    max_total_image_chars: int = 3_000_000  # ~3MB total base64 to allow more images
) -> List[Dict[str, Any]]:
    """
    Build multimodal content blocks for Claude API.

    Args:
        analyzed: List of analyzed attachments
        base_prompt: The main prompt text
        max_images: Maximum number of images to include as vision content
        max_total_image_chars: Maximum total base64 characters for all images

    Returns:
        List of content blocks suitable for Claude messages API
    """
    content = []
    image_count = 0
    total_image_chars = 0
    context_parts = []

    # First pass: collect images for vision (with size limit)
    for att in analyzed:
        if att.is_vision_content and att.base64_data and image_count < max_images:
            # Check if adding this image would exceed total size limit
            if total_image_chars + len(att.base64_data) > max_total_image_chars:
                logger.warning(f"[AttachmentAnalyzer] Skipping image {att.name} - would exceed total size limit")
                continue

            # Add instruction for images
            if image_count == 0:
                content.append({
                    "type": "text",
                    "text": """<reference_attachments>
The user has uploaded files. CAREFULLY ANALYZE these and incorporate them into your response:
- For images: Match the design, colors, layout, and style
- For data: Use the actual numbers and create appropriate visualizations
- For documents: Extract and use the relevant content
</reference_attachments>"""
                })

            # Add the image
            media_type = att.mime_type or 'image/png'
            if not media_type.startswith('image/'):
                media_type = 'image/png'

            content.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_type,
                    "data": att.base64_data
                }
            })
            content.append({
                "type": "text",
                "text": f"[Reference: {att.name}]"
            })
            image_count += 1
            total_image_chars += len(att.base64_data)
            logger.info(f"[AttachmentAnalyzer] Added image {att.name} ({len(att.base64_data)} chars, total: {total_image_chars})")

    # Second pass: collect text context and extracted images from non-image files
    for att in analyzed:
        if not att.is_vision_content:
            if att.file_type == FileType.SPREADSHEET and att.text_content:
                context_parts.append(f"\n<data_file name='{att.name}'>\n{att.text_content}\nSuggested visualization: {att.chart_suggestion}\n</data_file>")
            elif att.file_type in [FileType.DOCUMENT, FileType.PRESENTATION] and att.text_content:
                # Document content - if smart extraction was used, this is already optimized
                # Otherwise limit to 30k chars (~7.5k tokens)
                context_parts.append(f"\n<document name='{att.name}'>\n{att.text_content[:30000]}\n</document>")

            # Add extracted images from PDFs (charts, graphs, etc.)
            if att.extracted_images and image_count < max_images:
                for idx, (img_base64, img_mime) in enumerate(att.extracted_images):
                    if image_count >= max_images:
                        break
                    if total_image_chars + len(img_base64) > max_total_image_chars:
                        logger.warning(f"[AttachmentAnalyzer] Skipping PDF image - would exceed size limit")
                        break

                    content.append({
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": img_mime,
                            "data": img_base64
                        }
                    })
                    content.append({
                        "type": "text",
                        "text": f"[Chart/Image from {att.name} - image {idx + 1}]"
                    })
                    image_count += 1
                    total_image_chars += len(img_base64)
                    logger.info(f"[AttachmentAnalyzer] Added PDF image {idx + 1} from {att.name} ({len(img_base64)} chars)")

    # Add context section if we have non-image files
    if context_parts:
        content.append({
            "type": "text",
            "text": "<uploaded_file_content>" + "".join(context_parts) + "\n</uploaded_file_content>"
        })

    # Add the main prompt
    content.append({
        "type": "text",
        "text": base_prompt
    })

    return content


def get_attachment_context_summary(analyzed: List[AnalyzedAttachment]) -> str:
    """
    Get a text summary of all analyzed attachments for prompt context.

    Args:
        analyzed: List of analyzed attachments

    Returns:
        Human-readable summary string
    """
    if not analyzed:
        return ""

    parts = ["\n<user_uploaded_files>"]
    for att in analyzed:
        parts.append(f"- {att.summary}")
    parts.append("</user_uploaded_files>")

    return "\n".join(parts)
